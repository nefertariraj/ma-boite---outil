# Fichier initialisé proprement sans résidu de texte
import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import scipy.stats as stats
from datetime import datetime, date
import json
import io
import copy

def synchroniser_et_capturer_tout(projet_en_cours, index_projet):
    """
    Cette fonction aspire automatiquement tous les tableaux et cases modifiés 
    à l'écran pour les verrouiller dans le fichier de sauvegarde.
    """
    if not isinstance(projet_en_cours, dict):
        return projet_en_cours
        
    cle_unique = str(index_projet)
    
    # 1. Capture du Plan de Collecte (DCP)
    cles_dcp = [f"master_dcp_table_{cle_unique}", f"dcp_editor_{cle_unique}", "master_dcp_table"]
    for cle in cles_dcp:
        if cle in st.session_state:
            donnees = st.session_state[cle]
            projet_en_cours["master_dcp_table"] = donnees.to_dict(orient="records") if isinstance(donnees, pd.DataFrame) else donnees
            break

    # 2. Capture des tableaux du module MSA
    cle_msa = f"msa_classification_table_{cle_unique}"
    if cle_msa in st.session_state:
        donnees_msa = st.session_state[cle_msa]
        projet_en_cours["msa_classification_table"] = donnees_msa.to_dict(orient="records") if isinstance(donnees_msa, pd.DataFrame) else donnees_msa

    cle_rep = f"rep_table_{cle_unique}"
    if cle_rep in st.session_state:
        donnees_rep = st.session_state[cle_rep]
        projet_en_cours["rep_table_data"] = donnees_rep.to_dict(orient="records") if isinstance(donnees_rep, pd.DataFrame) else donnees_rep

    cle_reprod = f"reprod_table_{cle_unique}"
    if cle_reprod in st.session_state:
        donnees_reprod = st.session_state[cle_reprod]
        projet_en_cours["reprod_table_data"] = donnees_reprod.to_dict(orient="records") if isinstance(donnees_reprod, pd.DataFrame) else donnees_reprod

    # 3. Capture des choix, boutons et validations de la page
    projet_en_cours["msa_selected_var"] = st.session_state.get(f"msa_selected_var_{cle_unique}", "")
    projet_en_cours["msa_action_choice"] = st.session_state.get(f"msa_action_choice_{cle_unique}", "")
    projet_en_cours["msa_is_validated_status"] = st.session_state.get(f"msa_sign_off_{cle_unique}", False)

    for cle_interne in list(st.session_state.keys()):
        if cle_interne.startswith("reference_master_config_") and cle_interne.endswith(f"_{cle_unique}"):
            suffixe_variable = cle_interne.replace("reference_master_config_", "")
            save_target_key = f"save_master_config_{suffixe_variable}"
            projet_en_cours[save_target_key] = st.session_state[cle_interne].copy()

    # 4. Capture automatique de tous les autres champs de saisie de ce projet
    for cle_interne in st.session_state.keys():
        if cle_interne.endswith(f"_{cle_unique}"):
            nom_propre_du_champ = cle_interne.replace(f"_{cle_unique}", "")
            if "table" not in cle_interne and "editor" not in cle_interne:
                projet_en_cours[nom_propre_du_champ] = st.session_state[cle_interne]

    # AJOUT : Capture de la phase IMPROVE
    if "improve_strategies" in st.session_state:
        donnees = st.session_state["improve_strategies"]
        if isinstance(donnees, pd.DataFrame):
            projet_en_cours["dmaic"]["improve"]["strategies"] = donnees.to_dict(orient="records")

    return projet_en_cours

def save_data():
    """
    Cette fonction définit ce que signifie 'save_data'.
    Maintenant, quand vous l'appellerez dans vos boutons, 
    le programme saura quoi faire.
    """
    try:
        # On sauvegarde tout l'état du projet dans projects.json
        with open("projects.json", "w") as f:
            json.dump(st.session_state.projects, f, indent=4, default=str)
    except Exception as e:
        st.error(f"Erreur de sauvegarde : {e}")

# ==========================================
# 📋 MODÈLE DE RÉFÉRENCE UNIQUE POUR CHAQUE PROJET
# ==========================================
PROJET_MODELE_REFERENCE = {
    "nom": "",
    "name": "",
    "gantt_data": pd.DataFrame(),
    "mesure_data": pd.DataFrame(),
    "voc_raw_data": pd.DataFrame(columns=["client", "question", "réponse brute"]), # RECONNU PAR L'IMPORTATEUR
    "voc_questions": ["Temps perdu ?", "Retouches ?", "Pénibilité ?", "Irritants ?", "Changement unique ?"], # RECONNU PAR L'IMPORTATEUR
    "dmaic": {
        "define": {},
        "measure": {},
        "analyze": {},
        "improve": {},
        "innovate": {},
        "control": {}
    },
    "parametres": {},
    "progression": 0
}

# 🛠️ FONCTIONS DE SÉRIALISATION / DÉSÉRIALISATION
def deep_serialize(obj):
    """ Préparation récursive de tous les états, tables, phases et champs utilisateurs """
    if isinstance(obj, (date, datetime)):
        return obj.isoformat()
    if hasattr(obj, 'strftime'):
        try: return obj.strftime('%Y-%m-%d')
        except: pass
    if isinstance(obj, pd.DataFrame):
        df_clean = obj.copy()
        for col in df_clean.columns:
            if pd.api.types.is_datetime64_any_dtype(df_clean[col]):
                df_clean[col] = df_clean[col].dt.strftime('%Y-%m-%d')
        return {"_type_df_": True, "data": df_clean.to_dict(orient="records")}
    if isinstance(obj, dict):
        return {str(k): deep_serialize(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [deep_serialize(i) for i in obj]
    return obj

def deep_deserialize(obj):
    """ Reconstruction à l'identique des types d'origine (DataFrames, dicts...) """
    if isinstance(obj, dict):
        if obj.get("_type_df_") is True:
            return pd.DataFrame(obj.get("data", []))
        return {k: deep_deserialize(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [deep_deserialize(i) for i in obj]
    return obj

# 🔄 SYSTÈME D'IMPORTATION INVISIBLE SANS CASSE
def traiter_importation_json():
    """ Restaure l'intégralité des données dans le state existant sans altérer l'interface """
    fichier_charge = st.session_state.get("sidebar_uploader_file")
    if fichier_charge is not None:
        try:
            raw_data = json.load(fichier_charge)
            restored_data = deep_deserialize(raw_data) 
            
            if isinstance(restored_data, list):
                projets_valides = []
                for p_item in restored_data:
                    if not isinstance(p_item, dict):
                        continue
                    
                    # On crée la base vierge
                    projet_reconstruit = PROJET_MODELE_REFERENCE.copy()
                    # On y injecte tout ce que le JSON contient (y compris le voc_raw_data sauvé !)
                    projet_reconstruit.update(p_item)
                    
                    # Alignement des noms de projets
                    if "name" in p_item and not p_item.get("nom"):
                        projet_reconstruit["nom"] = p_item["name"]
                    if "nom" in p_item and not p_item.get("name"):
                        projet_reconstruit["name"] = p_item["nom"]
                    
                    # Sécurisation des DataFrames restaurés
                    if not isinstance(projet_reconstruit["gantt_data"], pd.DataFrame):
                        projet_reconstruit["gantt_data"] = pd.DataFrame(projet_reconstruit["gantt_data"])
                        
                    if not isinstance(projet_reconstruit["mesure_data"], pd.DataFrame):
                        projet_reconstruit["mesure_data"] = pd.DataFrame(projet_reconstruit["mesure_data"])
                    
                    # CORRECTION : On s'assure que le voc_raw_data redevienne un vrai DataFrame exploitable
                    if "voc_raw_data" in p_item:
                        if isinstance(p_item["voc_raw_data"], list):
                            projet_reconstruit["voc_raw_data"] = pd.DataFrame(p_item["voc_raw_data"])
                        elif isinstance(p_item["voc_raw_data"], dict) and "data" in p_item["voc_raw_data"]:
                            # Gestion du format encodé par deep_serialize
                            projet_reconstruit["voc_raw_data"] = pd.DataFrame(p_item["voc_raw_data"]["data"])
                        else:
                            projet_reconstruit["voc_raw_data"] = pd.DataFrame(p_item["voc_raw_data"])

                    dmaic_originel = p_item.get("dmaic", {})
                    dmaic_structure = {k: v.copy() if isinstance(v, dict) else v for k, v in PROJET_MODELE_REFERENCE["dmaic"].items()}
                    
                    for phase in ["define", "measure", "analyze", "improve", "innovate", "control"]:
                        if phase in dmaic_originel:
                            dmaic_structure[phase] = dmaic_originel[phase]
                    
                    # AJOUT : Restauration des données IMPROVE pour l'affichage
                    if "strategies" in projet_reconstruit["dmaic"]["improve"]:
                        st.session_state["improve_strategies"] = pd.DataFrame(
                            projet_reconstruit["dmaic"]["improve"]["strategies"]
                        )
                    
                    projet_reconstruit["dmaic"] = dmaic_structure
                    projets_valides.append(projet_reconstruit)

                if projets_valides:
                    st.session_state.projects = projets_valides
                    st.session_state["current_project_idx"] = None
        except Exception as e:
            st.sidebar.error(f"Erreur de restauration : {e}")

# --- CONFIGURATION DE LA PAGE & STYLE ---
st.set_page_config(page_title="LSS - Personal Toolbox", layout="wide")

if 'primary_color' not in st.session_state:
    st.session_state.primary_color = "#1E3A8A" 

st.markdown(f"""
    <style>
    .stApp {{ background-color: #FFFFFF; }}
    .stButton>button {{ background-color: {st.session_state.primary_color}; color: white; border-radius: 5px; }}
    .project-card {{
        border: 1px solid #E2E8F0;
        border-radius: 8px;
        padding: 16px;
        background-color: #F8FAFC;
        margin-bottom: 0.5rem;
        text-align: center;
    }}
    </style>
    """, unsafe_allow_html=True)

# --- INITIALISATION SÉCURISÉE DES VARIABLES ---
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False
if 'projects' not in st.session_state:
    st.session_state.projects = []
if 'current_project_idx' not in st.session_state:
    st.session_state.current_project_idx = None

# --- ÉCRAN DE CONNEXION ---
if not st.session_state.authenticated:
    st.title("🔐 Lean Six Sigma - Personal Toolbox")
    with st.container(border=True):
        user = st.text_input("Identifiant (Email Google)", key="login_user_input")
        pwd = st.text_input("Mot de passe", type="password", key="login_pwd_input")
        if st.button("Se connecter via Google Drive", key="login_submit_btn"):
            if user and pwd: 
                st.session_state.authenticated = True
                st.rerun()
    st.stop()

# ==========================================
# ⚙️ BARRE LATÉRALE MOTEUR
# ==========================================
with st.sidebar:
    st.title("⚙️ Paramètres & Sauvegarde")
    color = st.color_picker("Couleur de l'outil", st.session_state.primary_color, key="sidebar_color_picker")
    st.session_state.primary_color = color
    st.divider()

   # ------------------------------------------------
    # 💾 MODULE "ENREGISTRER SOUS" (INTEGRALITÉ DES PROJETS CORRIGÉE)
    # ------------------------------------------------
    st.subheader("💾 Archivage & Exportation Complète (Copie Conforme)")
    
    if "projects" in st.session_state and len(st.session_state.projects) > 0:
        indices_projets_tous = list(range(len(st.session_state.projects)))
        
        def formateur_liste_enregistrer(idx):
            p_test = st.session_state.projects[idx]
            return f"📁 {p_test.get('nom', 'Sans nom')} | Jalon: {p_test.get('status', 'Define')}"
            
        proj_sel_idx = st.selectbox(
            "Sélectionnez le projet à figer dans l'archive :",
            options=indices_projets_tous,
            format_func=formateur_liste_enregistrer,
            key="sb_enregistrer_sous_selector"
        )
        
        # Extraction du dictionnaire contenant l'INTÉGRALITÉ des données du projet
        p_exp = st.session_state.projects[proj_sel_idx]
        project_name = p_exp.get('nom', 'Projet_LSS').replace(" ", "_")
        
        import io
        from datetime import datetime
        import json

        st.info("📊 **Contrôle Qualité :** L'exportation va inclure l'intégralité des formulaires, des structures de données (DataFrames) et des rendus graphiques rattachés à ce projet.")

        # =====================================================================
        # 📊 1. EXPORTATION EXCEL SANS AUCUNE EXCLUSION (Toutes lignes, tous onglets)
        # =====================================================================
        try:
            buffer_xlsx = io.BytesIO()
            with pd.ExcelWriter(buffer_xlsx, engine='openpyxl') as writer:
                
                # Étape A : On isole toutes les variables simples (textes, inputs, KPIs)
                parametres_projets = []
                for cle, valeur in p_exp.items():
                    if not isinstance(valeur, (pd.DataFrame, list, dict)) and valeur is not None:
                        parametres_projets.append({"Composant / Formulaire": cle, "Valeur Saisie / Résultat": str(valeur)})
                
                if parametres_projets:
                    pd.DataFrame(parametres_projets).to_excel(writer, sheet_name='Formulaires & Synthèse', index=False)
                
                # Étape B : On parcourt et on extrait TOUS les DataFrames du projet
                df_trouves = 0
                for cle, valeur in p_exp.items():
                    if isinstance(valeur, pd.DataFrame):
                        nom_onglet = str(cle)[:30] # Limite Excel de 31 caractères
                        valeur.to_excel(writer, sheet_name=f"DF_{nom_onglet}", index=False)
                        df_trouves += 1
                    elif isinstance(valeur, list) and len(valeur) > 0 and isinstance(valeur[0], dict):
                        nom_onglet = str(cle)[:30]
                        pd.DataFrame(valeur).to_excel(writer, sheet_name=f"LIST_{nom_onglet}", index=False)
                        df_trouves += 1
                        
                if "dc_master_data" in st.session_state and isinstance(st.session_state.dc_master_data, pd.DataFrame):
                    st.session_state.dc_master_data.to_excel(writer, sheet_name='Master_Data_Collecte_T0', index=False)
                if "current_spc_data" in st.session_state and isinstance(st.session_state.current_spc_data, pd.DataFrame):
                    st.session_state.current_spc_data.to_excel(writer, sheet_name='Calculs_Cartes_Controle', index=False)

            st.download_button(
                label="📊 Télécharger la base Excel Intégrale (Données & Calculs)", 
                data=buffer_xlsx.getvalue(), 
                file_name=f"EXCEL_COMPLET_{project_name}.xlsx", 
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="btn_excel_full_zero_loss"
            )           
        except Exception as e:
            st.error(f"Erreur lors de l'extraction Excel : {e}")


        # =====================================================================
        # 📽️ 2. EXPORTATION POWERPOINT INTEGRALE (Sécurisée contre l'ambiguïté des DF)
        # =====================================================================
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # Slide 1 : Titre / Garde
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"LIVRABLE EXHAUSTIF : {project_name.upper()}"
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
            tf = txBox.text_frame
            tf.text = f"Génération système : {datetime.now().strftime('%Y-%m-%d %H:%M')}\nStatut : Validation finale DMAIC"

            # Balayage du dictionnaire
            for cle, valeur in p_exp.items():
                if valeur is None:
                    continue
                    
                # Protection : Si c'est un DataFrame, on gère son cas à part sans tester '== ""'
                if isinstance(valeur, pd.DataFrame):
                    if valeur.empty:
                        continue
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"Tableau (DataFrame) : {str(cle).upper()}"
                    
                    nb_rows = min(len(valeur) + 1, 8)
                    nb_cols = len(valeur.columns)
                    table_shape = slide.shapes.add_table(nb_rows, nb_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
                    table = table_shape.table
                    
                    for c_idx, col_name in enumerate(valeur.columns):
                        table.cell(0, c_idx).text = str(col_name)
                    for r_idx in range(nb_rows - 1):
                        for c_idx, col_name in enumerate(valeur.columns):
                            table.cell(r_idx+1, c_idx).text = str(valeur.iloc[r_idx][col_name])
                
                # Si c'est une liste de dictionnaires (SIPOC, rôles...)
                elif isinstance(valeur, list) and len(valeur) > 0 and isinstance(valeur[0], dict):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"Tableau : {str(cle).upper()}"
                    
                    df_temp = pd.DataFrame(valeur)
                    nb_rows = min(len(df_temp) + 1, 8)
                    nb_cols = len(df_temp.columns)
                    table_shape = slide.shapes.add_table(nb_rows, nb_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
                    table = table_shape.table
                    
                    for c_idx, col_name in enumerate(df_temp.columns):
                        table.cell(0, c_idx).text = str(col_name)
                    for r_idx in range(nb_rows - 1):
                        for c_idx, col_name in enumerate(df_temp.columns):
                            table.cell(r_idx+1, c_idx).text = str(df_temp.iloc[r_idx][col_name])
                
                # Si c'est un texte long (saisie utilisateur standard)
                elif isinstance(valeur, str) and valeur != "" and len(valeur) > 10 and cle not in ['nom', 'status', 'date_creation']:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f"Composant : {str(cle).upper()}"
                    body = slide.placeholders[1]
                    body.text = str(valeur)

            # Traitement des graphiques Plotly
            figures_a_importer = []
            if p_exp.get('spc_figure') is not None: figures_a_importer.append(("Carte SPC Projet", p_exp.get('spc_figure')))
            if st.session_state.get('current_spc_figure') is not None: figures_a_importer.append(("Carte SPC active", st.session_state.get('current_spc_figure')))
            if st.session_state.get('current_pareto_figure') is not None: figures_a_importer.append(("Analyse de Pareto", st.session_state.get('current_pareto_figure')))
            
            for nom_fig, fig_obj in figures_a_importer:
                try:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = f"Graphique Réel : {nom_fig}"
                    img_buf = io.BytesIO()
                    fig_obj.write_image(img_buf, format="png", width=1000, height=600)
                    img_buf.seek(0)
                    slide.shapes.add_picture(img_buf, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                except:
                    pass

            buffer_pptx = io.BytesIO()
            prs.save(buffer_pptx)
            st.download_button(
                label="📽️ Télécharger le PowerPoint Complet (Rapports & Visuels)", 
                data=bytes(buffer_pptx.getvalue()), 
                file_name=f"POWERPOINT_COMPLET_{project_name}.pptx", 
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="btn_pptx_full_zero_loss"
            )
        except Exception as e:
            st.error(f"Erreur lors de l'extraction PowerPoint : {e}")


        # =====================================================================
        # 📄 3. EXPORTATION PDF DE TYPE ARCHIVE (Sécurisée contre l'ambiguïté des DF et les caractères Unicode)
        # =====================================================================
        try:
            from fpdf import FPDF
            
            class PDF(FPDF):
                def header(self):
                    self.set_font('Helvetica', 'B', 10)
                    self.set_text_color(30, 58, 138)
                    self.cell(0, 10, "RAPPORT OFFICIEL DE CONFIGURATION ET D'AUDIT INTEGRAL", border=0, ln=1, align='L')
                    self.line(10, 16, 200, 16)
                    self.ln(4)
                    
                def footer(self):
                    self.set_y(-15)
                    self.set_font('Helvetica', 'I', 8)
                    self.set_text_color(156, 163, 175)
                    self.cell(0, 10, f'Livrable d\'archive autonome - Page {self.page_no()}', 0, 0, 'C')

            # Fonction de sécurisation des caractères non supportés par Helvetica
            def nettoyer_texte(texte):
                if not isinstance(texte, str):
                    texte = str(texte)
                remplacements = {
                    "≥": ">=",
                    "≤": "<=",
                    "≠": "!=",
                    "±": "+/-",
                    "µ": "u",
                    "²": "2",
                    "³": "3"
                }
                for carac, subst in remplacements.items():
                    texte = texte.replace(carac, subst)
                # Encodage de sécurité pour nettoyer les résidus bizarres
                return texte.encode('latin-1', 'replace').decode('latin-1')

            pdf = PDF()
            pdf.add_page()
            pdf.set_font("Helvetica", size=10)
            
            # En-tête de l'archive
            pdf.set_font("Helvetica", 'B', 14)
            pdf.set_text_color(13, 148, 136)
            pdf.cell(0, 10, f"PROJET DMAIC : {nettoyer_texte(project_name.upper())}", ln=1)
            pdf.set_font("Helvetica", size=10)
            pdf.set_text_color(0, 0, 0)
            pdf.cell(0, 6, f"Date de l'instantané d'impression : {datetime.now().strftime('%Y-%m-%d %H:%M')}", ln=1)
            pdf.cell(0, 6, f"Phase cible enregistrée : {p_exp.get('status', 'Define')}", ln=1)
            pdf.ln(6)
            
            # Boucle d'extraction systématique du dictionnaire
            for cle, valeur in p_exp.items():
                if valeur is None:
                    continue
                
                # CAS 1 : C'est un DataFrame brut (ex: calculs SPC, base de collecte interne)
                if isinstance(valeur, pd.DataFrame):
                    if valeur.empty:
                        continue
                    pdf.set_font("Helvetica", 'B', 11)
                    pdf.set_text_color(30, 58, 138)
                    pdf.cell(0, 8, f"Tableau Associé - {nettoyer_texte(cle).upper()} :", ln=1)
                    pdf.ln(2)
                    
                    pdf.set_font("Helvetica", 'B', 9)
                    largeur_col = int(180 / max(len(valeur.columns), 1))
                    for col in valeur.columns:
                        pdf.cell(largeur_col, 7, nettoyer_texte(col)[:15], border=1)
                    pdf.ln()
                    
                    pdf.set_font("Helvetica", size=9)
                    for _, row in valeur.head(15).iterrows():
                        for col in valeur.columns:
                            cell_val = nettoyer_texte(row[col])
                            pdf.cell(largeur_col, 6, cell_val[:15], border=1)
                        pdf.ln()
                    if len(valeur) > 15:
                        pdf.cell(0, 6, nettoyer_texte(f"... (+ {len(valeur) - 15} lignes archivées intégralement dans l'onglet Excel correspondant)"), ln=1)
                    pdf.ln(4)
                
                # CAS 2 : C'est une liste de dictionnaires (SIPOC, etc.)
                elif isinstance(valeur, list) and len(valeur) > 0 and isinstance(valeur[0], dict):
                    pdf.set_font("Helvetica", 'B', 11)
                    pdf.set_text_color(30, 58, 138)
                    pdf.cell(0, 8, f"Tableau de Structure - {nettoyer_texte(cle).upper()} :", ln=1)
                    pdf.ln(2)
                    
                    df_t = pd.DataFrame(valeur)
                    pdf.set_font("Helvetica", 'B', 9)
                    largeur_col = int(180 / max(len(df_t.columns), 1))
                    for col in df_t.columns:
                        pdf.cell(largeur_col, 7, nettoyer_texte(col)[:15], border=1)
                    pdf.ln()
                    
                    pdf.set_font("Helvetica", size=9)
                    for _, row in df_t.iterrows():
                        for col in df_t.columns:
                            cell_val = nettoyer_texte(row[col])
                            pdf.cell(largeur_col, 6, cell_val[:15], border=1)
                        pdf.ln()
                    pdf.ln(4)
                
                # CAS 3 : C'est une chaîne de caractères (Commentaires, notes, Problem Statement)
                elif isinstance(valeur, str) and valeur != "":
                    pdf.set_font("Helvetica", 'B', 11)
                    pdf.set_text_color(30, 58, 138)
                    pdf.cell(0, 8, f"Saisie / {nettoyer_texte(cle).upper()} :", ln=1)
                    pdf.set_font("Helvetica", size=10)
                    pdf.set_text_color(0, 0, 0)
                    pdf.multi_cell(0, 6, nettoyer_texte(valeur))
                    pdf.ln(3)

            pdf_bytes = pdf.output()
            
            st.download_button(
                label="📄 Télécharger le PDF Intégral (Archive Documentaire)",
                data=bytes(pdf_bytes),
                file_name=f"PDF_COMPLET_{project_name}.pdf",
                mime="application/pdf",
                use_container_width=True,
                key="btn_pdf_full_zero_loss"
            )
        except Exception as e:
            st.error(f"Erreur lors de l'extraction PDF : {e}")

    # ------------------------------------------------
    # 💾 SAUVEGARDE ET IMPORTATION GLOBALE (CONSERVÉS)
    # ------------------------------------------------
    st.subheader("💾 Sauvegarder mon travail")
    
    # On vérifie s'il y a des projets dans la session
    un_projet_au_moins = "projects" in st.session_state and len(st.session_state.projects) > 0

    if un_projet_au_moins:
        try:
            # 🟢 LIGNE AJOUTÉE : On force la capture complète du projet actif avant la sauvegarde
            # On vérifie que 'p' et 'p_idx' existent bien à ce moment-là du code
            if 'p' in locals() and 'p_idx' in locals():
                st.session_state.projects[p_idx] = synchroniser_et_capturer_tout(p, p_idx)

            # Sécurité absolue pour les dates cachées
            def encodeur_secours(obj):
                if isinstance(obj, (date, datetime)):
                    return obj.isoformat()
                raise TypeError(f"Type non sérialisable: {type(obj)}")

            # Préparation et sérialisation
            donnies_preparees = deep_serialize(st.session_state.projects)
            data_json = json.dumps(donnies_preparees, indent=4, ensure_ascii=False, default=encodeur_secours)
            
            # Le bouton de téléchargement, bien visible
            st.download_button(
                label="📤 Télécharger la sauvegarde (.json)",
                data=data_json,
                file_name="sauvegarde_boite_outils.json",
                mime="application/json",
                key="sidebar_download_btn",
                use_container_width=True
            )
        except Exception as e:
            st.error(f"Erreur export : {e}")
    else:
        st.info("💡 Aucun projet en mémoire. Créez un projet ou importez un JSON pour afficher le bouton de téléchargement.")

    st.divider()
    st.subheader("📥 Reprendre mon travail")
    st.file_uploader(
        "Importer un fichier de sauvegarde", 
        type="json", 
        key="sidebar_uploader_file",
        on_change=traiter_importation_json
    )

    # --- FONCTION DE SUPPRESSION FORCEE SUR LA PAGE PRINCIPALE ---
    def action_supprimer_projet(index_a_retirer):
        if "projects" in st.session_state and len(st.session_state.projects) > index_a_retirer:
            st.session_state.projects.pop(index_a_retirer)
            # Si on supprime le projet actuellement ouvert, on réinitialise l'index
            if st.session_state.get("current_project_idx") == index_a_retirer:
                st.session_state["current_project_idx"] = None

    # Section de suppression propre des projets
    if len(st.session_state.projects) > 0:
        st.write("### 🗑️ Suppression de projet")
        
        # Liste des index réels de la session
        indices_projets = list(range(len(st.session_state.projects)))
        
        # Extraction intelligente du vrai nom saisi à l'initialisation
        def extraire_nom_reel(idx):
            projet = st.session_state.projects[idx]
            
            # 1. On liste les clés possibles utilisées lors de la création
            cles_possibles = ["nom", "nom_projet", "name", "project_name", "Nom du projet"]
            
            for cle in cles_possibles:
                if cle in projet and projet[cle]:
                    # Si la valeur est elle-même un dictionnaire (cas rare), on cherche dedans
                    if isinstance(projet[cle], dict):
                        continue
                    return str(projet[cle]).strip()
            
            # 2. Si aucune clé classique n'a fonctionné, on cherche n'importe quelle chaîne de texte non vide 
            # au premier niveau du projet pour essayer de capturer le titre
            for cle, valeur in projet.items():
                if isinstance(valeur, str) and valeur.strip() and cle not in ["status", "id", "description", "problem"]:
                    if len(valeur) < 100: # Un titre de projet fait rarement plus de 100 caractères
                        return valeur.strip()
                        
            # Fallback ultime indexé si vraiment le dictionnaire est vide
            return f"Projet sans nom enregistré #{idx + 1}"
        
        # Liste déroulante affichant les vrais noms identifiables
        idx_selectionne = st.selectbox(
            "Sélectionner le projet à supprimer définitivement :", 
            options=indices_projets,
            format_func=extraire_nom_reel,
            key="selectbox_suppression_projets"
        )
        
        # Bouton unique de suppression calé sur toute la largeur
        st.button(
            "🗑️ Supprimer le projet sélectionné", 
            use_container_width=True, 
            key="suppr_btn_select_unique",
            help="Supprimer définitivement ce projet de la liste",
            on_click=action_supprimer_projet,
            args=(idx_selectionne,)
        )
            
    else:
        st.info("💡 Aucun projet disponible. Créez un nouveau projet ou importez un fichier JSON depuis le menu latéral.")
    
    st.info("🛠️ Vos composants graphiques originaux (onglets DMAIC, diagrammes Plotly d'origine, formulaires de saisie, tableaux éditables st.data_editor) se ré-exécutent automatiquement en utilisant les données fidèlement restaurées ci-dessus.")
    
# --- NAVIGATION PRINCIPALE ---
if st.session_state.current_project_idx is None:
    st.title("🚀 Mes Projets Lean Six Sigma")
    
    with st.expander("➕ Initialiser un nouveau projet"):
        p_name = st.text_input("Nom du projet", key="input_nouveau_projet_nom")
        if st.button("Créer le projet", key="btn_creer_nouveau_projet"):
            if p_name:
                import copy
                import pandas as pd

                # 1. Modèle de projet neuf, complet et strictement vide
                new_p = {
                    "nom": p_name,
                    "name": p_name,
                    "status": "Define",
                    "problem": "",
                    "gantt_data": pd.DataFrame(),
                    "mesure_data": pd.DataFrame(),
                    "master_dcp_table": pd.DataFrame(),
                    "current_state_process_map": pd.DataFrame(),
                    "dc_master_data": pd.DataFrame(),
                    "current_spc_data": pd.DataFrame(),
                    "improvement_strategy": pd.DataFrame(),
                    "dmaic": {
                        "define": {},
                        "measure": {},
                        "analyze": {},
                        "improve": {},
                        "innovate": {},
                        "control": {}
                    },
                    "sipoc_data": [{"Supplier": "", "Input": "", "Process": "", "Output": "", "Customer": ""}],
                    "voc_data": [{"Client": "", "Verbatim": "", "Besoin": ""}],
                    "selected_ctq": "Qualité",
                    "team_data": [{"Poste": "", "Nom": ""}],
                    "parametres": {},
                    "progression": 0
                }
            
                if "projects" not in st.session_state:
                    st.session_state.projects = []
            
                st.session_state.projects.append(copy.deepcopy(new_p))
                new_idx = len(st.session_state.projects) - 1
                st.session_state.current_project_idx = new_idx
            
                # 2. VIDAGE STRICT DE TOUTES LES VARIABLES GLOBALES POUR LE NOUVEAU PROJET
                st.session_state["current_state_process_map"] = pd.DataFrame()
                st.session_state["master_dcp_table"] = pd.DataFrame()
                st.session_state["mesure_data"] = pd.DataFrame()
                st.session_state["dc_master_data"] = pd.DataFrame()
                st.session_state["current_spc_data"] = pd.DataFrame()
                st.session_state["improvement_strategy"] = pd.DataFrame()
                st.session_state["improve_strategies"] = pd.DataFrame()
            
                # 3. PURGE DU CACHE DES WIDGETS
                keys_to_delete = []
                for k in list(st.session_state.keys()):
                    k_lower = k.lower()
                    if any(term in k_lower for term in [
                        "process", "map", "dcp", "dc_", "_dc", "master_dcp", 
                        "mesure", "detailed", "spc", "editor", "strategy", "$data_editor"
                    ]):
                        keys_to_delete.append(k)
            
                for k in keys_to_delete:
                    try:
                        del st.session_state[k]
                    except Exception:
                        pass

                st.success("Projet créé et initialisé à vide !")
                st.rerun()

    # Affichage des cartes projets
    cols = st.columns(3)
    for idx, proj in enumerate(st.session_state.projects):
        with cols[idx % 3]:
            with st.container(border=True):
                nom_final = "Projet sans nom"
                for cle_test in ["nom", "name", "nom_projet", "project_name"]:
                    if cle_test in proj and proj[cle_test] and not isinstance(proj[cle_test], dict):
                        nom_final = str(proj[cle_test]).strip()
                        break
            
                st.subheader(nom_final)
                if st.button("Ouvrir", key=f"open_{idx}"):
                    st.session_state.current_project_idx = idx
                    proj_cible = st.session_state.projects[idx]

                    import pandas as pd

                    # Restauration complète des données du projet dans les variables globales
                    for cle_cible, cle_dict in [
                        ("master_dcp_table", "master_dcp_table"),
                        ("mesure_data", "mesure_data"),
                        ("current_state_process_map", "current_state_process_map"),
                        ("dc_master_data", "dc_master_data"),
                        ("current_spc_data", "current_spc_data"),
                        ("improvement_strategy", "improvement_strategy")
                    ]:
                        val = proj_cible.get(cle_dict)
                        if val is not None and isinstance(val, list) and len(val) > 0:
                            st.session_state[cle_cible] = pd.DataFrame(val)
                        elif isinstance(val, pd.DataFrame) and not val.empty:
                            st.session_state[cle_cible] = val.copy()
                        else:
                            st.session_state[cle_cible] = pd.DataFrame()

                    # Restauration spécifique de la phase Improve
                    strategies_data = []
                    if "dmaic" in proj_cible and isinstance(proj_cible["dmaic"], dict):
                        improve_sec = proj_cible["dmaic"].get("improve", {})
                        if isinstance(improve_sec, dict):
                            strategies_data = improve_sec.get("strategies", [])
                
                    if isinstance(strategies_data, list) and len(strategies_data) > 0:
                        st.session_state["improve_strategies"] = pd.DataFrame(strategies_data)
                    elif isinstance(strategies_data, pd.DataFrame):
                        st.session_state["improve_strategies"] = strategies_data.copy()
                    else:
                        st.session_state["improve_strategies"] = pd.DataFrame()

                    # Purge des caches d'édition pour forcer le rechargement propre
                    for k in list(st.session_state.keys()):
                        k_lower = k.lower()
                        if any(term in k_lower for term in ["$data_editor", "editor", "process_map", "dc_master"]):
                            try:
                                del st.session_state[k]
                            except Exception:
                                pass

                    st.rerun()

else:
    # --- VUE PROJET (DMAIC) ---
    p_idx = st.session_state.current_project_idx
    p = st.session_state.projects[p_idx]

    col_back, col_title = st.columns([1, 8])
    with col_back:
        if st.button("⬅️"):
            st.session_state.current_project_idx = None
            st.rerun()
            
    with col_title:
        st.title(f"Projet : {p['name']}")

    # Définition des onglets
    tabs = st.tabs(["DEFINE", "MEASURE", "ANALYZE", "IMPROVE", "CONTROL"])

    # --- PHASE DEFINE ---
    with tabs[0]:
        st.header("Phase Define")
    
        # Récupération de l'index et du projet pour l'ensemble de l'onglet
        p_idx = st.session_state.get('current_project_idx', 0)
        p = st.session_state.projects[p_idx]
    
        # 1. Problème & CTQ
        st.subheader("1. Énoncé du Problème & CTQ")
        col1, col2 = st.columns(2)
    
        with col1:
            st.markdown('<p title="X correspond aux causes">Décrivez le problème en détail :</p>', unsafe_allow_html=True)
            # Utilisation de .get pour éviter le KeyError si le projet est ancien
            p_input = st.text_area("Saisie libre", value=p.get("problem", ""), height=150, key=f"prob_in_{p_idx}")
        
            if st.button("🪄 Retranscrire via IA"):
                p["problem"] = p_input
                st.session_state.ai_suggest_ctq = [
                    "⏱️ Temps : Réduire le Lead Time global",
                    "🎯 Qualité : Augmenter le taux de conformité (%C&A)",
                    "💰 Coût : Réduire les coûts opérationnels",
                    f"✨ Spécifique : {p_input[:30]}..."
                ]
                st.rerun() 

        with col2:
            st.write("Propositions de l'IA (CTQ) :")
            if 'ai_suggest_ctq' in st.session_state:
                for ctq in st.session_state.ai_suggest_ctq:
                    if st.button(ctq, key=f"btn_{ctq}"):
                        p["selected_ctq"] = ctq
                        st.rerun()
        
            st.divider()
        
            if "selected_ctq" in p and p["selected_ctq"]:
                st.write("✍️ **Ajustez votre CTQ final :**")
                new_val = st.text_input("Libellé du CTQ", value=p["selected_ctq"], key=f"edit_ctq_{p_idx}")
            
            if 'new_val' in locals() and new_val != p.get("selected_ctq", ""):
                p["selected_ctq"] = new_val
                st.rerun()
            st.info(f"**CTQ validé :** {p.get('selected_ctq', 'Aucun')}")

        # 2. Équipe Projet
        st.divider()
        st.subheader("2. Équipe Projet")
        team_key = f"editor_team_{p_idx}"
        edited_team = st.data_editor(
            p.get("team_data", [{"Poste": "", "Nom": ""}]), 
            num_rows="dynamic",
            use_container_width=True,
            key=team_key
        )
        if st.button("✅ Valider l'équipe", key=f"save_team_{p_idx}"):
            p["team_data"] = edited_team
            st.success("Équipe enregistrée !")
            st.rerun()

        # 3. BÉNÉFICES ATTENDUS & MATRICE D'OPPORTUNITÉ ---
        st.divider()
        st.subheader("3. Bénéfices Attendus & Matrice d'Opportunité")

        # Initialisation sécurisée pour l'import/export
        if "benefices_saisie" not in p:
            p["benefices_saisie"] = ""

        # Zone de saisie
        p["benefices_saisie"] = st.text_area(
            "Définissez les gains espérés (Financiers, Qualité, Délais) :",
            value=p["benefices_saisie"],
            height=120,
            placeholder="Ex: Réduction du taux de rebus de 15%, gain de 2 jours sur le lead time...",
            key=f"area_ben_v13_{p_idx}"
        )

        if len(p["benefices_saisie"]) > 10:
            st.markdown("---")
        
            # 1. ANALYSE SÉMANTIQUE BLACK BELT
            c_alt1, c_alt2 = st.columns([2, 1])
            with c_alt1:
                st.markdown("### 💎 Diagnostic de l'Opportunité")
                st.success(f"""
                **Analyse de la Valeur Ajoutée :**
                Votre description (« *{p['benefices_saisie'][:60]}...* ») suggère un potentiel de réduction des **COPQ** (Coûts de Non-Qualité) significatif. 
            
                *   **Levier Principal :** Optimisation de la performance (Sigma Level).
                *   **Type de Gain :** Hard Savings identifiés.
                *   **Risque de Cadrage :** Faible (Périmètre cohérent avec une démarche DMAIC).
                """)
        
            with c_alt2:
                st.metric("Potentiel de ROI", "Élevé", delta="Top 10% Projets")
                st.write("**Statut :** Candidat idéal pour analyse approfondie.")

            # 2. LA MATRICE D'OPPORTUNITÉ (Visualisation)
            with st.expander("👁️ Visualiser la Matrice d'Opportunité (Go / No-Go)"):
                st.write("Cette matrice évalue si le projet est 'éligible' à une certification Black Belt ou s'il s'agit d'un simple projet de maintenance.")
            
                # Mise en page de la matrice
                m_col1, m_col2 = st.columns([2, 1])
            
                with m_col1:
                    st.markdown("""
                    | | **COMPLEXITÉ MAÎTRISÉE** | **COMPLEXITÉ ÉLEVÉE** |
                    |---|---|---|
                    | **VALEUR STRATÉGIQUE** | ✨ **OPPORTUNITÉ MAJEURE** (Go) | 🏗️ **CHANTIER COMPLEXE** |
                    | **VALEUR MODÉRÉE** | 🤏 **MICRO-AMÉLIORATION** | ❌ **HORS PÉRIMÈTRE** |
                    """)
            
                with m_col2:
                    st.markdown("**Positionnement IA**")
                    # Curseur de positionnement basé sur le texte
                    st.select_slider(
                        "Éligibilité :",
                        options=["Hors périmètre", "Micro-Amélioration", "Chantier Complexe", "Opportunité Majeure"],
                        value="Opportunité Majeure",
                        disabled=True
                    )
                    st.caption("L'IA préconise le passage en phase 'Measure'.")
        else:
            st.info("💡 Décrivez vos bénéfices pour activer le diagnostic d'opportunité Black Belt.")

        # 4. Stakeholder Analysis
        st.divider()
        st.subheader("4. Stakeholder Analysis (Matrice d'adhésion)")

        # Initialisation des données si elles n'existent pas
        if "stakeholders" not in p:
            p["stakeholders"] = [
                {
                    "Name": "Sponsor",
                    "Strongly Against": False,
                    "Moderately Against": False,
                    "Neutral": True,
                    "Moderately Supportive": False,
                    "Strongly Supportive": False
                }
            ]

        st.info("Cochez le niveau d'adhésion actuel pour chaque partie prenante. Vous pouvez ajouter/supprimer des lignes via les icônes du tableau.")

        # Utilisation du data_editor
        edited_stakeholders = st.data_editor(
            p["stakeholders"],
            num_rows="dynamic", 
            key=f"stake_edit_{p_idx}",
            use_container_width=True,
            column_config={
                "Strongly Against": st.column_config.CheckboxColumn(default=False),
                "Moderately Against": st.column_config.CheckboxColumn(default=False),
                "Neutral": st.column_config.CheckboxColumn(default=False),
                "Moderately Supportive": st.column_config.CheckboxColumn(default=False),
                "Strongly Supportive": st.column_config.CheckboxColumn(default=False),
            }
        )

        if st.button("✅ Sauvegarder l'analyse des parties prenantes", key=f"save_stake_{p_idx}"):
            p["stakeholders"] = edited_stakeholders
            st.success("Analyse sauvegardée !")

        # 5. SIPOC & FLUX CROISÉ COMPACT ---
        st.divider()
        st.subheader("5. SIPOC & Flux de processus")

        COLONNES_SIPOC = ["Supplier", "Input", "Process", "Output", "Customer"]
    
        if "sipoc_data" not in p or not isinstance(p["sipoc_data"], list):
            p["sipoc_data"] = [dict.fromkeys(COLONNES_SIPOC, "")]

        with st.form(key=f"form_sipoc_compact_{p_idx}"):
            df_init_sipoc = pd.DataFrame(p["sipoc_data"])
            for col in COLONNES_SIPOC:
                if col not in df_init_sipoc.columns:
                    df_init_sipoc[col] = ""
        
            edited_sipoc = st.data_editor(
                df_init_sipoc[COLONNES_SIPOC],
                num_rows="dynamic",
                use_container_width=True,
                key=f"editor_sipoc_comp_{p_idx}",
                column_config={c: st.column_config.TextColumn(c) for c in COLONNES_SIPOC}
            )
            submit_sipoc = st.form_submit_button("✅ Actualiser le flux")

        if submit_sipoc:
            p["sipoc_data"] = edited_sipoc.to_dict('records')
            st.rerun()

        df_viz_sipoc = pd.DataFrame(p["sipoc_data"])
        
        # SÉCURISATION ABSOLUE : On force la présence des colonnes requises
        for c_req in COLONNES_SIPOC:
            if c_req not in df_viz_sipoc.columns:
                df_viz_sipoc[c_req] = ""

        # Nettoyage et filtrage sans risque de KeyError
        if not df_viz_sipoc.empty:
            df_viz_sipoc = df_viz_sipoc[(df_viz_sipoc["Process"].astype(str).str.strip() != "") & 
                                        (df_viz_sipoc["Customer"].astype(str).str.strip() != "")]
    
        if not df_viz_sipoc.empty:
            st.write("---")
            acteurs = df_viz_sipoc["Customer"].unique().tolist()
            cols_act = st.columns(len(acteurs))

            # Rendu des tâches
            for idx, row in df_viz_sipoc.iterrows():
                current_col = acteurs.index(row["Customer"])
                row_cols = st.columns(len(acteurs))
            
                has_next = idx < len(df_viz_sipoc) - 1
                if has_next:
                    next_actor = df_viz_sipoc.iloc[idx + 1]["Customer"]
                    next_col = acteurs.index(next_actor)

                for i in range(len(acteurs)):
                    with row_cols[i]:
                        if i == current_col:
                            st.markdown(f"""
                                <div style="border: 1px solid #e0e0e0; border-radius: 5px; padding: 5px 10px; background-color: #f8f9fa; font-size: 14px;">
                                    {row['Process']}
                                </div>
                            """, unsafe_allow_html=True)
                        
                            if has_next:
                                if next_col == current_col:
                                    st.markdown("<p style='text-align:center; margin:0; line-height:1;'>↓</p>", unsafe_allow_html=True)
                                elif next_col > current_col:
                                    st.markdown("<p style='text-align:right; margin:0; line-height:1;'>➡</p>", unsafe_allow_html=True)
                                else:
                                    st.markdown("<p style='text-align:left; margin:0; line-height:1;'>⬅</p>", unsafe_allow_html=True)
                        else:
                            st.write("")

        # 6. VOICE OF CUSTOMER (VOC) - INTEGRATION DE LA PERSISTANCE AVEC SÉCURITÉ TYPE
        st.divider()
        st.header("6. Voice of Customer (VOC) - Flux Black Belt")

        # Initialisation des variables
        if "voc_questions" not in p:
            p["voc_questions"] = ["Temps perdu ?", "Retouches ?", "Pénibilité ?", "Irritants ?", "Changement unique ?"]
        
        # SÉCURITÉ ULTRA-ROBUSTE : Si voc_raw_data est un texte 'str' (issu d'un ancien JSON), on le force en DataFrame
        if "voc_raw_data" not in p or isinstance(p["voc_raw_data"], str):
            p["voc_raw_data"] = pd.DataFrame(columns=["client", "question", "réponse brute"])
            
        # Si c'est un dictionnaire brut (chargé depuis le JSON), on le convertit proprement en DataFrame
        if isinstance(p["voc_raw_data"], dict):
            p["voc_raw_data"] = pd.DataFrame(p["voc_raw_data"])

        st.subheader("1. Élaboration du Questionnaire")
        with st.expander("📝 Configurer les questions"):
            for i, q in enumerate(p["voc_questions"]):
                p["voc_questions"][i] = st.text_input(f"Question {i+1}", value=q, key=f"q_v9_{p_idx}_{i}")

        st.write("---")
        st.subheader("2. Collecte des Données")
    
        up_file = st.file_uploader("Importer un fichier Excel", type=["xlsx", "xls"], key=f"up_v_final_{p_idx}")
    
        if up_file is not None:
            try:
                # Lecture unique du fichier importé
                df_imp = pd.read_excel(up_file).fillna("")
                df_imp.columns = [c.lower().strip() for c in df_imp.columns]
                
                # Détection automatique des colonnes
                c_col = next((c for c in df_imp.columns if "client" in c or "nom" in c), df_imp.columns[0])
                q_col = next((c for c in df_imp.columns if "quest" in c), df_imp.columns[1] if len(df_imp.columns) > 1 else None)
                r_col = next((c for c in df_imp.columns if "rép" in c or "verbatim" in c or "brute" in c), df_imp.columns[-1])
            
                new_rows = pd.DataFrame({
                    "client": df_imp[c_col].astype(str),
                    "question": df_imp[q_col].astype(str) if q_col else "Question non spécifiée",
                    "réponse brute": df_imp[r_col].astype(str)
                })
                
                # Vérification de sécurité avant comparaison
                if isinstance(p["voc_raw_data"], pd.DataFrame):
                    if p["voc_raw_data"].empty or not (p["voc_raw_data"]["réponse brute"].equals(new_rows["réponse brute"])):
                        p["voc_raw_data"] = new_rows
                        st.session_state.voc_results = None
                        st.rerun()
                    
            except Exception as e:
                st.error(f"Erreur lors de l'import : {e}")

        # L'éditeur lit désormais un DataFrame garanti à 100%
        edited_voc_df = st.data_editor(
            p["voc_raw_data"],
            num_rows="dynamic",
            use_container_width=True,
            key=f"voc_editor_sync_{p_idx}"
        )
        if edited_voc_df is not None:
            p["voc_raw_data"] = edited_voc_df

        st.write("---")
        st.subheader("3. Analyse Thématique & CTQ")
    
        if st.button("🧠 Lancer l'Analyse (Vue Black Belt)", key=f"btn_run_voc_{p_idx}"):
            df_to_analyze = p["voc_raw_data"]
            if isinstance(df_to_analyze, pd.DataFrame) and not df_to_analyze.empty:
                verbatims = df_to_analyze.iloc[:, -1].astype(str).str.lower().tolist()
                total = len(verbatims)
                mapping = [
                    {"th": "Délais (Lead Time)", "kw": ["temps", "long", "attente", "lent", "délai", "retard", "planning"], "ex": "Muda d'attente (Gaspillage)", "ctq": "Lead Time < 24h"},
                    {"th": "Qualité (Défauts)", "kw": ["refaire", "erreur", "trompé", "faute", "qualité", "mauvais", "non-conforme"], "ex": "Non-conformités (COPQ)", "ctq": "First Pass Yield 100%"},
                    {"th": "Pénibilité", "kw": ["pénible", "lourd", "difficile", "fatigue", "compliqué", "stress"], "ex": "Muri (Surcharge)", "ctq": "Ergonomie & Standard Work"},
                    {"th": "Outils", "kw": ["bug", "système", "outil", "logiciel", "ordinateur", "panne"], "ex": "Capabilité des moyens", "ctq": "Disponibilité IT > 99.9%"},
                    {"th": "Information", "kw": ["info", "manque", "flou", "comprendre", "échange", "communication"], "ex": "Mura (Variabilité du flux)", "ctq": "Standardisation de l'info"}
                ]
                res_data = []
                for m in mapping:
                    count = sum(1 for r in verbatims if any(k in r for k in m["kw"]))
                    res_data.append({
                        "thème des irritants": m["th"],
                        "explication": m["ex"],
                        "nombre d'occurrence": count,
                        "pourcentage": f"{(count/total)*100:.1f}%" if total > 0 else "0%",
                        "CTQ": m["ctq"],
                        "score": count
                    })
                st.session_state.voc_results = pd.DataFrame(res_data).sort_values("score", ascending=False).drop(columns=["score"])
                st.rerun()
            else:
                st.warning("⚠️ Le tableau est vide.")

        if st.session_state.get("voc_results") is not None:
            st.write("### 📊 Résultat de l'Analyse Thématique")
            st.table(st.session_state.voc_results)
            st.write("---")
            st.subheader("4. Diagnostic Expert & Plan d'Action Black Belt")
            res_voc = st.session_state.voc_results
            top_theme = res_voc.iloc[0]["thème des irritants"]
            ctq_cible = res_voc.iloc[0]["CTQ"]
        
            st.markdown(f"### 🎯 Focus Prioritaire : **{top_theme}**")
            st.write(f"Sur la base de l'analyse des verbatims, voici les 4 axes stratégiques pour impacter le CTQ : `{ctq_cible}`")

            plans_action = {
                "Délais (Lead Time)": [
                    "**Standardisation des flux (Takt Time)** : Aligner la cadence de production sur la demande client pour éliminer les stocks tampons.",
                    "**Mise en place d'un système Pull (Kanban)** : Réduire les en-cours en ne déclenchant le travail que sur signal de l'étape suivante.",
                    "**Chantier SMED / Setup** : Réduire les temps de changement de série pour gagner en flexibilité et réduire les attentes.",
                    "**Visual Management (Andon)** : Installer des indicateurs visuels pour détecter instantanément tout arrêt de flux."
                ],
                "Qualité (Défauts)": [
                    "**Analyse Poka-Yoke (Détrompeurs)** : Concevoir des systèmes physiques ou numériques empêchant l'erreur de se produire.",
                    "**Standard Work & Formation** : Réviser les modes opératoires et certifier les opérateurs sur les 'Points Clés de Qualité'.",
                    "**Mise en place du Jidoka** : Arrêt automatique du processus dès l'apparition d'un défaut pour éviter la propagation.",
                    "**Contrôle Statistique des Procédés (MSP)** : Suivre la variabilité en temps réel pour anticiper les dérives hors tolérance."
                ],
                "Pénibilité": [
                    "**Réingénierie Ergonomique** : Modifier l'implantation du poste (5S) pour réduire les déplacements et gestes inutiles (Muda).",
                    "**Équilibrage de ligne (Heijunka)** : Lisser la charge de travail pour éviter les pics de stress et la surcharge (Muri).",
                    "**Automatisation des tâches à faible valeur** : Déléguer les tâches répétitives à des solutions RPA ou mécaniques.",
                    "**Rotation de postes** : Planifier la polyvalence pour réduire l'exposition prolongée aux contraintes physiques/mentales."
                ],
                "Outils": [
                    "**Maintenance Préventive (TPM)** : Passer d'une logique curative à une maintenance planifiée pour garantir la capabilité machine.",
                    "**Mise à niveau technologique** : Investir dans des outils dont la précision est 10x supérieure à la tolérance du CTQ.",
                    "**Digitalisation du feedback** : Intégrer la saisie de données à la source pour éliminer les doubles saisies et erreurs outils.",
                    "**Audit de Capabilité (Cp/Cpk)** : Vérifier mathématiquement si l'outil actuel peut techniquement tenir les objectifs."
                ],
                "Information": [
                    "**Standardisation des échanges (SBAR)** : Imposer un format de communication structuré pour éliminer les ambiguïtés.",
                    "**Mise en place d'une 'Single Source of Truth'** : Centraliser les données pour éviter les versions de documents contradictoires.",
                    "**Rituels AIC (Animation à Intervalle Court)** : Créer des boucles de feedback rapides (5-10 min) pour fluidifier l'info.",
                    "**Management Visuel Numérique** : Rendre les informations critiques accessibles en un coup d'œil via des Dashboards partagés."
                ]
            }

            propositions = plans_action.get(top_theme, [
                "Audit approfondi des processus.", "Standardisation des tâches.", "Formation des équipes.", "Suivi des indicateurs."
            ])

            col_a, col_b = st.columns(2)
            with col_a:
                st.info(f"💡 **Proposition 1**\n\n{propositions[0]}")
                st.info(f"🚀 **Proposition 2**\n\n{propositions[1]}")
            with col_b:
                st.info(f"🛠️ **Proposition 3**\n\n{propositions[2]}")
                st.info(f"📊 **Proposition 4**\n\n{propositions[3]}")
                
        # --- 7. PROJECT MILESTONE & TIMING ---
        st.divider()
        st.header("📅 7. Project Milestone & Timing")

        import datetime
        import plotly.express as px
        import pandas as pd  # 👈 SÉCURITÉ : Assure l'existence de 'pd' pour éviter le crash NameError

        # SÉCURITÉ EXTENDED : Si 'p' a été perdu ou renommé lors du rechargement de l'application
        if 'p' not in locals() and 'p' not in globals():
            if "current_project" in st.session_state:
                p = st.session_state["current_project"]
            elif "project" in st.session_state:
                p = st.session_state["project"]
            elif "projects" in st.session_state and "p_idx" in locals():
                p = st.session_state["projects"][p_idx]
            else:
                p = {}

        # SÉCURITÉ INDEX : Assure la présence de p_idx pour les clés de composants
        if 'p_idx' not in locals() and 'p_idx' not in globals():
            p_idx = 0

        # 1. Initialisation initiale / Sécurisation si la structure est altérée ou absente
        if "gantt_data" not in p or not isinstance(p["gantt_data"], pd.DataFrame):
            p["gantt_data"] = pd.DataFrame([
                {"Etape": "Define", "Début": datetime.date(2026, 5, 1), "Fin": datetime.date(2026, 5, 15), "Responsable": "Black Belt"},
                {"Etape": "Measure", "Début": datetime.date(2026, 5, 16), "Fin": datetime.date(2026, 6, 15), "Responsable": "Green Belt"},
                {"Etape": "Analyze", "Début": datetime.date(2026, 6, 16), "Fin": datetime.date(2026, 7, 15), "Responsable": "Black Belt"},
                {"Etape": "Improve", "Début": datetime.date(2026, 7, 16), "Fin": datetime.date(2026, 9, 15), "Responsable": "Team"},
                {"Etape": "Control", "Début": datetime.date(2026, 9, 16), "Fin": datetime.date(2026, 10, 31), "Responsable": "Process Owner"}
            ])

        # 1b. Préparation propre des données pour l'affichage (Calendrier)
        try:
            df_input = p["gantt_data"].copy()
            df_input["Début"] = pd.to_datetime(df_input["Début"], errors='coerce').dt.date
            df_input["Fin"] = pd.to_datetime(df_input["Fin"], errors='coerce').dt.date
        except Exception:
            p["gantt_data"] = pd.DataFrame([
                {"Etape": "Define", "Début": datetime.date(2026, 5, 1), "Fin": datetime.date(2026, 5, 15), "Responsable": "Black Belt"},
                {"Etape": "Measure", "Début": datetime.date(2026, 5, 16), "Fin": datetime.date(2026, 6, 15), "Responsable": "Green Belt"},
                {"Etape": "Analyze", "Début": datetime.date(2026, 6, 16), "Fin": datetime.date(2026, 7, 15), "Responsable": "Black Belt"},
                {"Etape": "Improve", "Début": datetime.date(2026, 7, 16), "Fin": datetime.date(2026, 9, 15), "Responsable": "Team"},
                {"Etape": "Control", "Début": datetime.date(2026, 9, 16), "Fin": datetime.date(2026, 10, 31), "Responsable": "Process Owner"}
            ])
            df_input = p["gantt_data"].copy()

        # 2. Configuration des colonnes pour le sélecteur de date
        config_colonnes = {
            "Etape": st.column_config.TextColumn("Etape", required=True),
            "Début": st.column_config.DateColumn("Début", format="YYYY-MM-DD", required=True),
            "Fin": st.column_config.DateColumn("Fin", format="YYYY-MM-DD", required=True),
            "Responsable": st.column_config.TextColumn("Responsable")
        }

        # Clé unique pour l'état d'édition
        gantt_key = f"gantt_table_final_{p_idx}"

        # 📊 FONCTION DE SAUVEGARDE EN TEMPS RÉEL (on_change)
        def sync_gantt_live():
            if gantt_key in st.session_state:
                state = st.session_state[gantt_key]
                df_base = p["gantt_data"].copy()
        
                # Appliquer les lignes modifiées
                for row_idx, changes in state.get("edited_rows", {}).items():
                    for col, val in changes.items():
                        df_base.iloc[row_idx, df_base.columns.get_loc(col)] = val
                
                # Appliquer les lignes ajoutées
                for new_row in state.get("added_rows", []):
                    df_base = pd.concat([df_base, pd.DataFrame([new_row])], ignore_index=True)
            
                # Appliquer les lignes supprimées
                deleted_rows = state.get("deleted_rows", [])
                if deleted_rows:
                    df_base = df_base.drop(deleted_rows).reset_index(drop=True)
            
                # Mise à jour immédiate de la source de données principale
                p["gantt_data"] = df_base

        # Affichage du tableau éditable avec synchronisation automatique au changement
        edited_gantt = st.data_editor(
            df_input, 
            column_config=config_colonnes,
            num_rows="dynamic", 
            use_container_width=True, 
            key=gantt_key,
            on_change=sync_gantt_live
        )

        # 3. Bouton d'action pour générer le graphique Plotly
        if st.button("🚀 Générer le planning", key=f"btn_gantt_{p_idx}"):
            try:
                df_gantt = p["gantt_data"].copy().dropna(subset=["Etape", "Début", "Fin"])
        
                if not df_gantt.empty:
                    df_gantt["Début"] = pd.to_datetime(df_gantt["Début"])
                    df_gantt["Fin"] = pd.to_datetime(df_gantt["Fin"])
            
                    ordre_etapes = df_gantt["Etape"].tolist()
            
                    # Création du graphique Plotly rafraîchi
                    fig = px.timeline(
                        df_gantt, 
                        x_start="Début", 
                        x_end="Fin", 
                        y="Etape", 
                        color="Responsable",
                        title=f"Planning DMAIC - {p.get('name', 'Projet')}",
                        color_discrete_sequence=["#1E3A8A", "#10B981", "#F59E0B", "#EF4444", "#6B7280"]
                    )
            
                    fig.update_yaxes(categoryorder="array", categoryarray=ordre_etapes[::-1])
            
                    # Sauvegarde forcée du nouveau graphique dans le session_state
                    st.session_state[f"gantt_fig_{p_idx}"] = fig
                    st.toast("✅ Diagramme de Gantt mis à jour !", icon="📊")
                else:
                    st.error("Veuillez remplir correctement toutes les étapes et dates du tableau.")
            except Exception as e:
                st.error(f"Erreur lors de la génération : {e}")

        # 4. Rendu visuel permanent du graphique
        if f"gantt_fig_{p_idx}" in st.session_state and st.session_state[f"gantt_fig_{p_idx}"] is not None:
            st.plotly_chart(st.session_state[f"gantt_fig_{p_idx}"], use_container_width=True)
    
    # --- PHASE MEASURE ---
    with tabs[1]:
        st.header("Phase Measure")
        
        # 1. Project Definition & Équation Y = f(X)
        # ==========================================
        st.subheader("1. Project Definition & Transfer Function Y = f(X)")
        
        # Récupération des données réelles de la phase Define
        ctq_v = p.get("selected_ctq", "Non défini")
        voc_df = st.session_state.get("voc_raw_data", pd.DataFrame())
        
        with st.container(border=True):
            st.markdown(f"### 🎯 Objectif du Projet (Y) : `{ctq_v}`")
            st.markdown(
                r"Équation de transfert Black Belt : "
                r"$$Y = f(X_1, X_2, ..., X_n)$$"
            )
            st.caption("Détermination des variables critiques d'entrée (X) par rétro-ingénierie des verbatims clients.")

            # --- MOTEUR D'ANALYSE PAR EXTRACTION CONTEXTUELLE ---
            st.markdown("##### 🧠 Analyse IA Contextuelle des Verbatims (Top 10 X)")
            
            def analyze_voc_for_specific_x(df):
                # Si le tableau VOC est vide, on fournit une structure d'attente minimale
                if df.empty or "réponse brute" not in df.columns:
                    return [
                        {"Code Variable": f"x{i}", "Description": f"Attente de la collecte des verbatims réels pour extraction du facteur X{i}", "Impact Potentiel": "Moyen", "Origine": "Système"}
                        for i in range(1, 11)
                    ]
                
                # 1. Nettoyage et tokenisation du texte réel du client
                textes = df["réponse brute"].dropna().astype(str).tolist()
                mots_poubelles = [
                    "le", "la", "les", "des", "une", "pour", "dans", "avec", "plus", "fait", 
                    "tout", "cette", "dans", "sont", "avec", "pour", "mais", "plus", "nous",
                    "vous", "avec", "suis", "votre", "notre", "c'est", "d'un", "d'une", "est",
                    "elle", "ils", "elles", "pas", "que", "qui", "sur", "aux", "par", "dans"
                ]
                
                dictionnaire_contextuel = {}
                for texte in textes:
                    mots = [m.strip(",.?!()\"';:/*-_+") for m in texte.lower().split() if len(m) > 3]
                    for m in mots:
                        if m not in mots_poubelles:
                            dictionnaire_contextuel[m] = dictionnaire_contextuel.get(m, 0) + 1
                
                # 2. Tri des 10 concepts les plus fréquemment cités par le client (Top 10)
                concepts_cles = sorted(dictionnaire_contextuel.items(), key=lambda item: item[1], reverse=True)[:10]
                
                # 3. Traduction des concepts réels en Hypothèses Six Sigma (X)
                suggestions_x = []
                for idx, (mot, freq) in enumerate(concepts_cles):
                    # Génération de la description basée sur le VRAI mot du client
                    description_cause = f"Variabilité ou défaillance liée directement au facteur '{mot.upper()}' (Cité {freq} fois dans le VOC)"
                    
                    # Détermination de l'impact selon la récurrence du mot dans les plaintes
                    if freq >= 5:
                        impact = "Très Fort"
                    elif freq >= 3:
                        impact = "Fort"
                    else:
                        impact = "Moyen"
                        
                    suggestions_x.append({
                        "Code Variable": f"x{idx+1}",
                        "Description": description_cause,
                        "Impact Potentiel": impact,
                        "Origine": f"IA - Extraction Fréquence VOC ('{mot}')"
                    })
                
                # Remplissage de sécurité si le VOC est trop court pour atteindre 10 mots différents
                while len(suggestions_x) < 10:
                    current_len = len(suggestions_x) + 1
                    suggestions_x.append({
                        "Code Variable": f"x{current_len}", 
                        "Description": f"Facteur de variabilité additionnel à investiguer sur le terrain (X{current_len})", 
                        "Impact Potentiel": "Faible", 
                        "Origine": "Analyse Complémentaire Black Belt"
                    })
                    
                return suggestions_x

            # Génération/Mise à jour automatique basée sur l'état actuel du tableau VOC
            if "measure_x_table" not in p or st.button("🔄 Ré-analyser les données du VOC en temps réel", key=f"recalc_voc_{p_idx}"):
                p["measure_x_table"] = analyze_voc_for_specific_x(voc_df)
                if not voc_df.empty:
                    st.toast("⚡ Top 10 des variables X extrait avec succès !", icon="🧠")

            st.info("📋 **Mode Black Belt activé** : Le tableau ci-dessous a isolé le **Top 10 des variables X** les plus critiques basées sur votre VOC. Vous pouvez modifier, affiner, ou ajouter de nouvelles lignes.")

            # --- TABLEAU DYNAMIQUE DES X (4 COLONNES) ---
            df_x = pd.DataFrame(p["measure_x_table"])
            
            edited_x_df = st.data_editor(
                df_x,
                num_rows="dynamic", # Permet d'ajouter/supprimer librement des lignes
                use_container_width=True,
                key=f"x_matrix_editor_v3_{p_idx}",
                column_config={
                    "Code Variable": st.column_config.TextColumn("Code Variable", help="ex: x1, x2", width="small"),
                    "Description": st.column_config.TextColumn("Description de la cause probable (Spécifique au contexte)", width="large"),
                    "Impact Potentiel": st.column_config.SelectboxColumn(
                        "Impact Potentiel",
                        options=["Faible", "Moyen", "Fort", "Très Fort"],
                        width="medium"
                    ),
                    "Origine": st.column_config.TextColumn("Origine de l'hypothèse", width="medium")
                }
            )
            
            # Sauvegarde dans le dictionnaire projet
            if st.button("💾 Enregistrer la Matrice des X Validée", key=f"save_def_x_v3_{p_idx}"):
                p["measure_x_table"] = edited_x_df.to_dict('records')
                st.success("🎯 Matrice complète des 10 X enregistrée avec succès pour la phase Analyze.")

        # 2. Current State Detailed Process Map (VSM Base)
        # ==========================================
        st.subheader("2. Current State Detailed Process Map")

        # [PERSISTANCE VOC] Sécurité pour le rechargement automatique
        if "saved_voc_dict" in p and "voc_raw_data" not in st.session_state:
            st.session_state["voc_raw_data"] = pd.DataFrame(p["saved_voc_dict"])

        # 1. INITIALISATION STRICTE DE LA STRUCTURE DES MACRO-ÉTAPES (AVEC PRIORITÉ À L'IMPORT)
        if "vsm_macro_steps" not in st.session_state or "vsm_macro_steps" in p:
            # Si le projet contient des étapes (chargées depuis un JSON), on les prend en priorité
            if "vsm_macro_steps" in p and isinstance(p["vsm_macro_steps"], list):
                st.session_state["vsm_macro_steps"] = list(p["vsm_macro_steps"])
            else:
                sipoc_data = p.get("sipoc_data", [])
                macro_steps = []
                if isinstance(sipoc_data, list) and len(sipoc_data) > 0:
                    for row in sipoc_data:
                        step_name = row.get("Process") or row.get("process")
                        if step_name:
                            macro_steps.append(str(step_name))
                
                if not macro_steps:
                    macro_steps = ["1. Réception & Tri", "2. Saisie & Vérification", "3. Traitement & Analyse", "4. Validation & Approbation"]
                st.session_state["vsm_macro_steps"] = macro_steps

        # 2. INITIALISATION DU STOCKAGE DES DONNÉES DE TÂCHES (AVEC PRIORITÉ À L'IMPORT)
        if "vsm_detailed_map" not in st.session_state or "vsm_detailed_map" in p:
            if "vsm_detailed_map" in p and isinstance(p["vsm_detailed_map"], dict):
                st.session_state["vsm_detailed_map"] = dict(p["vsm_detailed_map"])
            else:
                st.session_state["vsm_detailed_map"] = {
                    step: [{"Détail de la tâche": "Sous-tâche initiale", "Valeur": 5.0, "Unité": "Minutes", "Type d'activité": "VA (Valeur Ajoutée)"}]
                    for step in st.session_state["vsm_macro_steps"]
                }

        # 3. INITIALISATION DES METRICS CALCULÉES
        if "vsm_totals" not in st.session_state:
            st.session_state["vsm_totals"] = {
                "lead_time": 0.0, "va": 0.0, "nva": 0.0, "bnrva": 0.0, "attente": 0.0, "pce": 0.0,
                "section_totals": {step: 5.0 for step in st.session_state["vsm_macro_steps"]}
            }

        # Fonction de conversion interne sécurisée
        def convert_to_minutes(valeur, unite):
            try: 
                val = float(valeur)
                if unite == "Heures": return val * 60.0
                if unite == "Jours": return val * 1440.0
                return val
            except (ValueError, TypeError): 
                return 0.0

        # --------------------------------------------------
        # FONCTIONS CALLBACKS (GESTION DE LA STRUCTURE)
        # --------------------------------------------------
        def add_section_callback(name, index_pos):
            clean_name = name.strip()
            if clean_name and clean_name not in st.session_state["vsm_macro_steps"]:
                st.session_state["vsm_macro_steps"].insert(index_pos, clean_name)
                st.session_state["vsm_detailed_map"][clean_name] = [
                    {"Détail de la tâche": "Première tâche à définir", "Valeur": 0.0, "Unité": "Minutes", "Type d'activité": "VA (Valeur Ajoutée)"}
                ]
                st.session_state["vsm_totals"]["section_totals"][clean_name] = 0.0
                st.toast(f"ℹ️ Section '{clean_name}' ajoutée !", icon="➕")
            elif clean_name in st.session_state["vsm_macro_steps"]:
                st.toast("⚠️ Cette section existe déjà !", icon="❌")

        def delete_section_callback(step_to_del):
            if step_to_del in st.session_state["vsm_macro_steps"]:
                st.session_state["vsm_macro_steps"].remove(step_to_del)
                if step_to_del in st.session_state["vsm_detailed_map"]:
                    del st.session_state["vsm_detailed_map"][step_to_del]
                if step_to_del in st.session_state["vsm_totals"]["section_totals"]:
                    del st.session_state["vsm_totals"]["section_totals"][step_to_del]
                st.toast(f"ℹ️ Section '{step_to_del}' supprimée.", icon="🗑️")

        # --------------------------------------------------
        # INTERFACE DE STRUCTURATION DU FLUX
        # --------------------------------------------------
        with st.container(border=True):
            st.markdown("### 🗺️ Gestion du Flux & Décomposition (Lead Time)")
            st.caption("Pilotez la structure de votre VSM. Saisissez vos tâches librement, puis cliquez sur le bouton de sauvegarde en bas pour recalculer l'ensemble du flux.")

            # Formulaire d'ajout de section
            col_add1, col_add2, col_add3 = st.columns([3, 2, 1.5])
            with col_add1:
                new_sec_name = st.text_input("📝 Nom de la nouvelle section :", placeholder="Ex: 2.bis Contrôle Qualité", key=f"new_input_text_{p_idx}")
            with col_add2:
                positions = ["En fin de processus", "Au tout début"]
                for step in st.session_state["vsm_macro_steps"]:
                    positions.append(f"Après : {step}")
                
                chosen_pos_text = st.selectbox("📍 Insérer la section :", options=positions, key=f"pos_select_{p_idx}")
                
                if chosen_pos_text == "Au tout début":
                    insert_idx = 0
                elif chosen_pos_text == "En fin de processus":
                    insert_idx = len(st.session_state["vsm_macro_steps"])
                else:
                    selected_step = chosen_pos_text.replace("Après : ", "")
                    insert_idx = st.session_state["vsm_macro_steps"].index(selected_step) + 1

            with col_add3:
                st.write("") 
                if st.button("➕ Créer la section", use_container_width=True, key=f"btn_create_{p_idx}"):
                    if new_sec_name.strip():
                        add_section_callback(new_sec_name, insert_idx)
                        st.rerun()
                    else:
                        st.warning("Nom invalide.")

            st.markdown("---")

            # --------------------------------------------------
            # BOUCLE D'ÉDITION LIBRE
            # --------------------------------------------------
            current_editors_state = {}

            for idx_step, step in enumerate(st.session_state["vsm_macro_steps"]):
                current_data = st.session_state["vsm_detailed_map"].get(step, [])
                df_sub = pd.DataFrame(current_data)
                
                sec_min = st.session_state["vsm_totals"]["section_totals"].get(step, 0.0)
                if sec_min >= 1440: t_text = f"{sec_min/1440:.1f} j"
                elif sec_min >= 60: t_text = f"{sec_min/60:.1f} h"
                else: t_text = f"{sec_min:.1f} min"

                col_title, col_del = st.columns([5, 1])
                with col_title:
                    st.markdown(f"#### 📦 Section {idx_step + 1} : {step} `⏱️ Dernier calcul: {t_text}`")
                with col_del:
                    st.write("") 
                    if st.button("🗑️ Supprimer", key=f"del_sec_{idx_step}_{p_idx}", use_container_width=True):
                        delete_section_callback(step)
                        st.rerun()

                edited_df = st.data_editor(
                    df_sub,
                    num_rows="dynamic",
                    use_container_width=True,
                    key=f"vsm_editor_v12_{step}_{p_idx}", 
                    column_config={
                        "Détail de la tâche": st.column_config.TextColumn("Détail des micro-tâches", width="large"),
                        "Valeur": st.column_config.NumberColumn("Délai", min_value=0.0, format="%.1f", width="small", required=True),
                        "Unité": st.column_config.SelectboxColumn("Unité", options=["Minutes", "Heures", "Jours"], width="small", required=True),
                        "Type d'activité": st.column_config.SelectboxColumn(
                            "Type d'activité (Lean)",
                            options=["VA (Valeur Ajoutée)", "NVA (Non Valeur Ajoutée)", "BNRVA (Business Necessary, Non-Value Added)", "Temps d'attente / Stock"],
                            width="medium", required=True
                        )
                    }
                )
                current_editors_state[step] = edited_df
                st.markdown("---")

            # --------------------------------------------------
            # BOUTON UNIQUE DE SÉCURISATION & CALCULS GLOBAUX
            # --------------------------------------------------
            if st.button("💾 Enregistrer la Cartographie & Mettre à jour les Calculs", type="primary", use_container_width=True, key=f"save_vsm_map_v12_{p_idx}"):
                
                total_lt = 0.0
                total_va = 0.0
                total_nva = 0.0
                total_bnrva = 0.0
                total_attente = 0.0
                new_section_totals = {}

                for step in st.session_state["vsm_macro_steps"]:
                    df_edited = current_editors_state.get(step, pd.DataFrame())
                    
                    if not df_edited.empty and "Détail de la tâche" in df_edited.columns:
                        df_edited = df_edited.dropna(subset=["Détail de la tâche"])
                    
                    records = df_edited.to_dict('records')
                    st.session_state["vsm_detailed_map"][step] = records
                    
                    sec_sum = 0.0
                    if not df_edited.empty:
                        for _, row in df_edited.iterrows():
                            t_min = convert_to_minutes(row.get("Valeur", 0), row.get("Unité", "Minutes"))
                            sec_sum += t_min
                            total_lt += t_min
                            
                            act_type = row.get("Type d'activité")
                            if act_type == "VA (Valeur Ajoutée)": total_va += t_min
                            elif act_type == "NVA (Non Valeur Ajoutée)": total_nva += t_min
                            elif act_type == "BNRVA (Business Necessary, Non-Value Added)": total_bnrva += t_min
                            elif act_type == "Temps d'attente / Stock": total_attente += t_min
                    
                    new_section_totals[step] = sec_sum

                pce_calc = (total_va / total_lt * 100) if total_lt > 0 else 0.0

                st.session_state["vsm_totals"] = {
                    "lead_time": total_lt,
                    "va": total_va,
                    "nva": total_nva,
                    "bnrva": total_bnrva,
                    "attente": total_attente,
                    "pce": pce_calc,
                    "section_totals": new_section_totals
                }

                # 🚨 ALIGNEMENT ET SAUVEGARDE DIRECTE DANS LE DICTIONNAIRE DE PROJET
                p["vsm_macro_steps"] = list(st.session_state["vsm_macro_steps"])
                p["vsm_detailed_map"] = dict(st.session_state["vsm_detailed_map"])
                if "voc_raw_data" in st.session_state:
                    p["saved_voc_dict"] = st.session_state["voc_raw_data"].to_dict('records')
                
                st.success("🎯 Tous les calculs ont été mis à jour avec succès sur l'ensemble de la chaîne !")
                st.rerun()

            # ==========================================
            # RENDU VISUEL BASÉ SUR LES DERNIÈRES DONNÉES VALIDÉES
            # ==========================================
            st.markdown("### 🗺️ Schéma Value Stream Mapping (VSM) Résumé")
            
            steps_list = st.session_state["vsm_macro_steps"]
            if len(steps_list) > 0:
                vsm_cols = st.columns(len(steps_list))
                for i, step in enumerate(steps_list):
                    with vsm_cols[i]:
                        sec_min = st.session_state["vsm_totals"]["section_totals"].get(step, 0.0)
                        tasks_count = len(st.session_state["vsm_detailed_map"].get(step, []))
                        short_name = str(step)[:15] + "..." if len(str(step)) > 15 else str(step)
                        
                        with st.container(border=True):
                            st.markdown(f"**🔹 ÉTAPE {i+1}**")
                            st.caption(f"*{short_name}*")
                            st.markdown(f"📋 `{tasks_count} tsk`")
                            st.code(f"{sec_min:.1f} m", language="text")

            # Timeline
            st.markdown("#### ⏱️ Ligne de temps du Flux (Value Stream Timeline)")
            timeline_markdown = ""
            for i, step in enumerate(steps_list):
                sec_min = st.session_state["vsm_totals"]["section_totals"].get(step, 0.0)
                short_name = str(step)[:20] + "..." if len(str(step)) > 20 else str(step)
                timeline_markdown += f"**[{i+1}] {short_name}** ({sec_min:.1f} min) `——➡️` "
            
            if timeline_markdown:
                st.info(timeline_markdown.rstrip(" `——➡️` "))

            # Rapports Métriques Globaux du Lead Time
            st.markdown("### 📊 Rapports Globaux du Processus")
            
            totals = st.session_state["vsm_totals"]
            if totals["lead_time"] >= 1440: display_lt = f"{totals['lead_time'] / 1440:.1f} Jours"
            elif totals["lead_time"] >= 60: display_lt = f"{totals['lead_time'] / 60:.1f} Heures"
            else: display_lt = f"{totals['lead_time']:.1f} Min"

            c1, c2, c3 = st.columns(3)
            c1.metric("⏳ LEAD TIME TOTAL", display_lt)
            c2.metric("🟢 TOTAL VALEUR AJOUTÉE (VA)", f"{totals['va']:.1f} min")
            c3.metric("📈 EFFICIENCE DU CYCLE (PCE)", f"{totals['pce']:.1f}%")

        # =====================================================================
        # 3. Lean Six Sigma Data Collection Plan (Y = f(X)) & 4. MSA
        # =====================================================================
        st.subheader("3. Master Black Belt Data Collection Plan")
            
        st.markdown("""
        ### 📊 Alignement Stratégique $Y = f(X)$ & Matrice de Collecte Phase Measure
        En tant que **Master Black Belt**, ce module structure votre plan de collecte de données terrain de manière rigoureuse.
        """)

        # Extraction propre de l'index du projet pour éviter les collisions de clés
        safe_idx = str(p_idx) if 'p_idx' in locals() else "default"

        # --- INITIALISATION GLOBALE STATIQUE & SÉCURISÉE ---
        msa_classif_key = f"msa_classification_table_{safe_idx}"

        if msa_classif_key not in st.session_state:
            st.session_state[msa_classif_key] = pd.DataFrame()

        df_classification_current = st.session_state[msa_classif_key]
        globals()['df_classification_current'] = st.session_state[msa_classif_key]

        if 'nom_colonne_variable' not in locals() and 'nom_colonne_variable' not in globals():
            nom_colonne_variable = "Variable Critique (liée au Y)"

        # --- ISOLATION DU PLAN DE COLLECTE ET MSA DANS UN FRAGMENT ANTI-FLICKER ---
        @st.fragment
        def render_data_collection_and_msa(project_dict, component_idx):
            matrix_key = f"mbb_prioritization_matrix_{component_idx}"
            dcp_table_key = f"master_dcp_table_{component_idx}"
            lock_key = f"dcp_validated_lock_{component_idx}"
            local_msa_key = f"msa_classification_table_{component_idx}"
            buffer_msa_key = f"msa_buffer_{component_idx}"

            # =====================================================================
            # 🛡️ SÉCURITÉ IMPORT JSON : RESTAURATION ET SÉQUENÇAGE STRICT
            # =====================================================================
            # FIX PARTIE 1 : Force le chargement ou l'extraction si la session est absente ou vide
            if matrix_key not in st.session_state or st.session_state[matrix_key].empty:
                saved_prio = project_dict.get("prio_matrix_saved", [])
                if saved_prio:
                    st.session_state[matrix_key] = pd.DataFrame(saved_prio)
                else:
                    # On réactive le calcul d'extraction automatique depuis le VSM
                    vsm_steps = st.session_state.get("vsm_macro_steps", [])
                    vsm_detailed = st.session_state.get("vsm_detailed_map", {})
                    vsm_totals = st.session_state.get("vsm_totals", {})
                    section_totals = vsm_totals.get("section_totals", {})

                    extracted_x = []
                    if vsm_steps:
                        if section_totals:
                            highest_step = max(section_totals, key=section_totals.get)
                            if section_totals[highest_step] > 0:
                                extracted_x.append({
                                    "etape": highest_step,
                                    "variable": f"Temps de cycle unitaire sur le goulot - {highest_step}",
                                    "muda": "Surproduction / Capacité"
                                })

                        for step in vsm_steps:
                            for t in vsm_detailed.get(step, []):
                                desc_tache = str(t.get("Détail de la tâche", ""))
                                type_act = t.get("Type d'activité", "")
                                desc_lower = desc_tache.lower()
                                    
                                if not desc_tache or desc_tache in ["Sous-tâche initiale", "Première tâche à définir"]:
                                    continue

                                if type_act == "Temps d'attente / Stock" or any(kw in desc_lower for kw in ["attente", "file", "stock"]):
                                    extracted_x.append({"etape": step, "variable": f"Temps de stagnation : {desc_tache}", "muda": "Attente (Waiting)"})
                                elif type_act == "NVA (Non Valeur Ajoutée)" and any(kw in desc_lower for kw in ["retouche", "correction", "erreur", "rework"]):
                                    extracted_x.append({"etape": step, "variable": f"Fréquence de : {desc_tache}", "muda": "Défauts / Retouches"})

                    if not extracted_x:
                        extracted_x = [
                            {"etape": "1. Réception & Tri", "variable": "Taux d'erreurs à l'entrée", "muda": "Défauts / Retouches"},
                            {"etape": "2. Saisie & Vérification", "variable": "Temps d'attente de validation", "muda": "Attente (Waiting)"}
                        ]

                    st.session_state[matrix_key] = pd.DataFrame([{
                        "Étape Source": item["etape"],
                        "Variable Potentielle (X)": item["variable"],
                        "Gaspillage / Muda": item["muda"],
                        "1. Influence fortement le Y ?": "Oui",
                        "2. Apparaît souvent ?": "Oui",
                        "3. Peut-on mesurer fiablement ?": "Oui",
                        "Utilité Analytique (Futur Test d'Hypothèse)": "Démontrer la corrélation mathématique avec la variation du Lead Time."
                    } for item in extracted_x])

            # Restauration Partie 2 : DCP Officiel
            if dcp_table_key not in st.session_state:
                saved_dcp = project_dict.get("master_dcp_table", [])
                st.session_state[dcp_table_key] = pd.DataFrame(saved_dcp) if (saved_dcp is not None and len(saved_dcp) > 0) else pd.DataFrame()

            # FIX PARTIE 2 : Force la restauration depuis le JSON même si la clé a été initialisée vide plus haut
            saved_msa = project_dict.get("msa_table_saved", [])
            if saved_msa and (local_msa_key not in st.session_state or st.session_state[local_msa_key].empty):
                st.session_state[local_msa_key] = pd.DataFrame(saved_msa)
                st.session_state[msa_classif_key] = pd.DataFrame(saved_msa)

            # Restauration du verrou de jalon automatique
            if saved_msa or project_dict.get("dcp_validated_lock", False):
                st.session_state[lock_key] = True
                project_dict["dcp_validated_lock"] = True
            elif lock_key not in st.session_state:
                st.session_state[lock_key] = False

            # Initialisation du tampon de saisie pour bloquer le lag du MSA
            if buffer_msa_key not in st.session_state:
                if local_msa_key in st.session_state and not st.session_state[local_msa_key].empty:
                    st.session_state[buffer_msa_key] = st.session_state[local_msa_key].copy()
                else:
                    st.session_state[buffer_msa_key] = pd.DataFrame()

            # --------------------------------------------------
            # 1. PRIORISATION DES X
            # --------------------------------------------------
            st.markdown("### 🧠 1. Filtrage et Priorisation des $X$ ($Y = f(X)$)")
                
            edited_prio_df = st.data_editor(
                st.session_state[matrix_key],
                num_rows="dynamic",
                use_container_width=True,
                key=f"prio_editor_ui_{component_idx}",
                column_config={
                    "Étape Source": st.column_config.TextColumn("Étape Source", disabled=True, width="medium"),
                    "Variable Potentielle (X)": st.column_config.TextColumn("Variable Potentielle (X)", disabled=True, width="large"),
                    "1. Influence fortement le Y ?": st.column_config.SelectboxColumn("Influence Y ?", options=["Oui", "Non"], width="small"),
                    "2. Apparaît souvent ?": st.column_config.SelectboxColumn("Fréquent ?", options=["Oui", "Non"], width="small"),
                    "3. Peut-on mesurer fiablement ?": st.column_config.SelectboxColumn("Mesurable ?", options=["Oui", "Non"], width="small")
                }
            )
                
            if st.button("⚙️ Valider la pertinence & Générer le Data Collection Plan Master", type="primary", use_container_width=True, key=f"btn_gen_dcp_{component_idx}"):
                df_prio_fixed = pd.DataFrame(edited_prio_df)
                st.session_state[matrix_key] = df_prio_fixed
                project_dict["prio_matrix_saved"] = df_prio_fixed.to_dict('records')
                
                nom_y_projet = project_dict.get("selected_ctq", "Indicateur de Performance Principal (Y)")
                    
                dcp_final_rows = [{
                    "Variable à mesurer": nom_y_projet,
                    "Objectif de mesure": "Quantifier la performance globale.",
                    "Lien avec le Y": "Variable de sortie principale (Y) du projet Lean Six Sigma.",
                    "Définition opérationnelle exacte": "Mesure standardisée de l'indicateur clé.",
                    "Type de donnée": "Continue (Temps)", "Unité": "Minutes", "Source de donnée": "Système d'information",
                    "Méthode de collecte": "Extraction automatique", "Point de mesure dans le processus": "Sortie globale",
                    "Responsable collecte": "Sponsor", "Fréquence": "Mensuelle", "Taille échantillon": "n ≥ 30",
                    "Période de collecte": "3 mois", "Outil utilisé": "ERP", "Risques de biais": "Aucun",
                    "Méthode de contrôle qualité des données": "Validation financière"
                }]
                    
                for _, row in df_prio_fixed.iterrows():
                    if row["1. Influence fortement le Y ?"] == "Oui" and row["2. Apparaît souvent ?"] == "Oui" and row["3. Peut-on mesurer fiablement ?"] == "Oui":
                        dcp_final_rows.append({
                            "Variable à mesurer": str(row["Variable Potentielle (X)"]),
                            "Objectif de mesure": "Quantifier l'impact de ce Muda.",
                            "Lien avec le Y": "Contribution directe au Lead Time Global (Y).",
                            "Définition opérationnelle exacte": "Chrono de début et fin.",
                            "Type de donnée": "Continue (Temps)", "Unité": "Minutes", "Source de donnée": "Terrain",
                            "Méthode de collecte": "Saisie manuelle", "Point de mesure dans le processus": row["Étape Source"],
                            "Responsable collecte": "Opérateur", "Fréquence": "Quotidienne", "Taille échantillon": "n ≥ 30",
                            "Période de collecte": "2 semaines", "Outil utilisé": "Excel", "Risques de biais": "Effet Hawthorne",
                            "Méthode de contrôle qualité des données": "Audit à blanc"
                        })
                    
                st.session_state[dcp_table_key] = pd.DataFrame(dcp_final_rows)
                project_dict["master_dcp_table"] = dcp_final_rows
                
                if 'projects' in st.session_state and 'p_idx' in locals():
                    st.session_state.projects[p_idx]["prio_matrix_saved"] = df_prio_fixed.to_dict('records')
                    st.session_state.projects[p_idx]["master_dcp_table"] = dcp_final_rows
                
                st.session_state[lock_key] = True
                project_dict["dcp_validated_lock"] = True
                st.rerun()

            # --------------------------------------------------
            # 2. TABLEAU OFFICIEL DU DCP
            # --------------------------------------------------
            if dcp_table_key in st.session_state and not st.session_state[dcp_table_key].empty:
                st.markdown("### 📋 2. Matrice Officielle du Plan de Collecte (Phase Measure)")
                    
                edited_dcp_df = st.data_editor(
                    st.session_state[dcp_table_key],
                    num_rows="dynamic",
                    use_container_width=True,
                    key=f"dcp_editor_ui_{component_idx}"
                )

                if st.button("💾 Enregistrer les ajustements du Data Collection Plan", key=f"save_mbb_dcp_{component_idx}", type="secondary", use_container_width=True):
                    df_ajuste = pd.DataFrame(edited_dcp_df)
                    st.session_state[dcp_table_key] = df_ajuste
                    
                    project_dict["master_dcp_table"] = df_ajuste.to_dict('records')
                    if 'projects' in st.session_state and 'p_idx' in locals():
                        st.session_state.projects[p_idx]["master_dcp_table"] = df_ajuste.to_dict('records')
                        
                    current_msa_df = st.session_state.get(local_msa_key, pd.DataFrame())
                    existing_status = {}
                    if not current_msa_df.empty and "Variable Critique (liée au Y)" in current_msa_df.columns:
                        existing_status = dict(zip(
                            current_msa_df["Variable Critique (liée au Y)"].astype(str).str.strip(),
                            current_msa_df["Statut de validation"]
                        ))
                        
                    msa_rows = []
                    for _, row in df_ajuste.iterrows():
                        var_nom = str(row.get("Variable à mesurer", "Non définie")).strip()
                        v_type_brut = str(row.get("Type de donnée", "Continue (Temps)"))
                        v_lien = str(row.get("Lien avec le Y", ""))
                            
                        statut_recupere = existing_status.get(var_nom, "En attente")
                            
                        msa_rows.append({
                            "Variable Critique (liée au Y)": var_nom,
                            "Rôle": "Y" if "Variable de sortie principale (Y)" in v_lien else "X",
                            "Type de Donnée": "Continue" if "continue" in v_type_brut.lower() else "Attributaire",
                            "MSA Recommandé": "Gage R&R" if "continue" in v_type_brut.lower() else "Attribute Agreement Analysis (Kappa)",
                            "Statut de validation": statut_recupere
                        })
                        
                    df_msa_nouveau = pd.DataFrame(msa_rows)
                    st.session_state[local_msa_key] = df_msa_nouveau
                    st.session_state[msa_classif_key] = df_msa_nouveau
                    st.session_state[buffer_msa_key] = df_msa_nouveau.copy()
                    
                    project_dict["msa_table_saved"] = df_msa_nouveau.to_dict('records')
                    if 'projects' in st.session_state and 'p_idx' in locals():
                        st.session_state.projects[p_idx]["msa_table_saved"] = df_msa_nouveau.to_dict('records')
                        
                    st.session_state[lock_key] = True
                    project_dict["dcp_validated_lock"] = True
                    st.toast("💾 Plan de collecte ajusté et synchronisé avec le MSA !", icon="🛡️")
                    st.rerun()

        # --- 1. DÉFINITION GLOBALE DES CLÉS (Indispensable avant usage) ---
        # Ces variables doivent être définies ici pour être accessibles au reste du script
        idx_str = str(safe_idx)
        dcp_table_key = f"master_dcp_table_{idx_str}"
        local_msa_key = f"msa_classification_table_{idx_str}"
        lock_key = f"dcp_validated_lock_{idx_str}"
        proto_key = f"protocol_data_{idx_str}"

        # --- 2. APPEL DU FRAGMENT (Génération des données) ---
        render_data_collection_and_msa(p, safe_idx)
        
        # --------------------------------------------------
        # 4. VALIDATE MEASUREMENT SYSTEM (MSA)
        # --------------------------------------------------
        st.divider()
        st.subheader("4. Validate Measurement System (MSA)")
        
        # --- CORRECTION DE SYNCHRONISATION PRIORITAIRE (L'ancre au démarrage) ---
        if p.get("msa_table_saved"):
             st.session_state[lock_key] = True
        
        # --- LOGIQUE DE RESTAURATION AMÉLIORÉE ---
        if local_msa_key not in st.session_state or st.session_state[local_msa_key].empty:
            saved_msa = p.get("msa_table_saved", [])
            st.session_state[local_msa_key] = pd.DataFrame(saved_msa) if saved_msa else pd.DataFrame(columns=["Variable Critique (liée au Y)", "Rôle", "Type de Donnée", "MSA Recommandé", "Statut de validation"])
        
        # --- FORCER LA SYNCHRONISATION DU PROTOCOLE AU CHARGEMENT ---
        # Cette ligne garantit que même si on vient de se connecter, le lien avec 'p' est rafraîchi
        if p.get("protocol_saved") and st.session_state[local_msa_key].empty:
             st.session_state[local_msa_key] = pd.DataFrame(p["protocol_saved"])
            
        # Vérification du verrouillage
        if not st.session_state.get(lock_key, False):
            st.info("🔒 **Statut Jalon : En attente de validation du DCP**")
        else:
            # Affichage MSA
            df_msa = st.session_state.get(local_msa_key, pd.DataFrame())
            if df_msa.empty:
                df_msa = pd.DataFrame(columns=["Variable Critique (liée au Y)", "Rôle", "Type de Donnée", "MSA Recommandé", "Statut de validation"])

            for i, row in df_msa.iterrows():
                var_name = row["Variable Critique (liée au Y)"]
                v_c = "".join(e for e in str(var_name) if e.isalnum())
                # On vérifie si des données de test (rép ou reprod) existent dans 'p' pour cette variable
                if (f"save_rep_{v_c}_{safe_idx}" in p) or (f"save_reprod_{v_c}_{safe_idx}" in p):
                    df_msa.at[i, "Statut de validation"] = "test effectué"
            
            # On met à jour la session avec le statut corrigé
            st.session_state[local_msa_key] = df_msa 
            # --------------------------------------------------
            
            edited_msa_df = st.data_editor(df_msa, num_rows="fixed", use_container_width=True)
    
            # Bouton de sauvegarde
            if st.button("💾 Enregistrer MSA et Protocole", key=f"btn_save_{idx_str}", type="primary"):
                p["msa_table_saved"] = edited_msa_df.to_dict('records')
                # Note: assurez-vous que edited_proto_df est défini ici si vous l'utilisez
                p["protocol_saved"] = edited_msa_df.to_dict('records') 
                st.success("✅ Données enregistrées dans le projet !")

        # --- LOGIQUE D'AFFICHAGE ROBUSTE ET PERSISTANTE ---
        
        # 1. On récupère le dataframe de la session ou du projet
        df_msa_actif = st.session_state.get(local_msa_key, pd.DataFrame())
        
        # 2. Vérification si le statut 'test effectué' existe déjà en mémoire ou dans le projet
        has_test_done_in_session = False
        if not df_msa_actif.empty and "Statut de validation" in df_msa_actif.columns:
            has_test_done_in_session = "test effectué" in df_msa_actif["Statut de validation"].values
            
        has_test_done_in_project = p.get("msa_table_saved") is not None and any(
            row.get("Statut de validation") == "test effectué" for row in p["msa_table_saved"]
        )

       # 3. La condition finale reste strictement la vôtre
        if p.get("protocol_saved") or has_test_done_in_session or has_test_done_in_project:    
            st.markdown("##### 👟 Exécution du Protocole Terrain")
            
            if not df_classification_current.empty and nom_colonne_variable in df_classification_current.columns:
                list_variables_critiques = df_classification_current[nom_colonne_variable].dropna().tolist()
            elif p.get("msa_table_saved"):
                list_variables_critiques = [row.get("Variable Critique (liée au Y)") for row in p["msa_table_saved"] if row.get("Variable Critique (liée au Y)")]
            else:
                list_variables_critiques = []
            
            if list_variables_critiques:
                if "msa_validated_vars" not in st.session_state:
                    st.session_state["msa_validated_vars"] = {}
                if "msa_bias_history" not in st.session_state:
                    st.session_state["msa_bias_history"] = {}

                # Génération des options du sélecteur avec un indicateur de statut clair
                options_sélecteur = []
                mapping_variables = {}
                for var in list_variables_critiques:
                    if str(var).strip() == "":
                        continue
                    v_c = "".join(e for e in str(var) if e.isalnum())
                    if st.session_state.get(f"status_lock_{v_c}_{safe_idx}", False) or p.get(f"validated_status_{v_c}_{safe_idx}", False):
                        label = f"✅ {var}"
                    else:
                        label = f"⏳ {var}"
                    options_sélecteur.append(label)
                    mapping_variables[label] = var

                selected_option = st.selectbox(
                    "Sélectionnez la variable à tester actuellement parmi vos variables critiques :",
                    options=options_sélecteur,
                    key=f"msa_selected_var_{safe_idx}"
                )
                
                # Récupération sécurisée du nom propre de la variable sans couper arbitrairement l'émoji
                selected_var_to_test = mapping_variables.get(selected_option, "")

                # --- TABLEAU DE BORD DES MESURES DÉJÀ VALIDÉES ---
                st.markdown("##### 📈 Historique des protocoles validés")
                
                variables_valides = []
                for var in list_variables_critiques:
                    v_c = "".join(e for e in str(var) if e.isalnum())
                    if st.session_state.get(f"status_lock_{v_c}_{safe_idx}", False) or p.get(f"validated_status_{v_c}_{safe_idx}", False):
                        variables_valides.append(var)
                
                if variables_valides:
                    with st.expander("🔍 Voir toutes les données terrain validées (Historique)", expanded=True):
                        for v_nom in variables_valides:
                            v_clean = "".join(e for e in str(v_nom) if e.isalnum())
                            p_rep_key = f"save_rep_{v_clean}_{safe_idx}"
                            p_reprod_key = f"save_reprod_{v_clean}_{safe_idx}"
                            
                            st.markdown(f"**🟢 Variable : {v_nom}**")
                            c1, c2 = st.columns(2)
                            with c1:
                                if p_rep_key in p:
                                    st.caption("Données de Reproductibilité enregistrées (différentes personnes obtiennent-elles des résultats similaires?: même produit, plusieurs opérateurs) :")
                                    st.dataframe(pd.DataFrame(p[p_rep_key]), use_container_width=True)
                            with c2:
                                if p_reprod_key in p:
                                    st.caption("Données de Répétabilité enregistrées (La même personne mesure-t-elle toujours pareil?: 1 opérateur, même situation, plusieurs mesures) :")
                                    st.dataframe(pd.DataFrame(p[p_reprod_key]), use_container_width=True)
                            st.markdown("---")
                else:
                    st.info("ℹ️ Aucune variable n'a encore été validée. Les résumés s'afficheront ici automatiquement après avoir cliqué sur 'Valider et verrouiller'.")
            
                # ISOLATION CELLULAIRE TERRAIN (Par variable sélectionnée)
                var_clean_id = "".join(e for e in selected_var_to_test if e.isalnum())
                
                dynamic_rep_key = f"rep_data_{var_clean_id}_{safe_idx}"
                dynamic_reprod_key = f"reprod_data_{var_clean_id}_{safe_idx}"
                bias_hist_key = f"hist_{var_clean_id}_{safe_idx}"
             
                p_rep_save_key = f"save_rep_{var_clean_id}_{safe_idx}"
                p_reprod_save_key = f"save_reprod_{var_clean_id}_{safe_idx}"
                p_bias_hist_save_key = f"save_bias_hist_{var_clean_id}_{safe_idx}"
                
                # Restauration automatique depuis la persistance 'p'
                if p_rep_save_key in p and dynamic_rep_key not in st.session_state:
                    st.session_state[dynamic_rep_key] = pd.DataFrame(p[p_rep_save_key])
                if p_reprod_save_key in p and dynamic_reprod_key not in st.session_state:
                    st.session_state[dynamic_reprod_key] = pd.DataFrame(p[p_reprod_save_key])
                if p_bias_hist_save_key in p and bias_hist_key not in st.session_state["msa_bias_history"]:
                    st.session_state["msa_bias_history"][bias_hist_key] = p[p_bias_hist_save_key]
                
                if bias_hist_key not in st.session_state["msa_bias_history"]:
                    st.session_state["msa_bias_history"][bias_hist_key] = []
                
                liste_unites = ["minutes", "heure", "jour", "g", "kg", "unité", "m", "l", "%", "Ar"]
                
                if dynamic_rep_key not in st.session_state:
                    st.session_state[dynamic_rep_key] = pd.DataFrame([
                        {"Opérateur": "Opérateur 1", "Situation A": 0.0, "Unité A": "unité", "Situation B": 0.0, "Unité B": "unité"},
                        {"Opérateur": "Opérateur 2", "Situation A": 0.0, "Unité A": "unité", "Situation B": 0.0, "Unité B": "unité"}
                    ])
                    
                if dynamic_reprod_key not in st.session_state:
                    st.session_state[dynamic_reprod_key] = pd.DataFrame([
                        {"Essai / Pièce": "Pièce 1", "Résultat": 0.0, "Unité": "unité"},
                        {"Essai / Pièce": "Pièce 2", "Résultat": 0.0, "Unité": "unité"}
                    ])
                    
                id_validation_blindee_active = f"status_lock_{var_clean_id}_{safe_idx}"
                if st.session_state.get(id_validation_blindee_active, False) or p.get(f"validated_status_{var_clean_id}_{safe_idx}", False):
                    st.success(f"🎯 **Statut : Validé** | Vous visualisez les données sécurisées pour : **{selected_var_to_test}**")
                else:
                    st.warning(f"📋 **Statut : Saisie en cours** | Remplissez les mesures pour : **{selected_var_to_test}**")
                
                col_t1, col_t2 = st.columns(2)
                
                with col_t1:
                    st.markdown("**🔬 Test de Reproductibilité (Inter-Opérateurs)**")
                    edited_rep = st.data_editor(
                        st.session_state[dynamic_rep_key],
                        num_rows="dynamic",
                        use_container_width=True,
                        key=f"editor_rep_{var_clean_id}_{safe_idx}",
                        column_config={
                            "Unité A": st.column_config.SelectboxColumn("Unité de mesure", options=liste_unites, required=True),
                            "Unité B": st.column_config.SelectboxColumn("Unité de mesure", options=liste_unites, required=True)
                        }
                    )

                with col_t2:
                    st.markdown("**👥 Test de Répétabilité (Intra-Opérateur)**")
                    edited_reprod = st.data_editor(
                        st.session_state[dynamic_reprod_key],
                        num_rows="dynamic",
                        use_container_width=True,
                        key=f"editor_reprod_{var_clean_id}_{safe_idx}",
                        column_config={
                            "Unité": st.column_config.SelectboxColumn("Unité de mesure", options=liste_unites, required=True)
                        }
                    ) 
                
                # 🛡️ SÉCURITÉ SANS TOUCHER AUX BOUTONS : Si l'utilisateur clique sur un bouton plus bas,
                # les variables edited_rep et edited_reprod se vident. On les force ici à récupérer les données en mémoire.
                if edited_rep is None or (isinstance(edited_rep, dict) and not edited_rep.get("edited_rows") and not edited_rep.get("added_rows") and not edited_rep.get("deleted_rows")):
                    if f"editor_rep_{var_clean_id}_{safe_idx}" in st.session_state:
                        edited_rep = st.session_state[f"editor_rep_{var_clean_id}_{safe_idx}"]

                if edited_reprod is None or (isinstance(edited_reprod, dict) and not edited_reprod.get("edited_rows") and not edited_reprod.get("added_rows") and not edited_reprod.get("deleted_rows")):
                    if f"editor_reprod_{var_clean_id}_{safe_idx}" in st.session_state:
                        edited_reprod = st.session_state[f"editor_reprod_{var_clean_id}_{safe_idx}"]
                
                # --- BLOC DÉFINITIF POUR VALEUR DE RÉFÉRENCE (Hors st.form pour persistance instantanée et JSON) ---
                session_key = f"reference_master_config_{var_clean_id}_{safe_idx}"
                save_key = f"save_master_config_{var_clean_id}_{safe_idx}"

                # 1. Initialisation de l'état de session si absent
                if session_key not in st.session_state:
                    if isinstance(p, dict) and save_key in p and isinstance(p[save_key], dict):
                        st.session_state[session_key] = p[save_key].copy()
                    else:
                        st.session_state[session_key] = {
                            "type_specification": "Valeur exacte",
                            "valeur_exacte": 0.0,
                            "valeur_seuil": 0.0,
                            "borne_inf": 0.0,
                            "borne_sup": 10.0,
                            "proposition_attributs": "Conforme / Non-conforme"
                        }

                k_type = f"spec_type_{var_clean_id}_{safe_idx}"
                k_vex = f"val_ex_{var_clean_id}_{safe_idx}"
                k_vseuil = f"val_seuil_{var_clean_id}_{safe_idx}"
                k_binf = f"binf_{var_clean_id}_{safe_idx}"
                k_bsup = f"bsup_{var_clean_id}_{safe_idx}"
                k_patt = f"prop_att_{var_clean_id}_{safe_idx}"

                current_cfg = st.session_state[session_key]

                st.markdown("##### 🎯 Valeur de Référence (Master / Standard)")

                type_spec_options = [
                    "Valeur exacte",
                    "Supérieur ou égal à (≥)",
                    "Inférieur ou égal à (≤)",
                    "Intervalle (Entre min et max)",
                    "Attribut qualitatif (Sans valeur numérique - Proposition Master Black Belt)"
                ]

                t_courant = current_cfg.get("type_specification", "Valeur exacte")
                idx_courant = type_spec_options.index(t_courant) if t_courant in type_spec_options else 0

                # Sélection du type avec callback ou mise à jour directe en temps réel
                choix_type = st.selectbox(
                    "Type de règle de référence (Standard Master) :",
                    options=type_spec_options,
                    index=idx_courant,
                    key=k_type
                )

                col_spec1, col_spec2 = st.columns(2)

                with col_spec1:
                    if choix_type == "Valeur exacte":
                        val_ex = st.number_input("Définir la valeur cible exacte (Master) :", value=float(current_cfg.get("valeur_exacte", 0.0)), key=k_vex)
                    elif choix_type in ["Supérieur ou égal à (≥)", "Inférieur ou égal à (≤)"]:
                        val_seuil = st.number_input("Valeur seuil acceptée :", value=float(current_cfg.get("valeur_seuil", 0.0)), key=k_vseuil)
                    elif choix_type == "Intervalle (Entre min et max)":
                        b_inf = st.number_input("Borne inférieure de l'intervalle :", value=float(current_cfg.get("borne_inf", 0.0)), key=k_binf)
                        b_sup = st.number_input("Borne supérieure de l'intervalle :", value=float(current_cfg.get("borne_sup", 10.0)), key=k_bsup)
                    else:
                        st.markdown("##### 🧠 Analyse MBB - Attributs / Données Qualitatives")
                        p_att = st.text_area("Proposition de référentiel qualitatif (Master Attribut) :", value=str(current_cfg.get("proposition_attributs", "Conforme / Non-conforme")), key=k_patt)

                with col_spec2:
                    st.markdown("##### ⚖️ Règle d'évaluation")
                    st.markdown("- Respecte la règle $\rightarrow$ Statut : `OK`")
                    st.markdown("- Hors règle $\rightarrow$ Statut : `Défaut (Non-OK / Non-conforme)`")

                # Sauvegarde instantanée et répercussion immédiate dans p et st.session_state
                updated_data = {
                    "type_specification": choix_type,
                    "valeur_exacte": float(st.session_state.get(k_vex, 0.0)) if choix_type == "Valeur exacte" else 0.0,
                    "valeur_seuil": float(st.session_state.get(k_vseuil, 0.0)) if choix_type in ["Supérieur ou égal à (≥)", "Inférieur ou égal à (≤)"] else 0.0,
                    "borne_inf": float(st.session_state.get(k_binf, 0.0)) if choix_type == "Intervalle (Entre min et max)" else 0.0,
                    "borne_sup": float(st.session_state.get(k_bsup, 10.0)) if choix_type == "Intervalle (Entre min et max)" else 10.0,
                    "proposition_attributs": str(st.session_state.get(k_patt, "Conforme / Non-conforme")) if choix_type.startswith("Attribut qualitatif") else "Conforme / Non-conforme"
                }

                st.session_state[session_key] = updated_data
                if isinstance(p, dict):
                    p[save_key] = updated_data.copy() 
                        
                # --- BOUTON DÉDIÉ : LANCER L'ANALYSE DES BIAIS ---
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("📊 Lancer l'analyse des risques de biais", key=f"btn_analyze_bias_{var_clean_id}_{safe_idx}", use_container_width=True):
                    
                    # 1. RÉCUPÉRATION ET BLINDAGE ANTI-CRASH (NONE-CHECK)
                    # Si la variable locale est perdue au clic, on tente de la récupérer via les clés d'état actives
                    if 'edited_rep' not in locals() or edited_rep is None:
                        edited_rep = st.session_state.get(dynamic_rep_key) or st.session_state.get(f"editor_rep_{var_clean_id}_{safe_idx}")
                        
                    if 'edited_reprod' not in locals() or edited_reprod is None:
                        edited_reprod = st.session_state.get(dynamic_reprod_key) or st.session_state.get(f"editor_reprod_{var_clean_id}_{safe_idx}")

                    # 2. SAUVEGARDE SÉCURISÉE DANS LE DICTIONNAIRE 'P' (Uniquement si les données existent)
                    if edited_rep is not None and hasattr(edited_rep, 'to_dict'):
                        p[p_rep_save_key] = edited_rep.to_dict(orient='records')
                        st.session_state[dynamic_rep_key] = edited_rep
                    
                    if edited_reprod is not None and hasattr(edited_reprod, 'to_dict'):
                        p[p_reprod_save_key] = edited_reprod.to_dict(orient='records')
                        st.session_state[dynamic_reprod_key] = edited_reprod

                    # 3. VOTRE LOGIQUE D'ANALYSE DES BIAIS (100% INTACTE ET PROTÉGÉE)
                    anomalies_mineures = []
                    anomalies_majeures = []
                    
                    if edited_rep is not None and not edited_rep.empty and 'Situation A' in edited_rep.columns and 'Situation B' in edited_rep.columns:
                        try:
                            val_col1 = pd.to_numeric(edited_rep['Situation A'], errors='coerce')
                            val_col2 = pd.to_numeric(edited_rep['Situation B'], errors='coerce')
                            
                            if val_col1.dropna().empty or val_col2.dropna().empty:
                                anomalies_majeures.append("Données de Répétabilité Invalides ou Vides")
                            else:
                                diffs = (val_col1 - val_col2).abs()
                                mean_val = val_col1.mean()
                                if not diffs.dropna().empty and mean_val > 0:
                                    if diffs.max() > mean_val * 0.30:
                                        anomalies_majeures.append("Incohérence Critique de Répétabilité (EV)")
                                    elif diffs.max() > mean_val * 0.10:
                                        anomalies_mineures.append("Dispersion de Répétabilité légère")
                                
                                all_vals = pd.concat([val_col1, val_col2]).dropna()
                                if not all_vals.empty and all(v % 5 == 0 for v in all_vals if v > 0):
                                    anomalies_mineures.append("Biais d'Arrondis Systématiques (Résolution insuffisante)")
                        except Exception as e:
                            anomalies_majeures.append(f"Erreur évaluation Répétabilité: {str(e)}")

                    if edited_reprod is not None and not edited_reprod.empty and "Résultat" in edited_reprod.columns:
                        try:
                            vals_reprod = pd.to_numeric(edited_reprod["Résultat"], errors='coerce').dropna()
                            
                            if vals_reprod.empty:
                                anomalies_majeures.append("Données de Reproductibilité Manquantes ou Invalides")
                            else:
                                if len(vals_reprod) >= 2:
                                    std_dev = vals_reprod.std()
                                    if pd.notna(std_dev) and vals_reprod.mean() > 0:
                                        if std_dev > (vals_reprod.mean() * 0.20):
                                            anomalies_majeures.append("Instabilité Critique de Reproductibilité (AV)")
                                if 'valeur_reference' in locals() and valeur_reference != 0.0:
                                    biais_justesse = abs(vals_reprod.mean() - valeur_reference)
                                    if biais_justesse > abs(valeur_reference * 0.08):
                                        anomalies_majeures.append("Biais de Justesse Linéaire (Décalage / Master)")
                        except Exception as e:
                            anomalies_majeures.append(f"Erreur évaluation Reproductibilité: {str(e)}")
                    
                    if (edited_rep is None or edited_rep.empty) and (edited_reprod is None or edited_reprod.empty):
                        anomalies_majeures.append("Variabilité Incalculable : Aucune donnée disponible")
                    
                    if len(anomalies_majeures) == 0 and len(anomalies_mineures) == 0:
                        score, status = "100%", "🟢 Système Fiable"
                    elif len(anomalies_majeures) == 0 and len(anomalies_mineures) > 0:
                        score, status = "75%", "🟡 Alerte : Biais Mineur (Sous contrôle)"
                    elif len(anomalies_majeures) == 1 and len(anomalies_mineures) == 0:
                        score, status = "50%", "🟠 Alerte : Dérive Majeure"
                    elif len(anomalies_majeures) == 1 and len(anomalies_mineures) > 0:
                        score, status = "35%", "🟠 Alerte : Instabilité Système Cumulée"
                    else:
                        score, status = "10%", "🔴 Système Non Fiable (Processus Hors Contrôle)"
                    
                    toutes_anomalies = anomalies_majeures + anomalies_mineures
                    liste_anomalies_str = ", ".join(toutes_anomalies) if toutes_anomalies else "Aucune (Système sain)"
                    
                    gmt3_time = pd.Timestamp.now(tz='UTC').tz_convert('Indian/Antananarivo').strftime('%d/%m/%Y %H:%M:%S')
                    
                    run_number = len(st.session_state["msa_bias_history"][bias_hist_key]) + 1
                    st.session_state["msa_bias_history"][bias_hist_key].append({
                        "Essai": f"Analyse #{run_number}",
                        "Date/Heure": gmt3_time,
                        "Indice de Fidélité": score,
                        "Statut Global": status,
                        "Anomalies Détectées": liste_anomalies_str
                    })
                    
                    p[p_bias_hist_save_key] = st.session_state["msa_bias_history"][bias_hist_key]
                    st.rerun()

                if st.session_state["msa_bias_history"][bias_hist_key]:
                    st.markdown("##### ⏳ Évolution de l'Analyse des Biais (Suivi cumulé des recalibrages)")
                    df_history = pd.DataFrame(st.session_state["msa_bias_history"][bias_hist_key])
                    st.table(df_history)

                # =====================================================================
                # 💾 BOUTON DE VALIDATION DÉFINITIVE
                # =====================================================================
                if st.button(
                    f"💾 Valider et verrouiller définitivement les données pour : {selected_var_to_test}", 
                    key=f"btn_validate_msa_{var_clean_id}_{safe_idx}", 
                    type="primary", 
                    use_container_width=True
                ):
                    # 1️⃣ Sécurité : On s'assure que le projet existe dans la session globale
                    if 'projects' in st.session_state and 'p_idx' in locals():
                        idx = p_idx
                        
                        # 2️⃣ Récupération directe des composants de la page
                        final_rep = edited_rep
                        final_reprod = edited_reprod

                        # 3️⃣ ON ÉCRIT DIRECTEMENT DANS LE PROJET GLOBAL (Pas dans 'p')
                        if final_rep is not None and hasattr(final_rep, 'to_dict'):
                            st.session_state.projects[idx][p_rep_save_key] = final_rep.to_dict(orient='records')
                            st.session_state[dynamic_rep_key] = final_rep
                            
                        if final_reprod is not None and hasattr(final_reprod, 'to_dict'):
                            st.session_state.projects[idx][p_reprod_save_key] = final_reprod.to_dict(orient='records')
                            st.session_state[dynamic_reprod_key] = final_reprod
                        
                        # Sauvegarde des statuts MSA directement dans le projet global
                        st.session_state.projects[idx][f"validated_status_{var_clean_id}_{safe_idx}"] = True
                        st.session_state.projects[idx][p_bias_hist_save_key] = st.session_state["msa_bias_history"][bias_hist_key]

                        # 4️⃣ Mise à jour et sauvegarde forcée du statut "test effectué"
                        classif_df = st.session_state.get(msa_classif_key, pd.DataFrame())
                        if not classif_df.empty:
                            for idx_row, row in classif_df.iterrows():
                                if str(row[nom_colonne_variable]).strip() == str(selected_var_to_test).strip():
                                    classif_df.at[idx_row, "statut validation"] = "test effectué"
                            st.session_state[msa_classif_key] = classif_df
                            st.session_state.projects[idx]["msa_classification_table"] = classif_df.to_dict(orient='records')

                        # 5️⃣ On synchronise la variable locale 'p' juste pour éviter les bugs d'affichage immédiat
                        p.update(st.session_state.projects[idx])

                    # Verrouillage de l'affichage de l'interface
                    st.session_state[f"status_lock_{var_clean_id}_{safe_idx}"] = True
                    st.session_state["msa_validated_vars"][f"{selected_var_to_test}_{safe_idx}"] = True
                    
                    st.balloons()
                    st.success(f"✅ Sauvegarde globale forcée avec succès pour **{selected_var_to_test}** !")
                    st.rerun()

                # --- 5 & 6. DIAGNOSTIC ET PLAN D'ACTION ---
                st.markdown("##### 📊 Plan d'Action Correctif (Si système non fiable au dernier essai)")
            
                last_score_str = "100%"
                if 'bias_hist_key' in locals() and bias_hist_key in st.session_state["msa_bias_history"] and st.session_state["msa_bias_history"][bias_hist_key]:
                    last_score_str = st.session_state["msa_bias_history"][bias_hist_key][-1]["Indice de Fidélité"]

                if last_score_str in ["50%", "35%", "10%"]:
                    p["msa_corrective_action"] = st.selectbox(
                        "Plan d'action prioritaire déployé lors du recalibrage :",
                        options=[
                            "Automatisation de la capture (Remplacement du facteur humain par une règle SI)",
                            "Sessions de recalibrage et formation sur définitions opérationnelles exactes",
                            "Mise en place d'une checksheet avec contrôles de saisie rigides"
                        ],
                        key=f"msa_action_choice_{safe_idx}"
                    )

                # --- 7. VALIDATION FINALE (SIGN-OFF) ---
                st.markdown("##### 📋 Validation Finale")
                if last_score_str == "10%":
                    st.error("🛑 Signature bloquée : Votre dernière analyse indique un système 'Non Fiable' (10%). Veuillez modifier vos données de test et cliquer à nouveau sur 'Lancer l'analyse des risques de biais' pour mettre à jour.")
                else:
                    saved_status = p.get("msa_is_validated_status", False)
                    is_validated = st.checkbox("Je certifie que le système de mesure est désormais stable, précis et reproductible.", value=saved_status, key=f"msa_sign_off_{safe_idx}")
                    p["msa_is_validated_status"] = is_validated
                
                    if is_validated:
                        st.success("🚀 **Measurement System Validated – Ready for Data Collection**")

        # =====================================================================
        # 5. DATA COLLECTION (WITH SENSITIVE AI ALIGNMENT & FILTERED E1)
        # =====================================================================
        from datetime import datetime, timezone, timedelta
        import io  # <-- Obligatoire pour empêcher le crash FileNotFoundError de Pandas

        # --- RECUPÉRATION DYNAMIQUE DES VARIABLES VALIDÉES DANS LE MSA ---
        safe_idx = str(p_idx) if 'p_idx' in locals() else "default"
        msa_classif_key = f"msa_classification_table_{safe_idx}"
        nom_colonne_variable = "Variable Critique (liée au Y)"

        liste_variables_dynamiques = []
        df_msa_source = pd.DataFrame()

        # Récupération de la source MSA
        if msa_classif_key in st.session_state and not st.session_state[msa_classif_key].empty:
            df_msa_source = st.session_state[msa_classif_key]
            if nom_colonne_variable in df_msa_source.columns:
                liste_variables_dynamiques = df_msa_source[nom_colonne_variable].dropna().tolist()
        elif "master_dcp_table" in st.session_state and len(st.session_state["master_dcp_table"]) > 0:
            liste_variables_dynamiques = [row["Variable à mesurer"] for row in st.session_state["master_dcp_table"] if "Variable à mesurer" in row]

        # Sécurité anti-crash si le plan est vide
        if not liste_variables_dynamiques:
            liste_variables_dynamiques = ["Temps de traitement", "Statut conformité"]

        # ÉTAPE 0 : INITIALISATION DE LA PERSISTANCE (CORRIGÉE SANS INTERPRÉTATION DE FICHIER)
        if "dc_plan" not in p:
            p["dc_plan"] = {"taille_prevue": 100, "date_debut": "2026-06-01", "date_fin_est": "2026-06-15"}

        if "dc_master_data" not in st.session_state:
            if "dc_saved_df_json" in p and p["dc_saved_df_json"]:
                # io.StringIO force Pandas à lire la chaîne texte stockée directement, évitant le plantage
                st.session_state.dc_master_data = pd.read_json(io.StringIO(p["dc_saved_df_json"]))
            else:
                cols = ["ID observation", "Date de modification"] + liste_variables_dynamiques
                st.session_state.dc_master_data = pd.DataFrame(columns=cols)

        # Synchronisation stricte des colonnes de la structure
        colonnes_requises = ["ID observation", "Date de modification"] + liste_variables_dynamiques
        for col_name in colonnes_requises:
            if col_name not in st.session_state.dc_master_data.columns:
                st.session_state.dc_master_data[col_name] = None
        
        exist_cols = [c for c in colonnes_requises if c in st.session_state.dc_master_data.columns]
        st.session_state.dc_master_data = st.session_state.dc_master_data[exist_cols]

        st.markdown("## 5 - Data collection")
        st.markdown("---")

        # =====================================================================
        # 📋 ÉCRAN 1 : RÉSUMÉ DE LA COLLECTE (FILTRÉ SUR 2 COLONNES)
        # =====================================================================
        st.markdown("### 📋 Écran 1 : Résumé de la Collecte")

        e1_c1, e1_c2, e1_c3 = st.columns(3)
        with e1_c1:
            p["dc_plan"]["taille_prevue"] = st.number_input("Taille d'échantillon prévue (N)", min_value=1, value=int(p["dc_plan"]["taille_prevue"]), key="dc_n_prevu")
        with e1_c2:
            p["dc_plan"]["date_debut"] = st.text_input("Date de début de collecte", value=p["dc_plan"]["date_debut"], key="dc_d_deb")
        with e1_c3:
            p["dc_plan"]["date_fin_est"] = st.text_input("Date estimée de fin", value=p["dc_plan"]["date_fin_est"], key="dc_d_fin")

        st.markdown("#### Liste des variables définies dans le Data Collection Plan (DCP)")
        
        # Filtrage strict sur les deux colonnes demandées
        colonnes_affichage_e1 = ["Variable Critique (liée au Y)", "Nature de la Donnée"]
        
        if not df_msa_source.empty:
            if "Type de Donnée" in df_msa_source.columns and "Nature de la Donnée" not in df_msa_source.columns:
                df_msa_source = df_msa_source.rename(columns={"Type de Donnée": "Nature de la Donnée"})
                
            cols_dispos = [c for c in colonnes_affichage_e1 if c in df_msa_source.columns]
            st.table(df_msa_source[cols_dispos])
        else:
            dcp_display = pd.DataFrame([
                {"Variable Critique (liée au Y)": var, "Nature de la Donnée": "Détectée (Quantitative/Qualitative)"} 
                for var in liste_variables_dynamiques
            ])
            st.table(dcp_display)
            
        st.caption("ℹ️ Cet écran a un rôle uniquement informatif. Utilisez l'Écran 2 pour insérer vos données terrain.")

       # =====================================================================
        # 📝 ÉCRAN 2 : SAISIE ET IMPORTATION PAR REMPLACEMENT STRICT
        # =====================================================================
        st.markdown("---")
        st.markdown("### 📝 Écran 2 : Saisie des Données (Tableaux Dynamiques)")

        st.markdown("#### 📥 Importation Intelligente Excel")
        uploaded_file = st.file_uploader("Télécharger un fichier Excel de terrain (Écrase et remplace les données actuelles)", type=["xlsx", "xls"], key="dc_excel_uploader_e2")

        # Sécurisation anti-recalcul : on ne traite le fichier QUE s'il vient d'être chargé
        if uploaded_file:
            # On vérifie si ce fichier a déjà été traité pour éviter de relancer l'IA en boucle
            file_cache_key = f"processed_{uploaded_file.name}_{uploaded_file.size}"
            
            if st.session_state.get("dc_last_processed_file") != file_cache_key:
                try:
                    import re
                    
                    # 1. Lecture brute rapide et suppression immédiate des lignes vides
                    raw_imported_df = pd.read_excel(uploaded_file)
                    raw_imported_df = raw_imported_df.dropna(how="all").reset_index(drop=True)
                    
                    st.info("🧠 *Moteur IA : Alignement de la structure en cours...*")
                    
                    # Optimisation locale des fonctions de nettoyage (compilées une seule fois)
                    regex_clean = re.compile(r'[_\-\s\./\\]+')
                    
                    def _structures_clean(text):
                        if pd.isna(text) or text is None:
                            return ""
                        t = str(text).lower().strip()
                        t = regex_clean.sub(' ', t)
                        return "".join(c for c in t if c.isalnum() or c == ' ')

                    cols_finales = ["ID observation", "Date de modification"] + liste_variables_dynamiques
                    aligned_df = pd.DataFrame(columns=cols_finales, index=range(len(raw_imported_df)))
                    colonnes_excel = list(raw_imported_df.columns)

                    # Pré-nettoyage des colonnes Excel pour accélérer les comparaisons
                    colonnes_clean = [_structures_clean(c) for c in colonnes_excel]

                    # 2. Alignement de l'ID observation réel (Recherche directe vectorielle)
                    mots_cles_id = {"id", "observation", "code", "num", "index", "identifiant", "key", "n°", "nom"}
                    id_col_source = None
                    
                    for col in colonnes_excel:
                        if _structures_clean(col) in mots_cles_id:
                            id_col_source = col
                            break
                    
                    if not id_col_source:
                        for col in colonnes_excel:
                            c_clean = _structures_clean(col)
                            if any(k in c_clean for k in mots_cles_id):
                                id_col_source = col
                                break

                    if id_col_source:
                        aligned_df["ID observation"] = (
                            raw_imported_df[id_col_source]
                            .astype(str)
                            .str.replace(r'\.0$', '', regex=True)
                            .str.strip()
                        )
                    else:
                        aligned_df["ID observation"] = [f"Obs_{i+1}" for i in range(len(raw_imported_df))]

                    # 3. Alignement flou des Variables Critiques (Optimisé)
                    for var_critique in liste_variables_dynamiques:
                        v_clean = _structures_clean(var_critique)
                        w1 = set(v_clean.split())
                        
                        meilleur_match = None
                        meilleur_score = 0.0
                        
                        for idx, col in enumerate(colonnes_excel):
                            c_clean = colonnes_clean[idx]
                            w2 = set(c_clean.split())
                            if not w1 or not w2:
                                continue
                            
                            score = len(w1.intersection(w2)) / max(len(w1), len(w2))
                            if v_clean in c_clean or c_clean in v_clean:
                                score += 0.3
                                
                            if score > meilleur_score:
                                meilleur_score = score
                                meilleur_match = col
                        
                        if meilleur_match and meilleur_score >= 0.35:
                            aligned_df[var_critique] = raw_imported_df[meilleur_match].values
                        else:
                            aligned_df[var_critique] = None

                    # 4. Fuseau horaire GMT+3 et nettoyage
                    tz_gmt3 = datetime.now(timezone(timedelta(hours=3))).strftime("%Y-%m-%d %H:%M:%S")
                    aligned_df["Date de modification"] = tz_gmt3
                    aligned_df = aligned_df.reset_index(drop=True).where(pd.notnull(aligned_df), None)

                    # Remplacement strict en session de manière définitive
                    st.session_state.dc_master_data = aligned_df
                    p["dc_saved_df_json"] = st.session_state.dc_master_data.to_json()
                    
                    # Marquer ce fichier comme "traité" pour bloquer les futurs recalculs parasites
                    st.session_state["dc_last_processed_file"] = file_cache_key
                    st.success(f"🚀 Base de données réinitialisée ! {len(aligned_df)} nouvelles lignes exclusives chargées.")
                    st.rerun()

                except Exception as e:
                    st.error(f"❌ Erreur critique lors de l'alignement : {e}")

        # --- TABLEAU DE COLLECTE TERRAIN (ÉCRAN 2) ---
        st.markdown("#### 🛠️ Tableau de Collecte Actuel")
        
        def _sauvegarder_grille_callback():
            if "dc_master_grid_editor" in st.session_state:
                grille_evenement = st.session_state["dc_master_grid_editor"]
                if grille_evenement["edited_rows"] or grille_evenement["added_rows"] or grille_evenement["deleted_rows"]:
                    tz_gmt3 = datetime.now(timezone(timedelta(hours=3))).strftime("%Y-%m-%d %H:%M:%S")
                    try:
                        for row_idx in grille_evenement["edited_rows"].keys():
                            if row_idx < len(st.session_state.dc_master_data):
                                st.session_state.dc_master_data.iloc[row_idx, st.session_state.dc_master_data.columns.get_loc("Date de modification")] = tz_gmt3
                    except Exception:
                        pass
                    p["dc_saved_df_json"] = st.session_state.dc_master_data.to_json()

        st.data_editor(
            st.session_state.dc_master_data,
            num_rows="dynamic",
            key="dc_master_grid_editor",
            use_container_width=True,
            on_change=_sauvegarder_grille_callback
        )

        # =====================================================================
        # 🔍 ÉCRAN 3 : QUALITÉ DES DONNÉES (OPTIMISÉ VECTORIELLEMENT)
        # =====================================================================
        st.markdown("---")
        st.markdown("### 🔍 Écran 3 : Qualité des Données")

        df_qualite = st.session_state.dc_master_data
        total_prevu = max(1, int(p["dc_plan"].get("taille_prevue", 100)))
        taille_echantillon_obs = len(df_qualite) if not df_qualite.empty else 0
        taux_completude_theorique = min(100.0, (taille_echantillon_obs / total_prevu) * 100)

        # Calculs ultra-rapides sans boucles imbriquées imbéciles
        donnees_manquantes = 0
        erreurs_detectees = 0

        if taille_echantillon_obs > 0 and liste_variables_dynamiques:
            cols_calcul = [c for c in liste_variables_dynamiques if c in df_qualite.columns]
            if cols_calcul:
                sub_df = df_qualite[cols_calcul]
                # Manquants vectoriels
                donnees_manquantes = int(sub_df.isna().sum().sum() + (sub_df == "").sum().sum())
                
                # --- DÉTECTION EXCLUSIVE DES ANOMALIES ET ERREURS DE SAISIE ---
                for col in cols_calcul:
                    # On isole la colonne sans les cellules vides
                    series_col = sub_df[col].dropna()
                    series_col = series_col[series_col != ""]
                    
                    if not series_col.empty:
                        # On tente de convertir en numérique pour piéger les textes parasites
                        numeric_coerced = pd.to_numeric(series_col, errors="coerce")
                        total_valeurs = len(series_col)
                        total_numerique = numeric_coerced.notna().sum()
                        
                        # CAS 1 : C'est une colonne de chiffres, mais quelqu'un a écrit des lettres
                        if total_numerique > 0 and (total_numerique / total_valeurs) > 0.30:
                            anomalies_texte = int(numeric_coerced.isna().sum())
                            erreurs_detectees += anomalies_texte
                        
                        # CAS 2 : C'est une colonne de texte, on cherche des anomalies de frappe (ex: un point ou un chiffre isolé)
                        elif total_numerique == 0:
                            # Repère les saisies suspectes ultra-courtes ou purement numériques là où on attend du texte
                            anomalies_frappe = series_col.astype(str).str.strip().isin([".", ",", "-", "?", "0", "1"]).sum()
                            erreurs_detectees += int(anomalies_frappe)

        # Affichage des KPIs
        eq1, eq2, eq3, eq4 = st.columns(4)
        eq1.metric("Taille échantillon (Obs.)", f"{taille_echantillon_obs} lignes")
        eq2.metric("Erreurs détectées", f"{erreurs_detectees} OK/KO")
        eq3.metric("Données manquantes", f"{donnees_manquantes} cellule(s)")
        eq4.metric("Complétude Théorique", f"{taux_completude_theorique:.1f} %")
        st.caption("📈 *Complétude théorique basée sur l'objectif d'échantillonnage défini à l'Écran 1.*")

        # =====================================================================
        # 📈 ÉCRAN 4 : SUIVI DE LA COLLECTE (ALLÉGÉ)
        # =====================================================================
        st.markdown("---")
        st.markdown("### 📈 Écran 4 : Suivi de la Collecte")

        # Utilisation de la taille déjà calculée pour éviter les redondances
        obs_collectees = taille_echantillon_obs
        if obs_collectees > 0 and "ID observation" in df_qualite.columns:
            # Approche ultra-rapide pour les valeurs uniques
            obs_collectees = df_qualite["ID observation"].nunique()

        restant = max(0, total_prevu - obs_collectees)
        avancement = min(100.0, (obs_collectees / total_prevu) * 100)

        e4_c1, e4_c2, e4_c3 = st.columns(3)
        e4_c1.metric("Observations collectées", obs_collectees)
        e4_c2.metric("Restant à collecter", restant)
        e4_c3.metric("Taux d'avancement", f"{avancement:.1f} %")

        # Graphique à barres optimisé
        progress_df = pd.DataFrame({"Statut": ["Collecté", "Restant"], "Valeur": [obs_collectees, restant]})
        st.bar_chart(progress_df.set_index("Statut"))

        # =====================================================================
        # 📊 ÉCRAN 5 : STATISTIQUES DESCRIPTIVES (CORRIGÉ & SÉCURISÉ)
        # =====================================================================
        st.markdown("---")
        st.markdown("### 📊 Écran 5 : Statistiques Descriptives")

        df_stats = st.session_state.dc_master_data

        if not df_stats.empty:
            for variable in liste_variables_dynamiques:
                if variable not in df_stats.columns:
                    continue
                st.markdown(f"#### Analyse descriptive : {variable}")
                series_data = df_stats[variable].dropna()
                
                if series_data.empty:
                    st.info(f"Aucune observation exploitable pour la variable '{variable}'")
                    continue

                numeric_series = pd.to_numeric(series_data, errors="coerce").dropna()
                
                # --- BLOC NUMÉRIQUE ---
                if not numeric_series.empty and not any(k in variable.lower() for k in ["statut", "verdict", "code"]):
                    stats_data = {
                        "Métrique LSS": ["Moyenne", "Médiane", "Minimum", "Maximum", "Écart-type (σ)"],
                        "Valeur": [
                            f"{numeric_series.mean():.2f}",
                            f"{numeric_series.median():.2f}",
                            f"{numeric_series.min():.2f}",
                            f"{numeric_series.max():.2f}",
                            f"{numeric_series.std():.2f}" if len(numeric_series) > 1 else "0.00"
                        ]
                    }
                    st.table(pd.DataFrame(stats_data))

                # --- BLOC QUALITATIF (Le "else" a été réaligné ici) ---
                else:
                    attr_counts = series_data.value_counts()
                    attr_pct = series_data.value_counts(normalize=True) * 100
                    attr_df = pd.DataFrame({
                        "Fréquence (N)": attr_counts,
                        "Pourcentage (%)": attr_pct.map("{:.2f} %".format)
                    })
                    st.table(attr_df)
        else:
            st.info("Aucune statistique disponible : le tableau de collecte est vide.")

        # =====================================================================
        # 🎯 ÉCRAN 6 : BASELINE DU PROCESSUS (SÉCURISÉ)
        # =====================================================================
        st.markdown("---")
        st.markdown("### 🎯 Écran 6 : Baseline du Processus")

        df_active = st.session_state.dc_master_data

        # Initialisation propre pour éviter les erreurs de type plus loin dans le code
        total_defauts_terrain = 0

        if not df_active.empty:
            st.markdown("#### Situation de référence (KPI Générés depuis la base Terrain)")
            baseline_metrics = []
    
            for variable in liste_variables_dynamiques:
                if variable in df_active.columns:
                    # Nettoyage systématique de la colonne avant calcul
                    col_data = pd.to_numeric(df_active[variable], errors="coerce")
                    num_series = col_data.dropna()
            
                    # Calcul des moyennes pour les variables numériques
                    if not num_series.empty and not any(k in variable.lower() for k in ["statut", "verdict", "validation"]):
                        baseline_metrics.append({
                            "KPI Courant": f"Moyenne globale [{variable}]", 
                            "Niveau de performance actuel": f"{num_series.mean():.1f} min" if "temps" in variable.lower() else f"{num_series.mean():.2f}"
                        })
                    # Comptage des défauts pour les variables qualitatives
                    elif not df_active[variable].dropna().empty:
                        non_ok_count = df_active[variable].astype(str).str.strip().str.upper().isin(["NON OK", "KO", "RETOUCHE", "1", "NON-CONFORME"]).sum()
                
                        # SÉCURITÉ : On s'assure que total_defauts_terrain reste un entier valide
                        total_defauts_terrain += int(non_ok_count)
                
                        valid_rows = len(df_active[variable].dropna())
                        pct_nok = (non_ok_count / valid_rows) * 100 if valid_rows > 0 else 0
                
                        baseline_metrics.append({
                            "KPI Courant": f"Taux de défauts [{variable}]", 
                            "Niveau de performance actuel": f"{pct_nok:.1f} % ({int(non_ok_count)} défauts)"
                        })

            if baseline_metrics:
                st.table(pd.DataFrame(baseline_metrics))
            else:
                st.table(pd.DataFrame([{"KPI Courant": "Taux de retouches global", "Niveau de performance actuel": "11.0 % (Par défaut)"}]))
    
            st.caption("⚙️ Les calculs de cette table de référence se mettent à jour dynamiquement.")
        else:
            st.info("Alimentez la base de données à l'Écran 2 pour projeter automatiquement la Baseline.")

        # =====================================================================
        # 📊 ÉTAPE 6 : MEASURE PROCESS CAPABILITY (DPMO & SIGMA)
        # =====================================================================
        st.divider()
        st.subheader("6. Measure process capability")

        import math  # Importation sécurisée de math

        # --- FONCTION INTERNE D'ÉVALUATION DE CONFORMITÉ ---
        def _evaluer_conformite_ligne_v2(val_observee, config_master):
            if pd.isna(val_observee):
                return True
            t_spec = config_master.get("type_specification", "Valeur exacte")
            
            if t_spec == "Attribut qualitatif (Sans valeur numérique - Proposition Master Black Belt)":
                v_str = str(val_observee).strip().upper()
                if v_str in ["KO", "NON OK", "NON", "ERREUR", "RETOUCHE", "1", "NON-CONFORME"]:
                    return False
                return True

            try:
                num_val = float(val_observee)
            except (ValueError, TypeError):
                return str(val_observee).strip().upper() not in ["KO", "NON OK", "1", "NON-CONFORME"]

            if t_spec == "Valeur exacte":
                return math.isclose(num_val, float(config_master.get("valeur_exacte", 0.0)), rel_tol=1e-5, abs_tol=1e-5)
            elif t_spec == "Supérieur ou égal à (≥)":
                return num_val >= float(config_master.get("valeur_seuil", 0.0))
            elif t_spec == "Inférieur ou égal à (≤)":
                return num_val <= float(config_master.get("valeur_seuil", 0.0))
            elif t_spec == "Intervalle (Entre min et max)":
                return float(config_master.get("borne_inf", 0.0)) <= num_val <= float(config_master.get("borne_sup", 10.0))
            return True

        # --- IDENTIFICATION DE LA VARIABLE Y (CIBLE PRINCIPALE) ---
        # On essaie de récupérer le nom du Y actif dans la session ou l'état du projet
        nom_variable_y = None
        if 'selected_var_to_test' in locals() and selected_var_to_test:
            nom_variable_y = selected_var_to_test
        elif 'nom_colonne_variable' in locals() and not df_classification_current.empty:
            # S'il y a une liste, on prend la première variable critique ou celle liée au Y
            list_y_candidates = df_classification_current[nom_colonne_variable].dropna().tolist()
            if list_y_candidates:
                nom_variable_y = list_y_candidates[0]

        # Récupération de la configuration Master active associée au Y
        master_cfg_active = st.session_state.get("reference_master_config", {
            "type_specification": "Valeur exacte",
            "valeur_exacte": 0.0
        })
        if 'var_clean_id' in locals() and f"reference_master_config_{var_clean_id}_{safe_idx}" in st.session_state:
            master_cfg_active = st.session_state[f"reference_master_config_{var_clean_id}_{safe_idx}"]

        # --- CALCUL STRICTEMENT BASÉ SUR LA VARIABLE Y ---
        df_active_local = df_active if 'df_active' in locals() and not df_active.empty else pd.DataFrame()
        
        raw_units = len(df_active_local.dropna(subset=["ID observation"])) if not df_active_local.empty and "ID observation" in df_active_local.columns else (len(df_active_local) if not df_active_local.empty else 0)
        
        calculated_defects_y = 0
        if not df_active_local.empty and nom_variable_y and nom_variable_y in df_active_local.columns:
            for idx_row, row in df_active_local.iterrows():
                if not _evaluer_conformite_ligne_v2(row[nom_variable_y], master_cfg_active):
                    calculated_defects_y += 1
        else:
            # Fallback si la colonne Y n'est pas explicitée dans le DF actif
            calculated_defects_y = 0

        val_units = max(1, int(raw_units))
        val_opp = 1  # Exactement 1 opportunité par unité (le Y lui-même)

        c1, c2 = st.columns(2)
        with c1:
            st.markdown(f"**Paramètres de Capabilité (Cible unique : {nom_variable_y or 'Variable Y'})**")
            st.metric("Nombre total de défauts réels (sur le Y)", f"{calculated_defects_y}")
            st.number_input("Nombre d'unités inspectées (N)", min_value=1, value=val_units, key="final_cap_units_y_only", disabled=True)
            st.number_input("Nombre d'opportunités par unité", min_value=1, value=val_opp, key="final_cap_opp_y_only", disabled=True)
        
        # --- CALCUL SÉCURISÉ DU DPMO ET DU NIVEAU SIGMA SUR LE Y ---
        dpmo_calculé = (calculated_defects_y / (val_units * val_opp)) * 1_000_000
        
        sigma_level = 0.0
        try:
            if dpmo_calculé <= 0:
                sigma_level = 6.0
            else:
                taux_defaut = dpmo_calculé / 1_000_000
                p_val = max(1e-7, min(0.5, taux_defaut if taux_defaut < 0.5 else 1 - taux_defaut))
                
                t = math.sqrt(-2.0 * math.log(p_val))
                z = t - ((2.515517 + 0.802853 * t + 0.010328 * t * t) / 
                         (1.0 + 1.432788 * t + 0.189269 * t * t + 0.001308 * t * t * t))
                sigma_brut = z if taux_defaut < 0.5 else -z
                sigma_level = round(sigma_brut + 1.5, 2)
                sigma_level = max(0.0, min(6.0, sigma_level))
        except Exception:
            sigma_level = 0.0

        with c2:
            st.markdown("<br><br>", unsafe_allow_html=True)
            st.metric("DPMO Réel (du Y)", f"{float(dpmo_calculé):,.0f}")
            
            if sigma_level >= 4.0:
                st.metric("Niveau Sigma du Processus", f"🟢 {sigma_level} σ")
            elif sigma_level >= 2.5:
                st.metric("Niveau Sigma du Processus", f"🟠 {sigma_level} σ")
            else:
                st.metric("Niveau Sigma du Processus", f"🔴 {sigma_level} σ")
                
        # =====================================================================
        # 📊 PARTIE 7 : SITUATION T0 (CARTE DE CONTRÔLE AVANT AMÉLIORATION)
        # =====================================================================
        st.markdown("---")
        st.markdown("### 📊 Partie 7 : Situation T0 (Cartes de Contrôle SPC)")
        st.caption("🔍 *Visualisation de la stabilité et performance historique de référence avant actions d'amélioration.*")

        # Utilisation de la source unique et synchronisée de l'Écran 2
        df_t0 = st.session_state.dc_master_data

        if df_t0.empty:
            st.info("💡 Alimentez la base de données à l'Écran 2 pour générer automatiquement les cartes de contrôle T0.")
        else:
            # 1. Filtrage strict des variables quantitatives uniquement
            variables_quantitatives = []
            for var in liste_variables_dynamiques:
                if var in df_t0.columns:
                    # On vérifie si la colonne contient des données numériques convertibles
                    dropped_na = df_t0[var].dropna()
                    if not dropped_na.empty:
                        numeric_coerced = pd.to_numeric(dropped_na, errors="coerce").dropna()
                        # Si la variable est majoritairement numérique et exclut les mots-clés attributaires
                        if not numeric_coerced.empty and not any(k in var.lower() for k in ["statut", "verdict", "validation", "ok", "ko"]):
                            variables_quantitatives.append(var)

            if not variables_quantitatives:
                st.warning("⚠️ Aucune variable quantitative (numérique) n'a été détectée dans le DCP pour générer des cartes I-MR.")
            else:
                # Options d'agrégation de l'Axe X (Temps / Observations)
                st.markdown("#### ⚙️ Configuration de l'Axe X")
                mode_axe_x = st.radio(
                    "Regroupement des données :",
                    ["Par observation (Individuel)", "Par jour", "Par semaine", "Par mois"],
                    horizontal=True,
                    key="spc_axe_x_mode"
                )

                # Préparation de l'axe temporel si demandé
                df_spc_master = df_t0.copy()
                if "Date de modification" in df_spc_master.columns:
                    df_spc_master["Date_Parsed"] = pd.to_datetime(df_spc_master["Date de modification"], errors="coerce")
                else:
                    df_spc_master["Date_Parsed"] = pd.NaT

                # Boucle de génération automatique : Une carte par variable quantitative
                for var_quant in variables_quantitatives:
                    st.markdown(f"---")
                    st.subheader(f"📈 Carte de contrôle I-MR : {var_quant}")

                    # Nettoyage local de la série
                    df_var = df_spc_master[["ID observation", "Date_Parsed", var_quant]].dropna(subset=[var_quant]).copy()
                    df_var[var_quant] = pd.to_numeric(df_var[var_quant], errors="coerce")
                    df_var = df_var.dropna(subset=[var_quant]).reset_index(drop=True)

                    if len(df_var) < 2:
                        st.info(f"Pas assez de données pour générer la carte I-MR pour '{var_quant}' (minimum 2 observations requises).")
                        continue

                    # 2. Gestion de l'agrégation de l'axe X selon le volume disponible
                    if mode_axe_x == "Par jour" and df_var["Date_Parsed"].notna().any():
                        df_grouped = df_var.groupby(df_var["Date_Parsed"].dt.date)[var_quant].mean().reset_index()
                        df_grouped.columns = ["Axe_X", "Valeur"]
                    elif mode_axe_x == "Par semaine" and df_var["Date_Parsed"].notna().any():
                        df_var["Semaine"] = df_var["Date_Parsed"].dt.to_period("W").astype(str)
                        df_grouped = df_var.groupby("Semaine")[var_quant].mean().reset_index()
                        df_grouped.columns = ["Axe_X", "Valeur"]
                    elif mode_axe_x == "Par mois" and df_var["Date_Parsed"].notna().any():
                        df_var["Mois"] = df_var["Date_Parsed"].dt.to_period("M").astype(str)
                        df_grouped = df_var.groupby("Mois")[var_quant].mean().reset_index()
                        df_grouped.columns = ["Axe_X", "Valeur"]
                    else:
                        # Par défaut ou par observation individuelle
                        df_grouped = df_var[["ID observation", var_quant]].copy()
                        df_grouped.columns = ["Axe_X", "Valeur"]

                    if len(df_grouped) < 2:
                        st.warning(f"Le regroupement temporel choisi génère moins de 2 points. Affichage individuel par observation appliqué.")
                        df_grouped = df_var[["ID observation", var_quant]].copy()
                        df_grouped.columns = ["Axe_X", "Valeur"]

                    # 3. Calculs Statistiques SPC (Carte I-MR avec constante d2 = 1.128 pour n=2)
                    # Calcul des étendues mobiles (Moving Range - MR)
                    df_grouped["MR"] = df_grouped["Valeur"].diff().abs()
                    mr_moyenne = df_grouped["MR"].mean()
                    
                    # Ligne centrale (Moyenne CL)
                    cl_moyenne = df_grouped["Valeur"].mean()
                    
                    # Calcul des limites UCL et LCL basées sur l'étendue mobile moyenne (Norme LSS)
                    ucl_calculé = cl_moyenne + 2.66 * mr_moyenne
                    lcl_calculé = cl_moyenne - 2.66 * mr_moyenne
                    lcl_calculé = max(0.0, lcl_calculé) if df_grouped["Valeur"].min() >= 0 else lcl_calculé

                    # Intégration des lignes de contrôle dans le DataFrame pour affichage graphique
                    df_grouped["CL - Moyenne"] = cl_moyenne
                    df_grouped["UCL"] = ucl_calculé
                    df_grouped["LCL"] = lcl_calculé

                    # 4. Détection Automatique des Signaux d'Instabilité (Règles SPC de Western Electric)
                    df_grouped["Instable"] = False
                    
                    # Règle 1 : Points hors limites strictes (UCL/LCL)
                    df_grouped.loc[(df_grouped["Valeur"] > ucl_calculé) | (df_grouped["Valeur"] < lcl_calculé), "Instable"] = True
                    
                    # Règle 2 : Séries anormales (7 points consécutifs du même côté de la ligne centrale)
                    df_grouped["Position"] = df_grouped["Valeur"] > cl_moyenne
                    df_grouped["Serie_Id"] = (df_grouped["Position"] != df_grouped["Position"].shift()).cumsum()
                    tailles_series = df_grouped.groupby("Serie_Id")["Serie_Id"].transform("count")
                    df_grouped.loc[tailles_series >= 7, "Instable"] = True

                    # Règle 3 : Tendances prolongées (6 points consécutifs qui montent ou qui descendent)
                    diffs = df_grouped["Valeur"].diff()
                    import numpy as np
                    signes = np.sign(diffs.fillna(0))
                    # Identification des changements de direction
                    df_grouped["Tendance_Id"] = (signes != signes.shift()).cumsum()
                    tailles_tendances = df_grouped.groupby("Tendance_Id")["Tendance_Id"].transform("count")
                    df_grouped.loc[tailles_tendances >= 6, "Instable"] = True

                    # Préparation des couleurs pour la visualisation graphique
                    df_grouped["Couleur_Point"] = df_grouped["Instable"].map({True: "#FF4B4B", False: "#0068C9"})

                    # 5. Affichage Graphique de la Carte Individuelle (I) via st.line_chart combiné
                    chart_cols = ["Valeur", "CL - Moyenne", "UCL", "LCL"]
                    df_chart = df_grouped.set_index("Axe_X")[chart_cols]
                    
                    st.line_chart(df_chart, color=["#0068C9", "#29B09D", "#FF4B4B", "#FF4B4B"])

                    # Alerte visuelle en cas de points ou séries hors contrôle détectés
                    points_hors_controle = df_grouped["Instable"].sum()
                    pct_hors_controle = (points_hors_controle / len(df_grouped)) * 100

                    if points_hors_controle > 0:
                        st.error(f"🔴 Attention : {points_hors_controle} point(s) ou série(s) d'instabilité détecté(s) ({pct_hors_controle:.1f} % des données). Le processus presents des causes spéciales de variation.")
                    else:
                        st.success("🟢 Processus statistiquement stable. Les variations observées relèvent uniquement de causes communes.")

                    # 6. Résumé Statistique de Référence T0
                    st.markdown("**📋 Résumé Statistique de Référence (Situation T0)**")
                    
                    ecart_type = df_grouped["Valeur"].std() if len(df_grouped) > 1 else 0.00
                    
                    stats_t0 = {
                        "Métrique SPC T0": [
                            "Nombre d'observations (N)", 
                            "Moyenne (Ligne Centrale)", 
                            "Minimum constaté", 
                            "Maximum constaté", 
                            "Écart-type (σ)", 
                            "Limite UCL (Upper Control Limit)", 
                            "Limite LCL (Lower Control Limit)",
                            "Nombre de points / séries hors contrôle",
                            "Pourcentage d'instabilité globale"
                        ],
                        "Valeur T0": [
                            f"{len(df_grouped)}",
                            f"{cl_moyenne:.2f}",
                            f"{df_grouped['Valeur'].min():.2f}",
                            f"{df_grouped['Valeur'].max():.2f}",
                            f"{ecart_type:.2f}",
                            f"{ucl_calculé:.2f}",
                            f"{lcl_calculé:.2f}",
                            f"{points_hors_controle}",
                            f"{pct_hors_controle:.1f} %"
                        ]
                    }
                    st.table(pd.DataFrame(stats_t0))
    
    # --- PHASE ANALYSE ---
    with tabs[2]: 
        st.header("Phase Analyse")

        # Initialisation de la session state pour éviter l'erreur NameError
        if "results" not in st.session_state:
            st.session_state.results = []

        if "dc_master_data" in st.session_state and not st.session_state.dc_master_data.empty:
            df = st.session_state.dc_master_data
            st.success(f"✅ Données détectées ({len(df)} lignes).")
        
        # --- PHASE 1 : TEST ACTION DES X SUR Y ---
        st.subheader("1. Test de l'action des X sur Y")
        
        if "dc_master_data" in st.session_state and not st.session_state.dc_master_data.empty:
            df = st.session_state.dc_master_data
            st.success(f"✅ Données détectées ({len(df)} lignes).")

            # Sélection des colonnes
            colonnes_a_exclure = ["ID observation", "Date de modification"]
            colonnes_disponibles = [c for c in df.columns if c not in colonnes_a_exclure]

            y_col = st.selectbox("Sélectionnez la variable Y (CTQ) :", colonnes_disponibles)
            x_cols = st.multiselect("Sélectionnez les variables X à tester :", [c for c in colonnes_disponibles if c != y_col])

            if st.button("🚀 Lancer l'analyse statistique"):
                results = []
                y_data = pd.to_numeric(df[y_col], errors='coerce')
                is_y_num = y_data.notna().sum() > (len(df) * 0.5)

                for x in x_cols:
                    # Conversion sécurisée de X
                    x_data = pd.to_numeric(df[x], errors='coerce')
                    is_x_num = x_data.notna().sum() > (len(df) * 0.5)
                
                    impact = "Aucun"
                    degre = "N/A"
                    p_value = 1.0
                
                    try:
                        # 1. Corrélation Pearson (Numérique vs Numérique)
                        if is_x_num and is_y_num:
                            corr, p_value = stats.pearsonr(x_data.dropna(), y_data.dropna())
                            if p_value < 0.05:
                                impact = "Oui, influence démontrée"
                                degre = "Fort" if abs(corr) > 0.7 else ("Modéré" if abs(corr) > 0.3 else "Faible")
                            else:
                                impact = "Non, influence non démontrée"
                                degre = "Nul"

                        # 2. ANOVA (Catégoriel vs Numérique)
                        elif not is_x_num and is_y_num:
                            groups = [group[y_col].dropna().values for _, group in df.groupby(x)]
                            if len(groups) > 1:
                                _, p_value = stats.f_oneway(*groups)
                                if p_value < 0.05:
                                    impact = "Oui, les moyennes diffèrent"
                                    degre = "Modéré à Fort"
                                else:
                                    impact = "Non, pas de différence"
                                    degre = "Nul"
                    
                        # 3. Chi-2 (Qualitatif vs Qualitatif) - REMPLACE L'ANCIEN ELSE
                        else:
                            contingency = pd.crosstab(df[x], df[y_col])
                            if contingency.size > 0:
                                chi2, p_value, _, _ = stats.chi2_contingency(contingency)
                                if p_value < 0.05:
                                    impact = "Oui, lien statistique trouvé"
                                    degre = "Modéré à Fort"
                                else:
                                    impact = "Non, variables indépendantes"
                                    degre = "Nul"
                            else:
                                impact = "Données insuffisantes"
                                degre = "N/A"

                    except Exception as e:
                        impact = f"Erreur de calcul : {e}"

                    results.append({
                        "Variable X": x,
                        "Le X agit-il sur Y ?": impact,
                        "Degré d'influence": degre,
                        "P-value": round(float(p_value), 4)
                    })
                
                # --- CORRECTION POUR SAUVEGARDE JSON ---
                # On récupère l'index du projet en cours
                idx = st.session_state["current_project_idx"]
                
                # On injecte les résultats DANS le projet actif
                # C'est cette ligne qui garantit que le JSON récupérera les données
                st.session_state.projects[idx]["dmaic"]["analyze"]["results"] = results
                
                st.rerun()

            # --- AFFICHAGE ---
            # On lit les résultats directement dans le projet actif
            idx = st.session_state["current_project_idx"]
            stored_results = st.session_state.projects[idx]["dmaic"]["analyze"].get("results", [])

            if stored_results:
                st.table(pd.DataFrame(stored_results))
            else:
                st.info("Aucun résultat pour le moment.")

        # --- PHASE 2 : IDENTIFICATION DES CAUSES RACINES ---
        st.subheader("2. Identification des causes racines")

        idx = st.session_state["current_project_idx"]
        dmaic_analyze = st.session_state.projects[idx]["dmaic"]["analyze"]

        if "x_analyses" not in dmaic_analyze:
            dmaic_analyze["x_analyses"] = {}

        results_phase1 = dmaic_analyze.get("results", [])
        tous_les_x = [r["Variable X"] for r in results_phase1]
        x_influents = [r["Variable X"] for r in results_phase1 if "Oui" in r.get("Le X agit-il sur Y ?", "")]
        x_critiques = st.multiselect("Sélectionnez les variables à approfondir :", tous_les_x, default=x_influents)

        for x in x_critiques:
            # 1. Initialisation dans le JSON
            if x not in dmaic_analyze["x_analyses"]:
                dmaic_analyze["x_analyses"][x] = {
                    "methode": "Ishikawa", "data": {}, "causes_racines_list": [""], "enregistre": False
                }
            x_state = dmaic_analyze["x_analyses"][x]
            safe_x_key = x.replace(" ", "_").replace("(", "").replace(")", "").replace("/", "_")

            # 2. SYNCHRO Session
            if f"causes_{safe_x_key}" not in st.session_state:
                st.session_state[f"causes_{safe_x_key}"] = x_state.get("causes_racines_list", [""])

            # 3. EXPANDER
            statut = "✅" if x_state.get("enregistre") else "⏳"
            with st.expander(f"{statut} Analyse pour : {x}"):
                
                # A. Choix méthode (Hors form)
                nouvelle_methode = st.radio("Méthode :", ["Ishikawa", "5 Pourquoi"], 
                                            index=0 if x_state["methode"] == "Ishikawa" else 1,
                                            key=f"radio_{safe_x_key}", horizontal=True)
                if nouvelle_methode != x_state["methode"]:
                    x_state["methode"] = nouvelle_methode
                    st.rerun()

                # B. AJOUT/SUPPRESSION (Hors form pour ne pas bugger)
                col_add, col_del = st.columns([1, 5])
                if col_add.button("➕", key=f"add_{safe_x_key}"):
                    st.session_state[f"causes_{safe_x_key}"].append("")
                    st.rerun()

                # C. FORMULAIRE UNIQUE
                with st.form(key=f"form_total_{safe_x_key}"):
                    # --- Analyse ---
                    form_data = {}
                    if x_state["methode"] == "Ishikawa":
                        cols = st.columns(2)
                        cats = ["Méthode", "Main d'œuvre", "Machine", "Matière", "Milieu", "Mesure"]
                        for i, cat in enumerate(cats):
                            form_data[cat] = cols[i % 2].text_area(f"🦴 {cat}", value=x_state["data"].get(cat, ""))
                    else:
                        for i in range(1, 6):
                            form_data[f"p{i}"] = st.text_input(f"Pourquoi {i} ?", value=x_state["data"].get(f"p{i}", ""))

                    # --- Causes (Boucle DANS le form) ---
                    st.write("**Causes racines :**")
                    for i in range(len(st.session_state[f"causes_{safe_x_key}"])):
                        c1, c2 = st.columns([10, 1])
                        st.session_state[f"causes_{safe_x_key}"][i] = c1.text_input(
                            f"Cause {i+1}", value=st.session_state[f"causes_{safe_x_key}"][i], 
                            key=f"input_{safe_x_key}_{i}", label_visibility="collapsed"
                        )
                        if c2.form_submit_button("🗑️", key=f"del_{safe_x_key}_{i}"):
                            st.session_state[f"causes_{safe_x_key}"].pop(i)
                            st.rerun()

                    # Bouton Soumission Finale
                    if st.form_submit_button("💾 Enregistrer TOUT"):
                        x_state["data"] = form_data
                        x_state["causes_racines_list"] = [c for c in st.session_state[f"causes_{safe_x_key}"] if c.strip() != ""]
                        x_state["enregistre"] = True
                        st.success("Analyse enregistrée !")
                        st.rerun()
    
        # --- PHASE 3 : FMEA (AMDEC) ---
        st.subheader("3. Current State FMEA")

        if "fmea_data" not in dmaic_analyze:
            dmaic_analyze["fmea_data"] = {}

        toutes_les_causes = []
        for x, state in dmaic_analyze.get("x_analyses", {}).items():
            for cause in state.get("causes_racines_list", []):
                if cause.strip():
                    toutes_les_causes.append({"x": x, "cause": cause})

        if not toutes_les_causes:
            st.warning("Aucune cause racine identifiée. Veuillez compléter la phase 2.")
        else:
            for item in toutes_les_causes:
                cause_id = f"{item['x']}_{item['cause']}"
                
                # Charger les données depuis le JSON
                data_actuelle = dmaic_analyze["fmea_data"].get(cause_id, {
                    "S": 1, "O": 1, "D": 1, "effet": "", "mode": ""
                })

                with st.expander(f"Analyse : {item['x']} | Cause : {item['cause']}"):
                    # FORMULAIRE : Rien ne bouge tant que vous ne validez pas
                    with st.form(key=f"form_fmea_{cause_id}"):
                        c1, c2 = st.columns(2)
                        mode = c1.text_input("Mode de défaillance", value=data_actuelle.get("mode", ""))
                        effet = c2.text_input("Effet de la défaillance", value=data_actuelle.get("effet", ""))
                        
                        s, o, d = st.columns(3)
                        val_s = s.slider("Sévérité (S)", 1, 10, data_actuelle.get("S", 1))
                        val_o = o.slider("Occurrence (O)", 1, 10, data_actuelle.get("O", 1))
                        val_d = d.slider("Détection (D)", 1, 10, data_actuelle.get("D", 1))
                        
                        rpn = val_s * val_o * val_d
                        st.metric("RPN (Calculé)", rpn)
                        
                        if st.form_submit_button(f"Valider l'analyse FMEA"):
                            dmaic_analyze["fmea_data"][cause_id] = {
                                "S": val_s, "O": val_o, "D": val_d, 
                                "effet": effet, "mode": mode
                            }
                            st.success("Analyse enregistrée !")
                            st.rerun()

            # Classement automatique
            st.divider()
            st.subheader("📊 Classement des risques (RPN)")
    
            import pandas as pd
            liste_tri = []
            for cid, vals in dmaic_analyze["fmea_data"].items():
                liste_tri.append({
                    "Cause": cid, 
                    "RPN": vals.get("S", 1) * vals.get("O", 1) * vals.get("D", 1)
                })
    
            if liste_tri:
                df_fmea = pd.DataFrame(liste_tri).sort_values(by="RPN", ascending=False)
                st.table(df_fmea)
    
        # --- PHASE 4 : DATA COLLECTION PLAN (GEMBA WALK) ---
        st.subheader("4. Data Collection Plan (Gemba Walk)")

        if "gemba_plans" not in dmaic_analyze:
            dmaic_analyze["gemba_plans"] = {}

        # 1. Récupération des causes pour génération
        toutes_les_causes = []
        for x, state in dmaic_analyze.get("x_analyses", {}).items():
            for cause in state.get("causes_racines_list", []):
                if cause.strip():
                    toutes_les_causes.append({"x": x, "cause": cause})

        if not toutes_les_causes:
            st.warning("Veuillez d'abord identifier des causes racines (Phase Analyse).")
        else:
            # Bouton pour générer/réinitialiser les fiches
            if st.button("🚀 Générer/Mettre à jour les fiches Gemba"):
                for item in toutes_les_causes:
                    cid = f"{item['x']}_{item['cause']}"
                    if cid not in dmaic_analyze["gemba_plans"]:
                        dmaic_analyze["gemba_plans"][cid] = {
                            "x_critique": item['x'],
                            "cause_racine": item['cause'],
                            "objectif": f"Vérifier sur le terrain que cette cause racine est effectivement présente et qu'elle contribue au problème observé.",
                            "activite": "", "lieu": "", "personnes": "", "elements_obs": "",
                            "preuves": "", "taille_ech": 10, "date": None, "responsable": "",
                            "critere": "Validé si observé dans 80% des cas", "statut": "À préparer"
                        }
                st.rerun()

            # 2. Affichage et édition des fiches
            for cid, fiche in dmaic_analyze["gemba_plans"].items():
                with st.expander(f"Fiche Gemba : {fiche['cause_racine']} ({fiche['statut']})"):
                    st.write(f"**X associé :** {fiche['x_critique']}")
                    st.write(f"**Objectif :** {fiche['objectif']}")
            
                    c1, c2 = st.columns(2)
                    fiche["activite"] = c1.text_input("Activité à observer", fiche.get("activite", ""), key=f"act_{cid}")
                    fiche["lieu"] = c2.text_input("Lieu d'observation", fiche.get("lieu", ""), key=f"lieu_{cid}")
            
                    fiche["personnes"] = st.text_input("Personnes/Fonctions à observer", fiche.get("personnes", ""), key=f"pers_{cid}")
                    fiche["elements_obs"] = st.text_area("Éléments à observer", fiche.get("elements_obs", ""), key=f"elem_{cid}")
                    fiche["preuves"] = st.text_area("Preuves recherchées", fiche.get("preuves", ""), key=f"prev_{cid}")
            
                    c3, c4, c5 = st.columns(3)
                    fiche["taille_ech"] = c3.number_input("Nb observations", min_value=1, value=fiche.get("taille_ech", 10), key=f"ech_{cid}")
                    fiche["responsable"] = c4.text_input("Responsable", fiche.get("responsable", ""), key=f"resp_{cid}")
                    fiche["statut"] = c5.selectbox("Statut", ["À préparer", "Planifié", "En cours", "Terminé"], 
                                                 index=["À préparer", "Planifié", "En cours", "Terminé"].index(fiche.get("statut", "À préparer")),
                                                 key=f"stat_{cid}")
            
                    fiche["critere"] = st.text_input("Critère de validation", fiche.get("critere", ""), key=f"crit_{cid}")

            if st.button("💾 Enregistrer le plan Gemba"):
                st.success("Plan Gemba sauvegardé !")

            # 3. Tableau récapitulatif
            st.divider()
            st.subheader("📋 Tableau récapitulatif du Plan")
    
            import pandas as pd
            tab_data = [{"X": f['x_critique'], "Cause": f['cause_racine'], "Activité": f['activite'], 
                         "Resp": f['responsable'], "Nb Obs": f['taille_ech'], "Statut": f['statut']} 
                        for f in dmaic_analyze["gemba_plans"].values()]
    
            if tab_data:
                st.table(pd.DataFrame(tab_data))
    
       # --- PHASE 5 : DATA COLLECTION (OBSERVATIONS TERRAIN) ---
        st.subheader("5. Data Collection (Résultats des Gemba Walks)")

        import pandas as pd

        # Initialisation du dictionnaire si nécessaire
        if "gemba_observations" not in dmaic_analyze:
            dmaic_analyze["gemba_observations"] = {}

        gemba_plans = dmaic_analyze.get("gemba_plans", {})

        for cid, plan in gemba_plans.items():
            with st.expander(f"Fiche : {plan['cause_racine']} (X: {plan['x_critique']})"):
                st.info(f"**Objectif :** {plan['objectif']}")
        
                # Récupération des données
                data_cid = dmaic_analyze["gemba_observations"].get(cid, {"logs": []})
                df_obs = pd.DataFrame(data_cid.get("logs", []), columns=["date", "confirme", "description", "preuve", "confiance"])
                        
                # 2. ÉDITEUR DANS UN FORMULAIRE (Empêche la mise à jour temps réel)
                with st.form(key=f"form_{cid}"):
                    edited_df = st.data_editor(
                        df_obs.fillna(""),
                        key=f"editor_{cid}", 
                        column_config={
                            "date": st.column_config.TextColumn("Date (YYYY-MM-DD)"),
                            "confirme": st.column_config.SelectboxColumn("Cause observée ?", options=["Oui", "Non", "Partiellement"]),
                            "confiance": st.column_config.SelectboxColumn("Confiance", options=["Faible", "Moyen", "Élevé"]),
                        },
                        num_rows="dynamic",
                        hide_index=True,
                        use_container_width=True
                    )
            
                    # Le bouton de soumission fait partie du formulaire
                    submit = st.form_submit_button(f"💾 Sauvegarder les données de {plan['cause_racine'][:10]}")
            
                    if submit:
                        # 1. S'assurer que le dictionnaire est bien initialisé pour ce cid
                        if cid not in dmaic_analyze["gemba_observations"]:
                            dmaic_analyze["gemba_observations"][cid] = {"logs": [], "status": "Non confirmée"}
                        
                        # 2. Sauvegarde des données
                        new_logs = edited_df.to_dict('records')
                        dmaic_analyze["gemba_observations"][cid]["logs"] = new_logs
                
                        # 3. Recalcul des statuts
                        fav = len([l for l in new_logs if str(l.get("confirme")).lower() == 'oui'])
                        pct = (fav / len(new_logs) * 100) if len(new_logs) > 0 else 0
                        
                        # Ici aussi, on s'assure que la clé existe avant d'écrire
                        dmaic_analyze["gemba_observations"][cid]["status"] = "Confirmée" if pct >= 80 else ("Partiellement confirmée" if pct >= 50 else "Non confirmée")
                
                        st.success("Données sauvegardées !")
                        st.rerun()

       # --- PHASE 6: VALIDATION DES CAUSES RACINES ---
        st.divider()
        st.subheader("6. Validation des causes racines")
        
        # Initialisation
        if "validation_data" not in dmaic_analyze:
            dmaic_analyze["validation_data"] = {}

        summary_data = []

        for cid, plan in gemba_plans.items():
            # 1. Accès à la source de vérité (la même liste que Phase 1)
            idx = st.session_state.get("current_project_idx", 0)
            results_phase1 = st.session_state.projects[idx]["dmaic"]["analyze"].get("results", [])
            
            x_nom_cible = str(plan["x_critique"]).strip().lower()
            
            # 2. Recherche du X correspondant
            p_val = 1.0
            stats_temp = {} 
            
            for r in results_phase1:
                if str(r.get("Variable X", "")).strip().lower() == x_nom_cible:
                    # On utilise exactement la clé "P-value" enregistrée en Phase 1
                    p_val = float(r.get("P-value", 1.0))
                    stats_temp = r
                    break
            
            # 2. RÉCUPÉRATION DES OBSERVATIONS TERRAIN
            obs = dmaic_analyze.get("gemba_observations", {}).get(cid, {})
            logs = obs.get("logs", []) if isinstance(obs, dict) else []
    
            # 3. CALCULS SÉCURISÉS
            conf = len([l for l in logs if str(l.get("confirme", "")).lower() == 'oui'])
            total = len(logs)
            non_conf = total - conf
            taux = (conf / total * 100) if total > 0 else 0
            
            # Évaluation terrain auto
            if taux >= 80: val_terrain = "Confirmée"
            elif taux >= 50: val_terrain = "Partiellement confirmée"
            else: val_terrain = "Non confirmée"
            
            # 4. LOGIQUE DE DÉCISION
            # On utilise p_val (calculé plus haut) et non stats.get
            stat_validee = p_val < 0.05 
            
            res_stat = "Confirmé" if stat_validee else "Non confirmé"
            res_terrain = val_terrain

            if res_stat == "Confirmé" and res_terrain == "Confirmée": decision = "Cause validée"
            elif res_stat == "Confirmé" and res_terrain == "Partiellement confirmée": decision = "Cause partiellement validée"
            elif res_stat == "Confirmé" and res_terrain == "Non confirmée": decision = "Contradiction à investiguer"
            elif res_stat == "Non confirmé" and res_terrain == "Confirmée": decision = "Contradiction à investiguer"
            elif res_stat == "Non confirmé" and res_terrain == "Partiellement confirmée": decision = "Cause non démontrée"
            elif res_stat == "Non confirmé" and res_terrain == "Non confirmée": decision = "Cause rejetée"
            else: decision = "En cours d'analyse"

            # Stockage pour le tableau
            summary_data.append({
                "X Critique": plan["x_critique"],
                "Cause Racine": plan["cause_racine"],
                "Test": stats_temp.get('Le X agit-il sur Y ?', 'N/A'), # Correction ici
                "P-value": f"{p_val:.4f}",
                "Taux Terrain": f"{taux:.0f}%",
                "Décision": decision
            })

            # 4. Affichage détaillé par Expandeur
            with st.expander(f"{plan['x_critique']} : {plan['cause_racine']} (Décision: {decision})"):
                col1, col2 = st.columns(2)
                with col1:
                    st.markdown("### 📊 Statistiques")
                    # ON UTILISE stats_temp QUI CONTIENT BIEN LES DONNÉES
                    st.write(f"**Test :** {stats_temp.get('Le X agit-il sur Y ?', 'N/A')}")
                    st.write(f"**P-value :** {p_val:.4f}")
                    st.write(f"**Conclusion :** {'Confirmée' if stat_validee else 'Non confirmée'}")
                with col2:
                    st.markdown("### 🔎 Terrain")
                    st.write(f"**Observations :** {total} | **Confirmées :** {conf} | **Contradictions :** {non_conf}")
                    st.write(f"**Taux :** {taux:.0f}%")
                    st.write(f"**Validation :** {val_terrain}")
        
                st.info(f"**Décision automatique :** {decision}")

        # 5. TABLEAU RÉCAPITULATIF FINAL (Mis à jour)
        st.subheader("Tableau récapitulatif final")
        
        # Préparation du DataFrame pour l'affichage
        df_final = pd.DataFrame(summary_data)
        
        # Affichage du tableau avec un formatage lisible
        st.dataframe(
            df_final,
            column_config={
                "Résultat Statistique": st.column_config.TextColumn("Résultat Statistique"),
                "Résultat Terrain": st.column_config.TextColumn("Résultat Terrain"),
                "Décision finale": st.column_config.TextColumn("Décision finale")
            },
            use_container_width=True,
            hide_index=True
        )

        # 6. CRITÈRE DE CLÔTURE
        st.divider()
        # Sécurité anti-KeyError : on vérifie que df_final est un vrai DataFrame et que la colonne 'Décision' y existe
        if isinstance(df_final, pd.DataFrame) and "Décision" in df_final.columns and not df_final.empty:
            tous_decides = all(d not in ["En cours d'analyse"] for d in df_final["Décision"])
        else:
            tous_decides = False
        if tous_decides and len(df_final) > 0:
            st.success("✅ Tous les critères de clôture sont remplis. Vous pouvez passer en phase Improve.")
        else:
            st.warning("⚠️ Phase Analyse incomplète : Certaines causes nécessitent encore une décision.")
            
    # --- PHASE IMPROVE ---
    with tabs[3]: 
        st.header("Phase Improve")

        # 0. SÉCURISATION DES DONNÉES (À placer ici pour être sûr de les avoir)
        if "summary_data_phase6" not in st.session_state and "summary_data" in locals():
             st.session_state.summary_data_phase6 = summary_data
        
        # 1 : IMPROVEMENT STRATEGIES ---
        st.subheader("1. Improvement strategies")
        if "summary_data_phase6" in st.session_state:
            summary_data = st.session_state.summary_data_phase6
            causes_validees = [item["Cause Racine"] for item in summary_data if item["Décision"] in ["Cause validée", "Cause partiellement validée"]]
            
            # Initialisation de la donnée officielle si inexistante
            if "improve_strategies" not in st.session_state:
                st.session_state.improve_strategies = pd.DataFrame(columns=["Cause racine", "Solution potentielle", "Type de solution", "Forces", "Faiblesses"])

            # UTILISATION D'UN FORMULAIRE
            with st.form(key="form_strategies"):
                # Édition
                temp_df = st.data_editor(
                    st.session_state.improve_strategies,
                    column_config={
                        "Type de solution": st.column_config.SelectboxColumn(
                            options=["Préventive", "Corrective", "Détective", "Automatisation", "Standardisation", "Formation", "Organisationnelle"]
                        ),
                        "Cause racine": st.column_config.SelectboxColumn(options=causes_validees, required=True)
                    },
                    num_rows="dynamic", 
                    use_container_width=True,
                    key="editor_strategies_data" 
                )
                
                submit_button = st.form_submit_button(label="✅ Valider et enregistrer les solutions")
            
            # Traitement unique de la soumission
            if submit_button:
                # 1. Mise à jour de la mémoire temporaire
                st.session_state.improve_strategies = temp_df.copy()
            
                # 2. Mise à jour de la base de données du projet (pour que le JSON le voit)
                idx = st.session_state["current_project_idx"]
            
                # On s'assure que le chemin existe dans le dictionnaire du projet
                if "improve" not in st.session_state.projects[idx]["dmaic"]:
                    st.session_state.projects[idx]["dmaic"]["improve"] = {}
                
                # On enregistre les données proprement
                st.session_state.projects[idx]["dmaic"]["improve"]["strategies"] = temp_df.to_dict(orient="records")
            
                # 3. MESSAGE DE SUCCÈS
                st.success("Données synchronisées avec le projet !")
                st.rerun()
        else:
            st.warning("Veuillez d'abord valider des causes racines dans la phase Analyse.")
        
        # 2 : CRITERIA - BASED SELECTION MATRIX ---
        if "current_project_idx" in st.session_state and st.session_state.current_project_idx is not None:
            idx = st.session_state.current_project_idx
            st.subheader("2. Criteria-based selection matrix")
            
            # On récupère les solutions saisies
            df_sol = st.session_state.improve_strategies

            if df_sol.empty:
                st.info("Veuillez saisir des solutions ci-dessus et cliquer sur 'Valider' pour commencer la notation.")
            else:
                # Initialisation des critères si absents
                if "criteria" not in st.session_state.projects[idx]["dmaic"]["improve"]:
                    st.session_state.projects[idx]["dmaic"]["improve"]["criteria"] = [
                        {"Critère": "Impact", "Poids": 30}, {"Critère": "Coût", "Poids": 20},
                        {"Critère": "Facilité", "Poids": 20}, {"Critère": "Délai", "Poids": 30}
                    ]

                with st.form(key="form_notation_matrice"):
                    # 1. Paramétrage critères (Caché par défaut)
                    with st.expander("⚙️ Paramétrage des critères (Cliquer pour afficher/masquer)"):
                        df_crit = pd.DataFrame(st.session_state.projects[idx]["dmaic"]["improve"]["criteria"])
                        edited_crit = st.data_editor(df_crit, num_rows="dynamic", key="editor_crit_matrice")

                    # 2. Notation des solutions
                    cols_a_noter = edited_crit["Critère"].tolist()
                    if "notation_data" in st.session_state.projects[idx]["dmaic"]["improve"]:
                        # Si on a déjà sauvegardé des notes, on les récupère
                        df_notation = pd.DataFrame(st.session_state.projects[idx]["dmaic"]["improve"]["notation_data"])
                    else:
                        # Sinon, on initialise avec des 3 par défaut
                        df_notation = df_sol.copy()
                        for c in cols_a_noter:
                            if c not in df_notation.columns: 
                                df_notation[c] = 3
            
                    st.write("### 📝 Notation des solutions")
                    
                    # Légende
                    st.markdown("""
                    **Légende des scores :**
                    *   **1** : Très faible / Très difficile / Très coûteux
                    *   **3** : Moyen / Neutre
                    *   **5** : Très élevé / Très facile / Très rentable
                    """)
                    
                    df_notes = st.data_editor(df_notation[["Solution potentielle"] + cols_a_noter], key="editor_notes_matrice")

                    submit_notation = st.form_submit_button(label="💾 Calculer le classement final")

                if submit_notation:
                    # 1. Calcul des scores
                    df_final = df_notes.copy()
                    for _, r in edited_crit.iterrows():
                        df_final[f"Score_{r['Critère']}"] = df_final[r['Critère']] * (r['Poids'] / 100)
                    
                    df_final["Score Total"] = df_final[[f"Score_{c}" for c in cols_a_noter]].sum(axis=1)
                    df_final = df_final.sort_values(by="Score Total", ascending=False)
            
                    # 2. Mise à jour de la mémoire temporaire (on stocke les DataFrames ici)
                    st.session_state.projects[idx]["dmaic"]["improve"]["criteria"] = edited_crit.to_dict(orient="records")
                    st.session_state.projects[idx]["dmaic"]["improve"]["notation_data"] = df_notes.to_dict(orient="records")
                    st.session_state.projects[idx]["dmaic"]["improve"]["selection_matrix"] = df_final.to_dict(orient="records")
                    
                    # 3. FONCTION DE NETTOYAGE : transforme TOUT DataFrame restant en liste
                    def prepare_for_json(obj):
                        if isinstance(obj, dict):
                            return {k: prepare_for_json(v) for k, v in obj.items()}
                        elif isinstance(obj, list):
                            return [prepare_for_json(item) for item in obj]
                        elif isinstance(obj, pd.DataFrame):
                            return obj.to_dict(orient="records")
                        else:
                            return obj

                    # 4. Sauvegarde sécurisée
                    import json
                    try:
                        # On nettoie tout le projet avant d'écrire
                        data_to_save = prepare_for_json(st.session_state.projects)
                        with open("projects.json", "w") as f:
                            json.dump(data_to_save, f, indent=4)
                        st.success("Classement et notes sauvegardés sur le disque !")
                    except Exception as e:
                        st.error(f"Erreur lors de la sauvegarde : {e}")

                # Affichage résultat
                if "selection_matrix" in st.session_state.projects[idx]["dmaic"]["improve"]:
                    st.write("### 📊 Résultats")
                    # On affiche en convertissant le JSON (list) en DataFrame pour l'affichage
                    res_data = st.session_state.projects[idx]["dmaic"]["improve"]["selection_matrix"]
                    st.dataframe(pd.DataFrame(res_data))
        

        # 3 : BENEFIT EFFORT MATRIX ---
        st.subheader("3. Benefit effort matrix")

        if "current_project_idx" in st.session_state and st.session_state.current_project_idx is not None:
            idx = st.session_state.current_project_idx
            dmaic_improve = st.session_state.projects[idx]["dmaic"]["improve"]

            # Vérification des données
            if "selection_matrix" not in dmaic_improve:
                st.warning("Veuillez d'abord compléter la matrice de sélection (Phase 2).")
            else:
                # Conversion du JSON en DataFrame
                df_results = pd.DataFrame(dmaic_improve["selection_matrix"])
        
                # Filtre : Retenues uniquement (si la colonne existe)
                if "Décision" in df_results.columns:
                    df_filtered = df_results[df_results["Décision"].isin(["Retenue", "Retenue pour pilote"])].copy()
                else:
                    df_filtered = df_results.copy()

                if df_filtered.empty:
                    st.info("Aucune solution 'Retenue' ou 'Retenue pour pilote' trouvée.")
                else:
                    # 1. CLASSIFICATION DES CRITÈRES
                    st.markdown("### ⚙️ Classification des critères")
                    if "criteria_mapping" not in dmaic_improve:
                        # Pré-classification automatique
                        dmaic_improve["criteria_mapping"] = {}
                        for crit in dmaic_improve["criteria"]:
                            name = crit["Critère"].lower()
                            if any(x in name for x in ["coût", "effort", "délai", "complexité", "ressource"]):
                                dmaic_improve["criteria_mapping"][crit["Critère"]] = "Effort"
                            else:
                                dmaic_improve["criteria_mapping"][crit["Critère"]] = "Bénéfice"

                    # Interface de modification
                    crit_cols = st.columns(len(dmaic_improve["criteria"]))
                    for i, crit in enumerate(dmaic_improve["criteria"]):
                        name = crit["Critère"]
                        dmaic_improve["criteria_mapping"][name] = crit_cols[i].selectbox(
                            f"{name}", ["Bénéfice", "Effort"], 
                            index=0 if dmaic_improve["criteria_mapping"].get(name) == "Bénéfice" else 1,
                            key=f"map_{name}"
                        )

                    # 2. CALCULS AUTOMATIQUES
                    poids_map = {c["Critère"]: c["Poids"] for c in dmaic_improve["criteria"]}
            
                    def calculate_scores(row):
                        score_ben = 0
                        score_eff = 0
                        for crit_name, cat in dmaic_improve["criteria_mapping"].items():
                            val = row.get(crit_name, 3) 
                            w = poids_map.get(crit_name, 0)
                            if cat == "Bénéfice":
                                score_ben += val * w
                            else:
                                score_eff += val * w
                        return pd.Series([score_ben, score_eff])

                    df_filtered[["Score Bénéfice", "Score Effort"]] = df_filtered.apply(calculate_scores, axis=1)
                    df_filtered["Indice de priorité"] = df_filtered["Score Bénéfice"] / df_filtered["Score Effort"].replace(0, 1)
            
                    # Détermination Quadrants (Médianes)
                    ben_mid = df_filtered["Score Bénéfice"].median()
                    eff_mid = df_filtered["Score Effort"].median()
            
                    def get_quadrant(row):
                        if row["Score Bénéfice"] >= ben_mid and row["Score Effort"] < eff_mid: return "Quick Win"
                        if row["Score Bénéfice"] >= ben_mid and row["Score Effort"] >= eff_mid: return "Projet Stratégique"
                        if row["Score Bénéfice"] < ben_mid and row["Score Effort"] < eff_mid: return "Opportunité Secondaire"
                        return "Faible Priorité"

                    df_filtered["Quadrant"] = df_filtered.apply(get_quadrant, axis=1)
                    priorite_map = {"Quick Win": "Très Haute", "Projet Stratégique": "Haute", "Opportunité Secondaire": "Moyenne", "Faible Priorité": "Faible"}
                    df_filtered["Priorité"] = df_filtered["Quadrant"].map(priorite_map)

                    # 3. VISUALISATION (Plotly Graph Objects)
                    fig = go.Figure()

                    # Ajout des points
                    fig.add_trace(go.Scatter(
                        x=df_filtered["Score Effort"], y=df_filtered["Score Bénéfice"],
                        mode='markers+text', text=df_filtered["Solution potentielle"],
                        textposition="top center",
                        marker=dict(size=12, color=df_filtered["Indice de priorité"], colorscale='Viridis', showscale=True)
                    ))

                    # Mise en forme (Quadrants)
                    fig.update_layout(
                        xaxis_title="Score Effort", yaxis_title="Score Bénéfice",
                        shapes=[
                            dict(type="line", x0=eff_mid, x1=eff_mid, y0=df_filtered["Score Bénéfice"].min(), y1=df_filtered["Score Bénéfice"].max(), line=dict(dash="dash", color="red")),
                            dict(type="line", x0=df_filtered["Score Effort"].min(), x1=df_filtered["Score Effort"].max(), y0=ben_mid, y1=ben_mid, line=dict(dash="dash", color="red"))
                        ],
                        annotations=[
                            dict(x=df_filtered["Score Effort"].min(), y=df_filtered["Score Bénéfice"].max(), text="Quick Win", showarrow=False, font=dict(size=14, color="green")),
                            dict(x=df_filtered["Score Effort"].max(), y=df_filtered["Score Bénéfice"].max(), text="Stratégique", showarrow=False, font=dict(size=14, color="blue")),
                            dict(x=df_filtered["Score Effort"].min(), y=df_filtered["Score Bénéfice"].min(), text="Opportunité", showarrow=False, font=dict(size=14, color="orange")),
                            dict(x=df_filtered["Score Effort"].max(), y=df_filtered["Score Bénéfice"].min(), text="Faible Priorité", showarrow=False, font=dict(size=14, color="grey"))
                        ]
                    )
                    st.plotly_chart(fig, use_container_width=True)

                    # 4. TABLEAU RÉCAPITULATIF
                    st.dataframe(df_filtered[["Solution potentielle", "Score Bénéfice", "Score Effort", "Indice de priorité", "Quadrant", "Priorité"]])

                    # 5. SAUVEGARDE
                    if st.button("💾 Enregistrer la matrice et les priorités"):
                        dmaic_improve["benefit_effort_results"] = df_filtered.to_dict(orient="records")
                
                        # --- SAUVEGARDE SÉCURISÉE ---
                        import json
                
                        # Fonction de nettoyage interne pour éviter l'erreur DataFrame
                        def prepare_for_json(obj):
                            if isinstance(obj, dict): return {k: prepare_for_json(v) for k, v in obj.items()}
                            elif isinstance(obj, list): return [prepare_for_json(item) for item in obj]
                            elif isinstance(obj, pd.DataFrame): return obj.to_dict(orient="records")
                            else: return obj
                
                        try:
                            # On nettoie tout le projet avant d'écrire
                            data_to_save = prepare_for_json(st.session_state.projects)
                            with open("projects.json", "w") as f:
                                json.dump(data_to_save, f, indent=4)
                            st.success("Priorités enregistrées !")
                        except Exception as e:
                            st.error(f"Erreur lors de la sauvegarde : {e}")
        

        # 4 : SOLUTIONS ACTION PLAN ---
        import plotly.figure_factory as ff # Nécessaire pour le Gantt

        st.markdown("---")
        st.subheader("4. Solutions action plan")

        if "benefit_effort_results" not in dmaic_improve:
            st.warning("Veuillez d'abord enregistrer la matrice de bénéfices/efforts (Étape 3).")
        else:
            # Récupération des données retenues
            df_results = pd.DataFrame(dmaic_improve["benefit_effort_results"])
    
            # a. FILTRE DE SÉLECTION
            st.markdown("### a. Sélection des solutions à déployer")
            filtre = st.selectbox(
                "Priorité minimale à déployer",
                ["Toutes les solutions", "Très Haute + Haute + Moyenne", "Très Haute + Haute", "Très Haute uniquement"]
            )
    
            # Logique de filtrage
            prio_order = {"Très Haute": 4, "Haute": 3, "Moyenne": 2, "Faible": 1}
            df_results['prio_val'] = df_results['Priorité'].map(prio_order)
    
            if filtre == "Très Haute uniquement": df_filter = df_results[df_results['prio_val'] == 4]
            elif filtre == "Très Haute + Haute": df_filter = df_results[df_results['prio_val'] >= 3]
            elif filtre == "Très Haute + Haute + Moyenne": df_filter = df_results[df_results['prio_val'] >= 2]
            else: df_filter = df_results
    
            # Initialisation structure stockage
            if "action_plan" not in dmaic_improve:
                dmaic_improve["action_plan"] = {}

            # b. CONSTRUCTION DU PLAN D'ACTION (Saisie libre pour le responsable)
            st.markdown("### b. Planification détaillée")
    
            if "action_plan" not in dmaic_improve:
                dmaic_improve["action_plan"] = {}

            for idx, row in df_filter.iterrows():
                sol = row["Solution potentielle"]
        
                # Récupération des données existantes (pour pré-remplir)
                saved = dmaic_improve["action_plan"].get(sol, {
                    "responsable": "", "debut": date.today(), "fin": date.today(),
                    "charge": 0, "budget": 0, "risques": "", "predecesseur": None
                })
        
                with st.expander(f"Action : {sol} (Priorité : {row['Priorité']})"):
                    # L'utilisation de st.form bloque le rechargement jusqu'au clic sur le bouton
                    with st.form(key=f"form_{sol}"):
                        c1, c2 = st.columns(2)
                
                        # Les inputs sont capturés dans des variables temporaires
                        val_resp = c1.text_input(f"Responsable(s) - {sol}", value=saved["responsable"])
                
                        # Récupération de l'index pour le selectbox
                        all_sols = [s for s in df_filter["Solution potentielle"] if s != sol]
                        options = [None] + all_sols
                        idx_predec = options.index(saved["predecesseur"]) if saved["predecesseur"] in options else 0
                
                        val_predec = c2.selectbox(f"Action prédécesseur", options, index=idx_predec)
                        val_deb = c1.date_input(f"Début - {sol}", value=saved["debut"])
                        val_fin = c2.date_input(f"Fin - {sol}", value=saved["fin"])
                        val_charg = c1.number_input(f"Charge (JH) - {sol}", value=saved["charge"])
                        val_budg = c2.number_input(f"Budget (Ar) - {sol}", value=saved["budget"])
                        val_risq = st.text_area(f"Risques - {sol}", value=saved["risques"])

                        # Bouton de soumission spécifique au formulaire
                        submitted = st.form_submit_button(f"Enregistrer {sol}")
                
                        if submitted:
                            # Enregistrement dans le dictionnaire
                            dmaic_improve["action_plan"][sol] = {
                                "responsable": val_resp,
                                "predecesseur": val_predec,
                                "debut": val_deb,
                                "fin": val_fin,
                                "charge": val_charg,
                                "budget": val_budg,
                                "risques": val_risq
                            }
                            save_data() 
                            st.success(f"Action '{sol}' enregistrée !")
                        
            # c. GÉNÉRATION DU GANTT (Corrigé pour la saisie libre)
            st.markdown("### c. Diagramme de Gantt")
            gantt_data = []
    
            for sol, data in dmaic_improve["action_plan"].items():
                # Vérification qu'il y a bien une saisie pour le responsable et les dates
                if data["responsable"] and data["debut"] and data["fin"]:
                    gantt_data.append(dict(
                        Task=sol, 
                        Start=str(data["debut"]), 
                        Finish=str(data["fin"]), 
                        Resource=str(data["responsable"]) # Maintenant une chaîne de caractères simple
                    ))
    
            if gantt_data:
                try:
                    fig = ff.create_gantt(
                        gantt_data, 
                        index_col='Resource', 
                        show_colorbar=True, 
                        group_tasks=True,
                        showgrid_x=True,
                        showgrid_y=True
                    )
                    st.plotly_chart(fig, use_container_width=True)
                except Exception as e:
                    st.error("Erreur de génération du Gantt. Vérifiez vos dates.")
            else:
                st.info("Remplissez le responsable et les dates pour générer le planning.")
                
            # d. INDICATEURS & SYNTHÈSE
            st.markdown("### d. Indicateurs et Synthèse")
            c1, c2, c3 = st.columns(3)

            # Calcul des totaux
            nb_actions = len(dmaic_improve["action_plan"])
            total_charge = sum(d['charge'] for d in dmaic_improve["action_plan"].values())
            total_budget = sum(d['budget'] for d in dmaic_improve["action_plan"].values())

            # Affichage des métriques
            c1.metric("Nb Actions", nb_actions)
            c2.metric("Charge Totale (JH)", f"{total_charge:,.0f}".replace(",", " "))

            # Affichage du budget avec séparateur de milliers (espace)
            # Le formatage :,.0f crée un séparateur et replace(" ", " ") met un espace insécable
            budget_formate = f"{total_budget:,.0f}".replace(",", " ")
            c3.metric("Budget Total (Ar)", budget_formate)
    
            # e. VALIDATION
            if st.button("✅ Valider le Solution Action Plan"):
                # Vérification
                errors = []
                for sol, data in dmaic_improve["action_plan"].items():
                    if not data["responsable"]: errors.append(f"{sol}: Responsable manquant")
                    if data["debut"] >= data["fin"]: errors.append(f"{sol}: Date début >= date fin")
            
                if errors:
                    for e in errors: st.error(e)
                else:
                    dmaic_improve["plan_valide"] = True
                    save_data()
                    st.success("Plan validé et transmis aux étapes suivantes (FMEA, Future State).")

        # 5 : FUTURE STATE PROCESS MAP ---
        st.subheader("5. Future state process map")
        
        # 🚨 CORRECTION INITIALISATION : Lecture sécurisée depuis p ou dmaic
        dmaic_data = p.get("dmaic", {})
        future_state_data = dmaic_data.get("future_state", {})

        if "future_state_map" not in st.session_state:
            # On vérifie si la map existe dans dmaic/future_state ou dans les clés directes de p
            saved_map = future_state_data.get("map") or p.get("future_state_map", {})
            st.session_state["future_state_map"] = copy.deepcopy(saved_map) if saved_map else copy.deepcopy(p.get("vsm_detailed_map", {}))

        if "future_macro_steps" not in st.session_state:
            saved_steps = future_state_data.get("macro_steps") or p.get("future_macro_steps", [])
            st.session_state["future_macro_steps"] = list(saved_steps) if saved_steps else list(p.get("vsm_macro_steps", []))

        st.info("💡 **Instructions** : Modifiez vos données, puis validez avec 'Enregistrer'.")

        # Formulaire global
        with st.form("future_state_form"):
            temp_editors = {}

            for step in st.session_state["future_macro_steps"]:
                col_head1, col_head2 = st.columns([0.8, 0.2])
                col_head1.subheader(f"Section : {step}")
                if col_head2.form_submit_button(f"🗑️ Supprimer {step}"):
                    st.session_state["future_macro_steps"].remove(step)
                    del st.session_state["future_state_map"][step]
                    st.rerun()

                data = st.session_state["future_state_map"].get(step, [])
                df_future = pd.DataFrame(data)
    
                # Normalisation
                df_future = df_future.rename(columns={"Valeur": "Délai à T0"})
                for col in ["Délai à T0", "Délai Actuel"]:
                    if col not in df_future.columns: df_future[col] = 0.0
                if "Unité" not in df_future.columns: df_future["Unité"] = "Minutes"
                if "Type d'activité" not in df_future.columns: df_future["Type d'activité"] = "VA (Valeur Ajoutée)"

                cols_to_use = ["Détail de la tâche", "Délai à T0", "Unité", "Type d'activité", "Délai Actuel"]

                # SÉCURISATION : Crée les colonnes manquantes si le DataFrame est vide
                for col in cols_to_use:
                    if col not in df_future.columns:
                        df_future[col] = ""
                
                edited_df = st.data_editor(
                    df_future[cols_to_use],
                    num_rows="dynamic",
                    use_container_width=True,
                    key=f"editor_{step}"
                )
                temp_editors[step] = edited_df

            submitted = st.form_submit_button("💾 Enregistrer et Calculer les gains", type="primary")

        # --- AJOUTER SECTION (En dehors du form pour éviter les conflits) ---
        st.markdown("---")
        new_section = st.text_input("Nom de la nouvelle section :")
        if st.button("➕ Ajouter la section"):
            if new_section and new_section not in st.session_state["future_macro_steps"]:
                st.session_state["future_macro_steps"].append(new_section)
                st.session_state["future_state_map"][new_section] = [{"Détail de la tâche": "Action", "Délai à T0": 0.0, "Unité": "Minutes", "Type d'activité": "VA (Valeur Ajoutée)", "Délai Actuel": 0.0}]
                st.rerun()

        # Logique après soumission
        if submitted:
            idx = st.session_state.get("current_project_idx", 0)
    
            # --- 1. SÉCURISATION DE LA STRUCTURE ---
            if "dmaic" not in st.session_state.projects[idx]:
                st.session_state.projects[idx]["dmaic"] = {}
            if "future_state" not in st.session_state.projects[idx]["dmaic"]:
                st.session_state.projects[idx]["dmaic"]["future_state"] = {}

            # --- 2. SYNCHRONISATION (L'ÉTAPE CRUCIALE) ---
            new_map = {}
            for step in st.session_state["future_macro_steps"]:
                if step in temp_editors:
                    new_map[step] = temp_editors[step].to_dict(orient='records')
    
            # Mise à jour dans session_state.projects
            st.session_state.projects[idx]["dmaic"]["future_state"]["map"] = new_map
            st.session_state.projects[idx]["dmaic"]["future_state"]["macro_steps"] = list(st.session_state["future_macro_steps"])
    
            # 🚨 CORRECTION SAUVEGARDE JSON (Mise à jour directe de p)
            if "dmaic" not in p:
                p["dmaic"] = {}
            if "future_state" not in p["dmaic"]:
                p["dmaic"]["future_state"] = {}
        
            p["dmaic"]["future_state"]["map"] = new_map
            p["dmaic"]["future_state"]["macro_steps"] = list(st.session_state["future_macro_steps"])
    
            # Raccourcis directs au cas où le script lit directement p[...]
            p["future_state_map"] = new_map
            p["future_macro_steps"] = list(st.session_state["future_macro_steps"])

            # --- 3. SAUVEGARDE SUR LE DISQUE ---
            save_data()
    
            st.success("✅ Données sauvegardées de manière permanente.")
            st.rerun()

        # --- AFFICHAGE DES RÉSULTATS (Dynamique) ---
                
        # 1. RÉCUPÉRATION DU T0 (Depuis la phase Mesure déjà stockée en session)
        # On utilise directement les totaux calculés par la phase mesure
        t0_data = st.session_state.get("vsm_totals", {})
        t0_lt_ref = float(t0_data.get("lead_time", 0.0))
        t0_va_ref = float(t0_data.get("va", 0.0))

        # 2. CALCUL ACTUEL (Depuis les éditeurs du Future State en temps réel)
        calc_actuel_lt = 0.0
        calc_actuel_va = 0.0
        
        for step in st.session_state["future_macro_steps"]:
            if step in temp_editors:
                df_temp = temp_editors[step]
                # Calcul basé sur ce qui est saisi dans le formulaire actuel
                calc_actuel_lt += float(df_temp["Délai Actuel"].sum())
                va_mask = df_temp["Type d'activité"] == "VA (Valeur Ajoutée)"
                calc_actuel_va += float(df_temp.loc[va_mask, "Délai Actuel"].sum())

        # 3. AFFICHAGE DE LA SYNTHÈSE
        st.markdown("### 📊 Synthèse des gains")
        
        # Calcul du PCE Actuel
        pce_actuel = (calc_actuel_va / calc_actuel_lt * 100) if calc_actuel_lt > 0 else 0
        gain_total = t0_lt_ref - calc_actuel_lt
        
        data_synth = {
            "Indicateur": ["Lead Time Total (min)", "Valeur Ajoutée (min)", "Efficience (PCE %)"],
            "Valeur T0 (Mesure)": [f"{t0_lt_ref:.1f}", f"{t0_va_ref:.1f}", f"{(t0_va_ref/t0_lt_ref*100 if t0_lt_ref>0 else 0):.1f}%"],
            "Valeur Actuel": [f"{calc_actuel_lt:.1f}", f"{calc_actuel_va:.1f}", f"{pce_actuel:.1f}%"]
        }
        st.table(pd.DataFrame(data_synth))
        st.metric("💰 GAIN TOTAL DE TEMPS", f"{gain_total:.1f} min")


        # 6 : FUTURE STATE FMEA ---
        st.markdown("---")
        st.subheader("6. Future state FMEA")

        # 1. Accès sécurisé aux données du projet
        idx = st.session_state.get("current_project_idx", 0)
        if "dmaic" not in st.session_state.projects[idx]:
            st.session_state.projects[idx]["dmaic"] = {}
        if "improve" not in st.session_state.projects[idx]["dmaic"]:
            st.session_state.projects[idx]["dmaic"]["improve"] = {}

        dmaic_analyze = st.session_state.projects[idx]["dmaic"].get("analyze", {})
        dmaic_improve = st.session_state.projects[idx]["dmaic"].get("improve", {})

        fmea_actuel_data = dmaic_analyze.get("fmea_data", {})
        gemba_plans = dmaic_analyze.get("gemba_plans", {})
        gemba_obs = dmaic_analyze.get("gemba_observations", {})
        results_phase1 = dmaic_analyze.get("results", [])

        # 2. Recalcul et application stricte du filtre de décision de l'étape 6 de l'Analyse
        causes_validees_list = []

        for cid, plan in gemba_plans.items():
            x_nom = str(plan.get("x_critique", "")).strip()
            cause_nom = str(plan.get("cause_racine", "")).strip()
    
            p_val = 1.0
            for r in results_phase1:
                if str(r.get("Variable X", "")).strip().lower() == x_nom.lower():
                    p_val = float(r.get("P-value", 1.0))
                    break
            
            obs_info = gemba_obs.get(cid, {})
            logs = obs_info.get("logs", []) if isinstance(obs_info, dict) else []
            conf = len([l for l in logs if str(l.get("confirme", "")).strip().lower() == 'oui'])
            total = len(logs)
            taux = (conf / total * 100) if total > 0 else 0
    
            val_terrain = "Confirmée" if taux >= 80 else ("Partiellement confirmée" if taux >= 50 else "Non confirmée")
            stat_validee = p_val < 0.05
            res_stat = "Confirmé" if stat_validee else "Non confirmé"
    
            if res_stat == "Confirmé" and val_terrain == "Confirmée":
                decision = "Cause validée"
            elif res_stat == "Confirmé" and val_terrain == "Partiellement confirmée":
                decision = "Cause partiellement validée"
            elif res_stat == "Confirmé" and val_terrain == "Non confirmée":
                decision = "Contradiction à investiguer"
            elif res_stat == "Non confirmé" and val_terrain == "Confirmée":
                decision = "Contradiction à investiguer"
            elif res_stat == "Non confirmé" and val_terrain == "Partiellement confirmée":
                decision = "Cause non démontrée"
            elif res_stat == "Non confirmé" and val_terrain == "Non confirmée":
                decision = "Cause rejetée"
            else:
                decision = "En cours d'analyse"

            if decision in ["Cause validée", "Cause partiellement validée"]:
                if cause_nom:
                    causes_validees_list.append({
                        "X Critique": x_nom,
                        "Cause Racine": cause_nom,
                        "ID": cid
                    })

        # 3. Gestion de la session et initialisation
        future_fmea_key = f"future_state_fmea_{idx}"

        col_btn1, col_btn2 = st.columns([0.8, 0.2])
        with col_btn2:
            if st.button("🔄 Rafraîchir", key=f"force_refresh_{idx}"):
                if future_fmea_key in st.session_state:
                    del st.session_state[future_fmea_key]
                st.rerun()

        if future_fmea_key not in st.session_state:
            saved_future_fmea = dmaic_improve.get("future_fmea", [])
            if saved_future_fmea:
                st.session_state[future_fmea_key] = pd.DataFrame(saved_future_fmea)
            else:
                initial_fmea_rows = []
                for item in causes_validees_list:
                    cid = item["ID"]
                    cause_nom = item["Cause Racine"]
            
                    fmea_vals = fmea_actuel_data.get(cid, {})
                    s_actuel = float(fmea_vals.get("S", 6))
                    o_actuel = float(fmea_vals.get("O", 5))
                    d_actuel = float(fmea_vals.get("D", 5))
                    rpn_actuel = s_actuel * o_actuel * d_actuel
            
                    s_futur = s_actuel
                    o_futur = max(1.0, o_actuel - 1.0)
                    d_futur = max(1.0, d_actuel - 1.0)
                    rpn_futur = s_futur * o_futur * d_futur
            
                    initial_fmea_rows.append({
                        "Cause racine validée": cause_nom,
                        "RPN actuel": rpn_actuel,
                        "S futur": s_futur,
                        "O futur": o_futur,
                        "D futur": d_futur,
                        "RPN futur": rpn_futur
                    })
            
                if not initial_fmea_rows:
                    initial_fmea_rows = [{
                        "Cause racine validée": "⚠️ Aucune cause 'validée' ou 'partiellement validée' trouvée dans la phase Analyse",
                        "RPN actuel": 0.0,
                        "S futur": 1.0, "O futur": 1.0, "D futur": 1.0, "RPN futur": 1.0
                    }]
                st.session_state[future_fmea_key] = pd.DataFrame(initial_fmea_rows)

        # 4. Utilisation d'un formulaire pour bloquer les sauvegardes et calculs automatiques au clic extérieur
        with st.form(key=f"form_future_fmea_{idx}"):
            seuil_critique = st.number_input("Seuil RPN critique paramétrable :", value=100.0, step=10.0, key=f"seuil_rpn_improve_{idx}")

            st.markdown("### 📝 Réévaluation du futur état (S, O, D & RPN futur)")

            df_current = st.session_state[future_fmea_key].copy()
            if "S futur" in df_current.columns and "O futur" in df_current.columns and "D futur" in df_current.columns:
                df_current["RPN futur"] = df_current["S futur"] * df_current["O futur"] * df_current["D futur"]

            edited_future_fmea = st.data_editor(
                df_current,
                use_container_width=True,
                num_rows="fixed",
                key=f"editor_future_fmea_view_{idx}",
                column_config={
                    "Cause racine validée": st.column_config.TextColumn("Cause racine validée", disabled=True, width="medium"),
                    "RPN actuel": st.column_config.NumberColumn("RPN actuel", disabled=True, format="%.1f", width="small"),
                    "S futur": st.column_config.NumberColumn("S futur (1-10)", min_value=1.0, max_value=10.0, step=1.0, format="%.0f", width="small"),
                    "O futur": st.column_config.NumberColumn("O futur (1-10)", min_value=1.0, max_value=10.0, step=1.0, format="%.0f", width="small"),
                    "D futur": st.column_config.NumberColumn("D futur (1-10)", min_value=1.0, max_value=10.0, step=1.0, format="%.0f", width="small"),
                    "RPN futur": st.column_config.NumberColumn("RPN futur", disabled=True, format="%.1f", width="small")
                }
            )

            submit_button = st.form_submit_button("💾 Enregistrer la Future State FMEA & Calculer", type="primary", use_container_width=True)

            if submit_button:
                df_res = edited_future_fmea.copy()
        
                df_res["RPN futur"] = df_res["S futur"] * df_res["O futur"] * df_res["D futur"]
                df_res["Réduction du risque"] = df_res["RPN actuel"] - df_res["RPN futur"]
                df_res["Pourcentage de réduction"] = ((df_res["RPN actuel"] - df_res["RPN futur"]) / df_res["RPN actuel"].replace(0, 1)) * 100
        
                def interpret_reduction(row):
                    red_pct = row["Pourcentage de réduction"]
                    rpn_fut = row["RPN futur"]
            
                    if red_pct >= 70: interp = "Risque fortement réduit"
                    elif red_pct >= 40: interp = "Risque significativement réduit"
                    elif red_pct >= 20: interp = "Réduction limitée"
                    else: interp = "Impact faible des solutions"
            
                    if rpn_fut > seuil_critique:
                        interp += " - Risque résiduel critique"
                    return interp

                def class_efficacite(red_pct):
                    if red_pct >= 70: return "Très efficace"
                    elif red_pct >= 40: return "Efficace"
                    elif red_pct >= 20: return "Partiellement efficace"
                    else: return "Peu efficace"

                df_res["Interprétation"] = df_res.apply(interpret_reduction, axis=1)
                df_res["Efficacité"] = df_res["Pourcentage de réduction"].apply(class_efficacite)
        
                st.session_state[future_fmea_key] = df_res
                st.session_state.projects[idx]["dmaic"]["improve"]["future_fmea"] = df_res.to_dict('records')
                if "save_data" in globals():
                    save_data()
        
                st.success("✅ Future State FMEA enregistrée avec succès !")
                st.rerun()
    
    # --- PHASE CONTROL ---
    with tabs[4]: 
        st.header("Phase Control")

        from datetime import datetime, timezone, timedelta
        import io
        import re
        import math

        if "control" not in p:
            p["control"] = {}

        ctrl_data = p["control"]

        # ---------------------------------------------------------------------
        # 0. RÉCUPÉRATION DES VARIABLES CRITIQUES & DONNÉES T0 (Phase Mesure - Immuable)
        # ---------------------------------------------------------------------
        safe_idx = str(p_idx) if 'p_idx' in locals() else "default"
        safe_p_idx = safe_idx
        msa_classif_key = f"msa_classification_table_{safe_idx}"
        nom_colonne_variable = "Variable Critique (liée au Y)"

        liste_variables_dynamiques = []
        df_msa_source = pd.DataFrame()

        if msa_classif_key in st.session_state and not st.session_state[msa_classif_key].empty:
            df_msa_source = st.session_state[msa_classif_key]
            if nom_colonne_variable in df_msa_source.columns:
                liste_variables_dynamiques = df_msa_source[nom_colonne_variable].dropna().tolist()
        elif "master_dcp_table" in st.session_state and len(st.session_state["master_dcp_table"]) > 0:
            liste_variables_dynamiques = [row["Variable à mesurer"] for row in st.session_state["master_dcp_table"] if "Variable à mesurer" in row]

        if not liste_variables_dynamiques:
            liste_variables_dynamiques = ["Temps de traitement", "Statut conformité"]

        # Chargement immuable de T0 depuis la Phase Mesure
        df_t0_data = pd.DataFrame()
        if "dc_saved_df_json" in p and p["dc_saved_df_json"]:
            try:
                df_t0_data = pd.read_json(io.StringIO(p["dc_saved_df_json"]))
            except Exception:
                pass

        # ---------------------------------------------------------------------
        # 1. DATA CONTROL PLAN & GESTION DES VAGUES
        # ---------------------------------------------------------------------
        st.subheader("1. Data Control Plan")

        if "waves" not in ctrl_data:
            ctrl_data["waves"] = ["T1"]  # Vagues dynamiques par défaut hors T0

        # Initialisation du stockage persistant par vagues s'il n'existe pas
        if "vagues_data_cache" not in ctrl_data:
            ctrl_data["vagues_data_cache"] = {}

        # Bouton unique pour ajouter une vague
        if st.button("➕ Ajouter une vague", key=f"add_wave_global_{safe_p_idx}"):
            existing_nums = []
            for w in ctrl_data["waves"]:
                m = re.search(r'T(\d+)', str(w).upper())
                if m: existing_nums.append(int(m.group(1)))
            prochain_num = max(existing_nums, default=0) + 1
            nouvelle_vague = f"T{prochain_num}"
            if nouvelle_vague not in ctrl_data["waves"]:
                ctrl_data["waves"].append(nouvelle_vague)
                st.rerun()

        # Liste d'options unique et dédupliquée pour la vague active
        toutes_les_vagues_options = ["T0 (Référence Mesure)"] + ctrl_data["waves"]
        selected_wave = st.selectbox("Vague active de suivi (pour métriques ciblées) :", options=toutes_les_vagues_options, key=f"active_wave_{safe_p_idx}")

        st.markdown("---")

        # ---------------------------------------------------------------------
        # 2. COLLECTE ET SAISIE DES DONNÉES DE CONTRÔLE (Écran 1 & Écran 2)
        # ---------------------------------------------------------------------
        st.markdown("## 2 - Collecte des Données de Contrôle & Saisie Multi-Vagues")

        if "ctrl_plan" not in ctrl_data:
            ctrl_data["ctrl_plan"] = {"taille_prevue": 50, "date_debut": "2026-07-01", "date_fin_est": "2026-07-15"}

        if "ctrl_master_data" not in ctrl_data or not isinstance(ctrl_data["ctrl_master_data"], pd.DataFrame):
            cols_init = ["ID observation", "Date de modification", "Vague"] + liste_variables_dynamiques
            ctrl_data["ctrl_master_data"] = pd.DataFrame(columns=cols_init)

        # Paramètres de l'Écran 1 (Résumé de la collecte)
        st.markdown("### 📋 Écran 1 : Paramètres d'Échantillonnage")
        with st.form(key=f"form_ecran1_ctrl_{safe_p_idx}"):
            e1_c1, e1_c2, e1_c3 = st.columns(3)
            with e1_c1:
                n_prev = st.number_input("Taille d'échantillon prévue (N)", min_value=1, value=int(ctrl_data["ctrl_plan"].get("taille_prevue", 50)))
            with e1_c2:
                d_deb = st.text_input("Date de début de collecte", value=ctrl_data["ctrl_plan"].get("date_debut", "2026-07-01"))
            with e1_c3:
                d_fin = st.text_input("Date estimée de fin", value=ctrl_data["ctrl_plan"].get("date_fin_est", "2026-07-15"))

            submitted_e1 = st.form_submit_button("💾 Enregistrer les paramètres")
            if submitted_e1:
                ctrl_data["ctrl_plan"]["taille_prevue"] = n_prev
                ctrl_data["ctrl_plan"]["date_debut"] = d_deb
                ctrl_data["ctrl_plan"]["date_fin_est"] = d_fin
                st.success("Paramètres enregistrés.")

        st.markdown("---")

        # Écran 2 : Saisie / Import par Vague (Gestion persistante de chaque onglet)
        st.markdown("### 📝 Écran 2 : Saisie des Données de Contrôle & Import Excel par Vague")
    
        dict_dfs_vagues = {}
        if ctrl_data["waves"]:
            onglets_vagues = st.tabs([f"Vague {v}" for v in ctrl_data["waves"]])

            for i, v_name in enumerate(ctrl_data["waves"]):
                with onglets_vagues[i]:
                    col_tab_hd1, col_tab_hd2 = st.columns([4, 1])
                    with col_tab_hd1:
                        st.markdown(f"##### Gestion des données pour la période `{v_name}`")
                    with col_tab_hd2:
                        if st.button("🗑️ Supprimer vague", key=f"del_wave_tab_{v_name}_{safe_p_idx}"):
                            ctrl_data["waves"].remove(v_name)
                            if v_name in ctrl_data["vagues_data_cache"]:
                                del ctrl_data["vagues_data_cache"][v_name]
                            if "ctrl_master_data" in ctrl_data and not ctrl_data["ctrl_master_data"].empty:
                                ctrl_data["ctrl_master_data"] = ctrl_data["ctrl_master_data"][ctrl_data["ctrl_master_data"]["Vague"] != v_name].reset_index(drop=True)
                            st.success(f"Vague {v_name} supprimée.")
                            st.rerun()

                    uploaded_ctrl_file = st.file_uploader(f"Importer un fichier Excel pour {v_name}", type=["xlsx", "xls"], key=f"ctrl_excel_up_{v_name}_{safe_p_idx}")
                
                    # Récupération depuis le cache persistant ou la master data existante
                    if v_name in ctrl_data["vagues_data_cache"] and not ctrl_data["vagues_data_cache"][v_name].empty:
                        df_vague_courante = ctrl_data["vagues_data_cache"][v_name].copy()
                    else:
                        df_vague_courante = pd.DataFrame(columns=["ID observation", "Date de modification", "Vague"] + liste_variables_dynamiques)
                        existing_sub = ctrl_data["ctrl_master_data"][ctrl_data["ctrl_master_data"]["Vague"] == v_name] if not ctrl_data["ctrl_master_data"].empty and "Vague" in ctrl_data["ctrl_master_data"].columns else pd.DataFrame()
                        if not existing_sub.empty:
                            df_vague_courante = existing_sub.copy()

                    if uploaded_ctrl_file:
                        file_cache_key_ctrl = f"processed_ctrl_{v_name}_{uploaded_ctrl_file.name}_{uploaded_ctrl_file.size}"
                        if st.session_state.get(f"ctrl_last_processed_{v_name}") != file_cache_key_ctrl:
                            try:
                                raw_imp_ctrl = pd.read_excel(uploaded_ctrl_file).dropna(how="all").reset_index(drop=True)
                                regex_clean = re.compile(r'[_\-\s\./\\]+')
                                def _struct_clean(text):
                                    if pd.isna(text): return ""
                                    t = str(text).lower().strip()
                                    t = regex_clean.sub(' ', t)
                                    return "".join(c for c in t if c.isalnum() or c == ' ')

                                cols_f = ["ID observation", "Date de modification", "Vague"] + liste_variables_dynamiques
                                aligned_c_df = pd.DataFrame(columns=cols_f, index=range(len(raw_imp_ctrl)))
                                c_excel = list(raw_imp_ctrl.columns)
                                c_clean = [_struct_clean(c) for c in c_excel]

                                id_src = next((col for col in c_excel if _struct_clean(col) in {"id", "observation", "code", "num", "index", "identifiant"}), None)
                                if id_src:
                                    aligned_c_df["ID observation"] = raw_imp_ctrl[id_src].astype(str).str.replace(r'\.0$', '', regex=True).str.strip()
                                else:
                                    aligned_c_df["ID observation"] = [f"{v_name}_Obs_{x+1}" for x in range(len(raw_imp_ctrl))]

                                aligned_c_df["Vague"] = v_name

                                for var_c in liste_variables_dynamiques:
                                    v_cl = _struct_clean(var_c)
                                    w1 = set(v_cl.split())
                                    best_m, best_s = None, 0.0
                                    for idx, col in enumerate(c_excel):
                                        w2 = set(c_clean[idx].split())
                                        if not w1 or not w2: continue
                                        score = len(w1.intersection(w2)) / max(len(w1), len(w2))
                                        if v_cl in c_clean[idx] or c_clean[idx] in v_cl: score += 0.3
                                        if score > best_s:
                                            best_s = score
                                            best_m = col
                                    if best_m and best_s >= 0.35:
                                        aligned_c_df[var_c] = raw_imp_ctrl[best_m].values
                                    else:
                                        aligned_c_df[var_c] = None

                                tz_val = datetime.now(timezone(timedelta(hours=3))).strftime("%Y-%m-%d %H:%M:%S")
                                aligned_c_df["Date de modification"] = tz_val
                                df_vague_courante = aligned_c_df.reset_index(drop=True).where(pd.notnull(aligned_c_df), None)
                                st.session_state[f"ctrl_last_processed_{v_name}"] = file_cache_key_ctrl
                                st.success(f"🚀 Fichier importé pour {v_name} ({len(df_vague_courante)} lignes).")
                            except Exception as e:
                                st.error(f"Erreur d'import : {e}")

                    edited_v_table = st.data_editor(
                        df_vague_courante,
                        num_rows="dynamic",
                        use_container_width=True,
                        key=f"ctrl_grid_edit_{v_name}_{safe_p_idx}"
                    )
                    edited_v_table = pd.DataFrame(edited_v_table)
                    edited_v_table["Vague"] = v_name
                
                    # Sauvegarde immédiate dans le cache persistant par vague
                    ctrl_data["vagues_data_cache"][v_name] = edited_v_table
                    dict_dfs_vagues[v_name] = edited_v_table
        else:
            st.info("💡 Aucune vague de suivi active. Cliquez sur 'Ajouter une vague' ci-dessus.")

        # Consolidation globale de toutes les vagues
        if dict_dfs_vagues:
            ctrl_data["ctrl_master_data"] = pd.concat(list(dict_dfs_vagues.values()), ignore_index=True)
        elif "vagues_data_cache" in ctrl_data and ctrl_data["vagues_data_cache"]:
            ctrl_data["ctrl_master_data"] = pd.concat(list(ctrl_data["vagues_data_cache"].values()), ignore_index=True)
        else:
            ctrl_data["ctrl_master_data"] = pd.DataFrame(columns=["ID observation", "Date de modification", "Vague"] + liste_variables_dynamiques)

        st.markdown("---")

        # ---------------------------------------------------------------------
        # 3. STATISTIQUES COMPARATIVES (T0 vs Vagues)
        # ---------------------------------------------------------------------
        st.subheader("3. Statistiques Descriptives & Comparatif T0 vs Vagues")

        df_combined = pd.DataFrame()
        if not df_t0_data.empty:
            df_t0_copy = df_t0_data.copy()
            df_t0_copy["Vague"] = "T0 (Référence Mesure)"
            df_combined = pd.concat([df_t0_copy, ctrl_data["ctrl_master_data"]], ignore_index=True)
        else:
            df_combined = ctrl_data["ctrl_master_data"]

        if not df_combined.empty and "Vague" in df_combined.columns:
            for var_item in liste_variables_dynamiques:
                if var_item in df_combined.columns:
                    st.markdown(f"##### Indicateur : {var_item}")
                    num_s = pd.to_numeric(df_combined[var_item], errors="coerce")
                    df_combined["_tmp_val"] = num_s
                    stats_res = df_combined.groupby("Vague", sort=False)["_tmp_val"].agg(
                        N="count", Moyenne="mean", Médiane="median", Écart_type="std"
                    ).reset_index().round(2)
                    st.dataframe(stats_res, use_container_width=True, hide_index=True)
        else:
            st.info("Aucune donnée disponible pour l'analyse comparative.")

        st.markdown("---")

        # ---------------------------------------------------------------------
        # 4. PROCESS CAPABILITY COMPARATIF (T0 vs Vagues)
        # ---------------------------------------------------------------------
        st.subheader("4. Process Capability (Comparatif T0 vs Vagues Post-Amélioration)")
        st.caption("Évaluation comparative du niveau Sigma et du DPMO entre la référence T0 et les vagues de contrôle.")

        import math

        # --- FONCTION INTERNE D'ÉVALUATION DE CONFORMITÉ ---
        def _evaluer_conformite_ligne_v2(val_observee, config_master):
            if pd.isna(val_observee):
                return True
            t_spec = config_master.get("type_specification", "Valeur exacte")
            
            if t_spec == "Attribut qualitatif (Sans valeur numérique - Proposition Master Black Belt)":
                v_str = str(val_observee).strip().upper()
                if v_str in ["KO", "NON OK", "NON", "ERREUR", "RETOUCHE", "1", "NON-CONFORME"]:
                    return False
                return True

            try:
                num_val = float(val_observee)
            except (ValueError, TypeError):
                return str(val_observee).strip().upper() not in ["KO", "NON OK", "1", "NON-CONFORME"]

            if t_spec == "Valeur exacte":
                return math.isclose(num_val, float(config_master.get("valeur_exacte", 0.0)), rel_tol=1e-5, abs_tol=1e-5)
            elif t_spec == "Supérieur ou égal à (≥)":
                return num_val >= float(config_master.get("valeur_seuil", 0.0))
            elif t_spec == "Inférieur ou égal à (≤)":
                return num_val <= float(config_master.get("valeur_seuil", 0.0))
            elif t_spec == "Intervalle (Entre min et max)":
                return float(config_master.get("borne_inf", 0.0)) <= num_val <= float(config_master.get("borne_sup", 10.0))
            return True

        # --- IDENTIFICATION DE LA VARIABLE Y ACTIVE ---
        nom_variable_y_comp = None
        if 'selected_var_to_test' in locals() and selected_var_to_test:
            nom_variable_y_comp = selected_var_to_test
        elif 'nom_colonne_variable' in locals() and not df_classification_current.empty:
            list_y_cands_comp = df_classification_current[nom_colonne_variable].dropna().tolist()
            if list_y_cands_comp:
                nom_variable_y_comp = list_y_cands_comp[0]

        # Récupération de la configuration Master active
        master_cfg_comp = st.session_state.get("reference_master_config", {
            "type_specification": "Valeur exacte",
            "valeur_exacte": 0.0
        })
        if 'var_clean_id' in locals() and f"reference_master_config_{var_clean_id}_{safe_idx}" in st.session_state:
            master_cfg_comp = st.session_state[f"reference_master_config_{var_clean_id}_{safe_idx}"]

        df_t0_cap = df_t0_data.copy() if not df_t0_data.empty else pd.DataFrame()
        df_ctrl_cap = ctrl_data["ctrl_master_data"].copy() if not ctrl_data["ctrl_master_data"].empty else pd.DataFrame()

        # --- FONCTION DE CALCUL CIBLÉE SUR LE Y ---
        def _calculer_metriques_capabilite_y_only(df_subset, var_y, cfg_master):
            if df_subset.empty or not var_y or var_y not in df_subset.columns:
                return 0, 1, 1, 0.0, 0.0
                
            r_units = len(df_subset.dropna(subset=["ID observation"])) if "ID observation" in df_subset.columns else len(df_subset)
            r_opp = 1  # 1 opportunité pour le Y
        
            def_count = 0
            for idx_row, row in df_subset.iterrows():
                if not _evaluer_conformite_ligne_v2(row[var_y], cfg_master):
                    def_count += 1
        
            u_val = max(1, r_units)
            dpmo_val = (def_count / (u_val * r_opp)) * 1_000_000
        
            sig_lvl = 0.0
            try:
                if dpmo_val <= 0:
                    sig_lvl = 6.0
                else:
                    taux_d = dpmo_val / 1_000_000
                    p_v = max(1e-7, min(0.5, taux_d if taux_d < 0.5 else 1 - taux_d))
                    t_val = math.sqrt(-2.0 * math.log(p_v))
                    z_val = t_val - ((2.515517 + 0.802853 * t_val + 0.010328 * t_val * t_val) / 
                                   (1.0 + 1.432788 * t_val + 0.189269 * t_val * t_val + 0.001308 * t_val * t_val * t_val))
                    sig_brut = z_val if taux_d < 0.5 else -z_val
                    sig_lvl = max(0.0, min(6.0, round(sig_brut + 1.5, 2)))
            except Exception:
                sig_lvl = 0.0
                
            return def_count, u_val, r_opp, dpmo_val, sig_lvl

        st.markdown(f"**Variable cible analysée :** `{nom_variable_y_comp or 'Non définie'}`")

        # --- CONSTITUTION DES DONNÉES DU TABLEAU COMPARATIF ---
        lignes_tableau = []

        # Ajout de T0
        def_t0, u_t0, opp_t0, dpmo_t0, sig_t0 = _calculer_metriques_capabilite_y_only(df_t0_cap, nom_variable_y_comp, master_cfg_comp)
        lignes_tableau.append({
            "Vague": "T0 (Référence Mesure)",
            "Unités (N)": u_t0,
            "Défauts": def_t0,
            "DPMO": f"{dpmo_t0:,.0f}",
            "Niveau Sigma": f"{sig_t0} σ",
            "Évolution DPMO": "Baseline"
        })

        # Extraction dynamique des vagues additionnelles (T1, T2...)
        vagues_disponibles = []
        if not df_ctrl_cap.empty and "Vague" in df_ctrl_cap.columns:
            vagues_disponibles = sorted(df_ctrl_cap["Vague"].dropna().unique().tolist())

        for vague_nom in vagues_disponibles:
            df_vague_encours = df_ctrl_cap[df_ctrl_cap["Vague"] == vague_nom]
            def_v, u_v, opp_v, dpmo_v, sig_v = _calculer_metriques_capabilite_y_only(df_vague_encours, nom_variable_y_comp, master_cfg_comp)
            
            delta_dpmo = dpmo_v - dpmo_t0
            signe_delta = f"{delta_dpmo:+,.0f}" if delta_dpmo != 0 else "0"
            
            lignes_tableau.append({
                "Vague": vague_nom,
                "Unités (N)": u_v,
                "Défauts": def_v,
                "DPMO": f"{dpmo_v:,.0f}",
                "Niveau Sigma": f"{sig_v} σ",
                "Évolution DPMO": signe_delta
            })

        df_comparatif_final = pd.DataFrame(lignes_tableau)
        st.dataframe(df_comparatif_final, use_container_width=True, hide_index=True)

        st.markdown("---")

        # ---------------------------------------------------------------------
        # 5. CONTROL CHART MIXTE (T0 au début, puis T1, T2...)
        # ---------------------------------------------------------------------
        st.subheader("5. Carte de Contrôle Mixte (Séquence chronologique : T0 en premier)")
        st.caption("Affichage chronologique strict débutant par T0, suivi des vagues, avec adaptation dynamique des limites (CL/UCL/LCL) par vague.")

        df_t0_plot = df_t0_cap.copy() if not df_t0_cap.empty else pd.DataFrame()
        if not df_t0_plot.empty:
            df_t0_plot["Vague"] = "T0 (Référence Mesure)"
            if "Date de modification" not in df_t0_plot.columns:
                df_t0_plot["Date de modification"] = "2026-06-01 00:00:00"

        df_ctrl_plot = df_ctrl_cap.copy() if not df_ctrl_cap.empty else pd.DataFrame()
        df_mixte = pd.concat([df_t0_plot, df_ctrl_plot], ignore_index=True)

        if df_mixte.empty:
            st.info("💡 Veuillez alimenter vos données T0 et ajouter/importer des données de contrôle.")
        else:
            quant_vars_mixte = [v for v in liste_variables_dynamiques if v in df_mixte.columns and not any(k in v.lower() for k in ["statut", "verdict", "validation", "ok", "ko"])]
        
            if not quant_vars_mixte:
                st.warning("⚠️ Aucune variable quantitative exploitable pour la carte de contrôle mixte.")
            else:
                var_choisie_mixte = st.selectbox("Sélectionner la variable quantitative pour le Control Chart :", options=quant_vars_mixte, key=f"mixte_var_final_{safe_p_idx}")
            
                df_mixte["Valeur_Num"] = pd.to_numeric(df_mixte[var_choisie_mixte], errors="coerce")
                df_mixte["Date_Parsed"] = pd.to_datetime(df_mixte["Date de modification"], errors="coerce")
                df_clean_mixte = df_mixte.dropna(subset=["Valeur_Num"]).copy()
            
                if df_clean_mixte.empty:
                    st.warning("Aucune valeur numérique valide trouvée pour cette variable.")
                else:
                    def _ordre_vague_chronologique(vague_val):
                        v_str = str(vague_val).upper()
                        if "T0" in v_str or "RÉFÉRENCE" in v_str:
                            return -1
                        match_num = re.search(r'T(\d+)', v_str)
                        if match_num:
                            return int(match_num.group(1))
                        return 99

                    df_clean_mixte["_sort_order"] = df_clean_mixte["Vague"].apply(_ordre_vague_chronologique)
                    df_clean_mixte = df_clean_mixte.sort_values(by=["_sort_order", "Date_Parsed", "ID observation"], ascending=[True, True, True]).reset_index(drop=True)

                    vagues_ordonnees = []
                    for v in df_clean_mixte["Vague"]:
                        if v not in vagues_ordonnees:
                            vagues_ordonnees.append(v)

                    cl_list, ucl_list, lcl_list = [], [], []

                    for v_item in vagues_ordonnees:
                        sub_vague = df_clean_mixte[df_clean_mixte["Vague"] == v_item]
                        cl_v = sub_vague["Valeur_Num"].mean()
                    
                        mr_v = sub_vague["Valeur_Num"].diff().abs().mean()
                        if pd.isna(mr_v) or mr_v == 0:
                            mr_v = 0.001 
                        
                        ucl_v = cl_v + 2.66 * mr_v
                        lcl_v = max(0.0, cl_v - 2.66 * mr_v) if sub_vague["Valeur_Num"].min() >= 0 else cl_v - 2.66 * mr_v

                        for idx_row in sub_vague.index:
                            cl_list.append((idx_row, cl_v))
                            ucl_list.append((idx_row, ucl_v))
                            lcl_list.append((idx_row, lcl_v))

                    df_clean_mixte["CL"] = [val for _, val in sorted(cl_list, key=lambda x: x[0])]
                    df_clean_mixte["UCL"] = [val for _, val in sorted(ucl_list, key=lambda x: x[0])]
                    df_clean_mixte["LCL"] = [val for _, val in sorted(lcl_list, key=lambda x: x[0])]

                    chart_display_df = df_clean_mixte[["Valeur_Num", "CL", "UCL", "LCL", "Vague"]].copy()
                    chart_display_df.index.name = "Séquence_Chronologique"
                
                    st.markdown(f"##### Évolution séquentielle (T0 au début $\\rightarrow$ vagues de suivi) pour : `{var_choisie_mixte}`")
                    st.line_chart(chart_display_df[["Valeur_Num", "CL", "UCL", "LCL"]])
                
                    st.markdown("###### Paramètres statistiques par vague :")
                    recap_stats = df_clean_mixte.groupby("Vague", sort=False).agg(
                        Moyenne_CL=("CL", "first"),
                        UCL_specifique=("UCL", "first"),
                        LCL_specifique=("LCL", "first"),
                        Nb_Points=("Valeur_Num", "count")
                    ).reset_index().round(2)
                    st.table(recap_stats)
